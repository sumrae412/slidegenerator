"""
Data Intelligence Module for Slide Generator

Detects numerical data in content and suggests appropriate chart visualizations
with automatic data extraction and PowerPoint chart generation.

Features:
- AI-powered chart type detection (bar, line, pie, scatter, column, area)
- Automatic data extraction from text
- PowerPoint chart generation using python-pptx
- Cost tracking for API usage
"""

import json
import logging
from typing import List, Dict, Optional, Any
from dataclasses import dataclass

# PowerPoint chart generation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

logger = logging.getLogger(__name__)


# Chart type definitions with use cases
CHART_TYPES = {
    'bar': 'Comparing categories or showing rankings',
    'column': 'Comparing values across categories',
    'line': 'Showing trends over time',
    'pie': 'Showing proportions or percentages',
    'scatter': 'Showing correlations between variables',
    'area': 'Showing cumulative trends'
}


@dataclass
class VisualizationConfig:
    """Configuration for a chart visualization"""
    should_visualize: bool
    chart_type: Optional[str] = None
    chart_title: Optional[str] = None
    data: Optional[Dict[str, Any]] = None
    x_label: Optional[str] = None
    y_label: Optional[str] = None
    confidence: float = 0.0
    reasoning: Optional[str] = None
    cost: float = 0.0

    def to_dict(self) -> dict:
        """Convert to dictionary"""
        return {
            'should_visualize': self.should_visualize,
            'chart_type': self.chart_type,
            'chart_title': self.chart_title,
            'data': self.data,
            'x_label': self.x_label,
            'y_label': self.y_label,
            'confidence': self.confidence,
            'reasoning': self.reasoning,
            'cost': self.cost
        }


class DataIntelligence:
    """
    Analyzes content for data visualization opportunities and generates charts.

    Uses AI to detect numerical data patterns and suggest appropriate chart types,
    then automatically extracts data and generates PowerPoint charts.
    """

    def __init__(self, client=None, cost_tracker=None):
        """
        Initialize DataIntelligence.

        Args:
            client: Anthropic client for AI analysis
            cost_tracker: Optional CostTracker instance for tracking API costs
        """
        self.client = client
        self.cost_tracker = cost_tracker

        # Chart type mapping for python-pptx
        self.chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
            'area': XL_CHART_TYPE.AREA
        }

    def suggest_visualization(
        self,
        text: str,
        bullets: List[str],
        slide_title: str = ""
    ) -> VisualizationConfig:
        """
        Detect if content should be visualized and extract chart configuration.

        Args:
            text: Original paragraph text
            bullets: Current bullet points
            slide_title: Slide title for context

        Returns:
            VisualizationConfig with chart details or should_visualize=False
        """
        if not self.client:
            logger.warning("No AI client provided - visualization detection disabled")
            return VisualizationConfig(should_visualize=False, reasoning="No AI client available")

        try:
            # Build analysis prompt
            prompt = self._build_visualization_prompt(text, bullets, slide_title)

            # Call Claude API
            response = self.client.messages.create(
                model="claude-sonnet-4-5-20250929",
                max_tokens=2000,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )

            # Extract response text
            response_text = response.content[0].text.strip()

            # Track cost
            cost = self._calculate_cost(response.usage)
            if self.cost_tracker:
                self.cost_tracker.add_cost('visualization_detection', cost)

            # Parse JSON response
            viz_data = self._parse_visualization_response(response_text)
            viz_data['cost'] = cost

            # Validate data structure
            if viz_data.get('should_visualize'):
                viz_data = self._validate_visualization_data(viz_data)

            return VisualizationConfig(**viz_data)

        except Exception as e:
            logger.error(f"Error in visualization detection: {e}")
            return VisualizationConfig(
                should_visualize=False,
                reasoning=f"Error: {str(e)}"
            )

    def _build_visualization_prompt(
        self,
        text: str,
        bullets: List[str],
        slide_title: str
    ) -> str:
        """Build the AI prompt for visualization detection"""
        bullets_text = "\n".join(f"- {b}" for b in bullets)

        return f"""Analyze this presentation content for data visualization opportunities.

Slide Title: {slide_title}

Original Text:
{text}

Current Bullets:
{bullets_text}

Your Task:
1. Determine if this content contains numerical data that would be better presented as a chart
2. If yes, identify the most effective chart type from: {', '.join(CHART_TYPES.keys())}
3. Extract the specific data points (labels and values)
4. Define appropriate axis labels and chart title

Chart Type Guidelines:
{chr(10).join(f"- {k}: {v}" for k, v in CHART_TYPES.items())}

Response Format:
Return ONLY valid JSON with this exact structure:

If visualization IS appropriate:
{{
  "should_visualize": true,
  "chart_type": "bar",
  "chart_title": "Customer Satisfaction Scores",
  "data": {{
    "labels": ["Support", "Product", "Documentation", "Onboarding"],
    "series": [
      {{
        "name": "Satisfaction Score",
        "values": [85, 92, 78, 88]
      }}
    ]
  }},
  "x_label": "Department",
  "y_label": "Score (0-100)",
  "confidence": 0.95,
  "reasoning": "Clear numerical comparison across categories - bar chart is ideal"
}}

If visualization is NOT appropriate:
{{
  "should_visualize": false,
  "reasoning": "Content is qualitative/conceptual without numerical data"
}}

Important:
- ONLY return JSON, no other text
- Ensure all numbers in "values" are numeric (int or float), not strings
- For pie charts, ensure values sum to a meaningful total (often 100 for percentages)
- Confidence should be 0.0 to 1.0
- Be conservative - only suggest visualization if data is clear and would genuinely improve understanding

Return your JSON analysis:"""

    def _parse_visualization_response(self, response_text: str) -> dict:
        """Parse the JSON response from Claude"""
        try:
            # Extract JSON from response (handle markdown code blocks)
            if '```json' in response_text:
                json_start = response_text.find('```json') + 7
                json_end = response_text.find('```', json_start)
                response_text = response_text[json_start:json_end].strip()
            elif '```' in response_text:
                json_start = response_text.find('```') + 3
                json_end = response_text.find('```', json_start)
                response_text = response_text[json_start:json_end].strip()

            data = json.loads(response_text)
            return data

        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse visualization JSON: {e}")
            logger.error(f"Response text: {response_text[:500]}")
            return {
                'should_visualize': False,
                'reasoning': f"JSON parsing error: {str(e)}"
            }

    def _validate_visualization_data(self, viz_data: dict) -> dict:
        """Validate and clean visualization data structure"""
        try:
            # Ensure required fields exist
            if not viz_data.get('chart_type') in self.chart_type_map:
                logger.warning(f"Invalid chart type: {viz_data.get('chart_type')}")
                viz_data['should_visualize'] = False
                viz_data['reasoning'] = f"Invalid chart type: {viz_data.get('chart_type')}"
                return viz_data

            # Validate data structure
            data = viz_data.get('data', {})
            if not data.get('labels') or not data.get('series'):
                logger.warning("Missing labels or series in chart data")
                viz_data['should_visualize'] = False
                viz_data['reasoning'] = "Invalid data structure - missing labels or series"
                return viz_data

            # Ensure values are numeric
            for series in data['series']:
                values = series.get('values', [])
                try:
                    # Convert to float to ensure numeric
                    series['values'] = [float(v) for v in values]
                except (ValueError, TypeError) as e:
                    logger.warning(f"Non-numeric values in series: {e}")
                    viz_data['should_visualize'] = False
                    viz_data['reasoning'] = "Data contains non-numeric values"
                    return viz_data

            # Validate labels match values count
            labels_count = len(data['labels'])
            for series in data['series']:
                if len(series['values']) != labels_count:
                    logger.warning(f"Mismatch between labels ({labels_count}) and values ({len(series['values'])})")
                    viz_data['should_visualize'] = False
                    viz_data['reasoning'] = "Mismatch between labels and values count"
                    return viz_data

            return viz_data

        except Exception as e:
            logger.error(f"Error validating visualization data: {e}")
            viz_data['should_visualize'] = False
            viz_data['reasoning'] = f"Validation error: {str(e)}"
            return viz_data

    def _calculate_cost(self, usage) -> float:
        """Calculate API cost from usage"""
        # Claude Sonnet 4.5 pricing (as of Oct 2024)
        input_cost_per_mtok = 3.00  # $3 per million input tokens
        output_cost_per_mtok = 15.00  # $15 per million output tokens

        input_tokens = getattr(usage, 'input_tokens', 0)
        output_tokens = getattr(usage, 'output_tokens', 0)

        cost = (input_tokens / 1_000_000 * input_cost_per_mtok +
                output_tokens / 1_000_000 * output_cost_per_mtok)

        return cost

    def create_chart_slide(
        self,
        viz_config: VisualizationConfig,
        prs
    ):
        """
        Create PowerPoint slide with chart using python-pptx.

        Args:
            viz_config: Visualization configuration from suggest_visualization()
            prs: PowerPoint presentation object (pptx.Presentation)

        Returns:
            Created slide with embedded chart, or None if creation failed
        """
        if not viz_config.should_visualize:
            logger.info("Visualization not recommended - skipping chart creation")
            return None

        try:
            # Get chart type
            chart_type = viz_config.chart_type
            if chart_type not in self.chart_type_map:
                logger.error(f"Unsupported chart type: {chart_type}")
                return None

            xl_chart_type = self.chart_type_map[chart_type]

            # Create chart data
            chart_data = CategoryChartData()
            chart_data.categories = viz_config.data['labels']

            # Add series
            for series in viz_config.data['series']:
                chart_data.add_series(series['name'], series['values'])

            # Create blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            # Add slide title
            if viz_config.chart_title:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
                )
                title_frame = title_box.text_frame
                title_frame.text = viz_config.chart_title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(28)
                title_para.font.bold = True

            # Add chart
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
            chart_shape = slide.shapes.add_chart(
                xl_chart_type, x, y, cx, cy, chart_data
            )

            chart = chart_shape.chart

            # Set chart title (if different from slide title)
            chart.has_title = False  # Use slide title instead

            # Add axis labels
            if chart_type not in ['pie']:  # Pie charts don't have axes
                try:
                    if viz_config.x_label:
                        chart.category_axis.has_title = True
                        chart.category_axis.axis_title.text_frame.text = viz_config.x_label

                    if viz_config.y_label:
                        chart.value_axis.has_title = True
                        chart.value_axis.axis_title.text_frame.text = viz_config.y_label
                except AttributeError:
                    # Some chart types might not support axis titles
                    logger.warning(f"Could not set axis labels for chart type: {chart_type}")

            logger.info(f"✅ Successfully created {chart_type} chart slide: {viz_config.chart_title}")
            return slide

        except Exception as e:
            logger.error(f"Error creating chart slide: {e}")
            logger.exception(e)
            return None

    def analyze_slide_for_visualization(
        self,
        slide_content,
        original_text: str = ""
    ) -> Optional[VisualizationConfig]:
        """
        Convenience method to analyze a SlideContent object for visualization.

        Args:
            slide_content: SlideContent object
            original_text: Original paragraph text (if available)

        Returns:
            VisualizationConfig if visualization is recommended, None otherwise
        """
        # Extract text from slide
        text = original_text if original_text else "\n".join(slide_content.content)
        bullets = slide_content.content
        title = slide_content.title

        # Analyze for visualization
        viz_config = self.suggest_visualization(text, bullets, title)

        # Only return if visualization is recommended
        if viz_config.should_visualize and viz_config.confidence >= 0.7:
            return viz_config

        return None
