"""
PowerPoint Presentation Generator

Handles generation of PowerPoint presentations with support for
visual prompts, AI-generated images, and intelligent slide layouts.

Extracted from file_to_slides.py (lines 7590-10789)
"""

import os
import logging
import random
from typing import List, Dict, Optional, Tuple, Any
from datetime import datetime
from math import cos, sin
from pathlib import Path
import io

# Presentation generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Image generation for sketches
from PIL import Image, ImageDraw, ImageFont
import requests

# Import data models
from .data_models import SlideContent, DocumentStructure

# Setup logging
logger = logging.getLogger(__name__)

# Constants
EXPORT_FOLDER = 'exports'


class SlideGenerator:
    """Handles generation of presentation slides"""
    
    def __init__(self, openai_client=None):
        """Initialize with optional OpenAI client for AI image generation"""
        self.client = openai_client

    def _insert_ai_image(self, slide, image_path: str, left: float = 5.2, top: float = 1.5,
                        width: float = 4.3, height: Optional[float] = None) -> bool:
        """
        Insert AI-generated image into a slide.

        Args:
            slide: PowerPoint slide object
            image_path: Path to the image file (local or URL)
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Optional height in inches (maintains aspect ratio if None)

        Returns:
            True if image was successfully inserted, False otherwise
        """
        try:
            # Handle remote URLs
            if image_path.startswith('http://') or image_path.startswith('https://'):
                logger.info(f"Downloading image from URL: {image_path[:100]}...")
                response = requests.get(image_path, timeout=30)
                response.raise_for_status()
                image_stream = io.BytesIO(response.content)
            else:
                # Local file path
                image_path_obj = Path(image_path)
                if not image_path_obj.exists():
                    logger.warning(f"Image file not found: {image_path}")
                    return False
                image_stream = str(image_path_obj)

            # Insert image into slide
            if height:
                slide.shapes.add_picture(image_stream, Inches(left), Inches(top),
                                        width=Inches(width), height=Inches(height))
            else:
                slide.shapes.add_picture(image_stream, Inches(left), Inches(top),
                                        width=Inches(width))

            logger.info(f"✅ Successfully inserted AI image into slide")
            return True

        except Exception as e:
            logger.error(f"Failed to insert AI image: {e}")
            return False

    def create_powerpoint(self, doc_structure: DocumentStructure, skip_visuals: bool = False) -> str:
        """Generate PowerPoint presentation with learner-focused content and optional visual prompts"""
        prs = Presentation()
        
        # Organize slides by heading hierarchy first to check for H1 headings
        organized_slides = self._organize_slides_by_hierarchy(doc_structure.slides)
        
        # Check if there's any H1 heading (presentation_title type) to decide on title slide
        has_h1_heading = any(section.get('type') == 'presentation_title' for section in organized_slides)
        
        # Only create title slide if there's a Heading 1 in the document
        if has_h1_heading:
            logger.info("Found H1 heading - creating title slide")
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = doc_structure.title
            subtitle.text = f"Generated from {doc_structure.metadata['filename']}\n{datetime.now().strftime('%B %d, %Y')}\n\nOptimized for Google Slides Import"
        else:
            logger.info("No H1 heading found - skipping title slide creation")
        
        logger.info(f"PowerPoint generation: Input has {len(doc_structure.slides)} slides")
        logger.info(f"After organization: {len(organized_slides)} organized sections")
        
        # Count organized slides by type
        slide_types = {}
        for section in organized_slides:
            slide_type = section.get('type', 'unknown')
            slide_types[slide_type] = slide_types.get(slide_type, 0) + 1
        
        logger.info(f"Organized slide types: {slide_types}")
        
        # Create slides based on hierarchy: H1=presentation title, H2=section, H3=subsection, H4=slide titles, cells=content
        pptx_slides_created = 0
        for section in organized_slides:
            if section['type'] == 'presentation_title':
                # This is handled by the main title slide above, skip
                continue
                
            elif section['type'] == 'divider':
                # Section divider - visual break slide
                divider_slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
                slide = prs.slides.add_slide(divider_slide_layout)
                pptx_slides_created += 1
                logger.info(f"Created section divider slide: '{section['title']}'")

                # Center the section title on a minimal slide
                # Add centered title shape
                title_shape = slide.shapes.add_textbox(
                    left=Inches(1), top=Inches(3),
                    width=Inches(8), height=Inches(1.5)
                )
                title_frame = title_shape.text_frame
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.text = section['title']
                title_paragraph.font.size = Pt(54)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(66, 66, 66)  # Dark gray
                title_paragraph.alignment = PP_ALIGN.CENTER

                # Optional: Add a subtle horizontal line above title for visual effect
                line_shape = slide.shapes.add_shape(
                    1,  # Line shape type
                    left=Inches(2), top=Inches(2.7),
                    width=Inches(6), height=Inches(0)
                )
                line_shape.line.color.rgb = RGBColor(189, 189, 189)  # Light gray
                line_shape.line.width = Pt(2)

            elif section['type'] == 'section_title':
                # H2 = Section title page
                section_slide_layout = prs.slide_layouts[5]  # Title and content layout
                slide = prs.slides.add_slide(section_slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                pptx_slides_created += 1
                logger.info(f"Created section title slide: '{section['title']}'")

                # Style as section title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(48)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(25, 118, 210)  # Blue color

                # Add overview if available
                if section.get('overview'):
                    content_shape = None
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            content_shape = placeholder
                            break
                    if content_shape:
                        content_shape.text = section['overview']
                        
            elif section['type'] == 'subsection_title':
                # H3 = Subsection title page
                subsection_slide_layout = prs.slide_layouts[5]  # Title and content layout
                slide = prs.slides.add_slide(subsection_slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                pptx_slides_created += 1
                logger.info(f"Created subsection title slide: '{section['title']}'")
                
                # Style as subsection title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(40)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(102, 126, 234)  # Lighter blue
                
                # Add overview if available
                if section.get('overview'):
                    content_shape = None
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            content_shape = placeholder
                            break
                    if content_shape:
                        content_shape.text = section['overview']
                        
            elif section['type'] == 'slide_title':
                # H4 = Individual slide with content
                slide_layout = prs.slide_layouts[5]  # Title and content layout
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                pptx_slides_created += 1
                logger.info(f"Created H4 slide: '{section['title']}'")
                
                # Style as individual slide title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(36)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(55, 71, 79)  # Dark gray
                
                # Add content if available
                if section.get('content'):
                    content_shape = None
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            content_shape = placeholder
                            break
                    if content_shape:
                        # Convert content to bullet points if it's a list
                        if isinstance(section['content'], list):
                            content_text = '\n'.join([f"• {item}" for item in section['content']])
                        else:
                            # Process text content into bullet points
                            _, bullet_points = self._create_bullet_points(str(section['content']), skip_visuals)
                            content_text = '\n'.join([f"• {bullet}" for bullet in bullet_points])
                        
                        content_shape.text = content_text
                        
            elif section['type'] == 'script':
                # Script content - two-column layout with bullet points on left and sketch suggestion on right
                script_slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
                slide = prs.slides.add_slide(script_slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(
                    left=Inches(0.5), top=Inches(0.3), 
                    width=Inches(9), height=Inches(1)
                )
                title_frame = title_shape.text_frame
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.text = section['title']
                title_paragraph.font.size = Pt(32)
                title_paragraph.font.bold = True
                title_paragraph.alignment = PP_ALIGN.CENTER
                
                # Get the current slide number BEFORE incrementing for unique image generation
                current_slide_number = pptx_slides_created + 1
                pptx_slides_created += 1
                logger.info(f"Created script slide #{current_slide_number}: '{section['title']}' with {len(section.get('content', []))} bullet points")
                
                # Left side - Subheader and Bullet points
                if section.get('content') or section.get('subheader'):
                    bullet_shape = slide.shapes.add_textbox(
                        left=Inches(0.5), top=Inches(1.5),
                        width=Inches(4.5), height=Inches(5)
                    )
                    bullet_frame = bullet_shape.text_frame
                    bullet_frame.word_wrap = True
                    bullet_frame.clear()

                    para_index = 0

                    # Add subheader if present (bold topic sentence above bullets)
                    if section.get('subheader'):
                        p = bullet_frame.paragraphs[0]
                        p.text = section['subheader']
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.bold = True
                        p.space_after = Pt(12)
                        p.space_before = Pt(0)
                        para_index += 1

                    # Add bullet points
                    for i, bullet_point in enumerate(section.get('content', [])):
                        p = bullet_frame.paragraphs[para_index] if (i == 0 and para_index == 0) else bullet_frame.add_paragraph()
                        p.text = f"• {bullet_point}"
                        p.level = 0
                        p.font.size = Pt(16)
                        p.space_after = Pt(8)
                        p.space_before = Pt(4)
                
                # Right side - AI-generated image or visual prompt text
                # Priority: 1) AI image, 2) Visual prompt text, 3) Nothing (fast mode)

                # Check if AI-generated image is available
                visual_image_path = section.get('visual_image_path') or section.get('visual_image_url')

                if visual_image_path:
                    # Insert AI-generated image
                    logger.info(f"Inserting AI-generated image for slide {current_slide_number}: {visual_image_path[:100]}")
                    image_inserted = self._insert_ai_image(
                        slide=slide,
                        image_path=visual_image_path,
                        left=5.2,
                        top=0.5,
                        width=4.3
                    )

                    if image_inserted and section.get('visual_prompt'):
                        # Add visual prompt as small caption below image
                        try:
                            caption_textbox = slide.shapes.add_textbox(
                                left=Inches(5.2),
                                top=Inches(5.5),
                                width=Inches(4.3),
                                height=Inches(1.0)
                            )
                            caption_frame = caption_textbox.text_frame
                            caption_frame.word_wrap = True
                            caption_p = caption_frame.paragraphs[0]
                            caption_p.text = f"🎨 {section['visual_prompt'][:150]}..."
                            caption_p.font.size = Pt(8)
                            caption_p.font.italic = True
                            caption_p.font.color.rgb = RGBColor(120, 120, 120)
                        except Exception as e:
                            logger.warning(f"Failed to add image caption: {e}")

                elif not skip_visuals:
                    # No AI image - generate visual prompt text instead
                    logger.info(f"Generating visual prompt text for slide {current_slide_number}: '{section['title']}' with {len(section.get('content', []))} bullet points")
                    visual_prompt_text = self._generate_drawing_prompt(section.get('content', []), section['title'], current_slide_number)
                    logger.info(f"Generated visual prompt text for slide {current_slide_number}")

                    # Add visual prompt as copyable text on the right side
                    if visual_prompt_text:
                        try:
                            # Add the visual prompt as copyable text box on the right side, aligned with slide top
                            prompt_textbox = slide.shapes.add_textbox(
                                left=Inches(5.2),
                                top=Inches(0.5),
                                width=Inches(4.3),
                                height=Inches(6)
                            )
                            prompt_frame = prompt_textbox.text_frame
                            prompt_frame.clear()
                            prompt_frame.margin_left = Inches(0.2)
                            prompt_frame.margin_right = Inches(0.2)
                            prompt_frame.margin_top = Inches(0.2)
                            prompt_frame.margin_bottom = Inches(0.2)
                            prompt_frame.word_wrap = True

                            # Add header
                            header_p = prompt_frame.paragraphs[0]
                            header_p.text = "🎨 Visual Prompt"
                            header_p.font.size = Pt(14)
                            header_p.font.name = 'Arial'
                            header_p.font.bold = True
                            header_p.font.color.rgb = RGBColor(50, 100, 200)
                            header_p.alignment = PP_ALIGN.LEFT

                            # Add separator line
                            sep_p = prompt_frame.add_paragraph()
                            sep_p.text = "─" * 40
                            sep_p.font.size = Pt(8)
                            sep_p.font.color.rgb = RGBColor(180, 180, 180)

                            # Add the full visual prompt text
                            prompt_p = prompt_frame.add_paragraph()
                            prompt_p.text = visual_prompt_text
                            prompt_p.font.size = Pt(9)
                            prompt_p.font.name = 'Arial'
                            prompt_p.font.color.rgb = RGBColor(60, 60, 60)
                            prompt_p.alignment = PP_ALIGN.LEFT

                            # Add copy instruction
                            instruction_p = prompt_frame.add_paragraph()
                            instruction_p.text = "\n💡 Copy this text to use with AI image generators (DALL-E, Midjourney, Stable Diffusion)"
                            instruction_p.font.size = Pt(8)
                            instruction_p.font.name = 'Arial'
                            instruction_p.font.italic = True
                            instruction_p.font.color.rgb = RGBColor(120, 120, 120)
                            instruction_p.alignment = PP_ALIGN.LEFT

                            logger.info(f"Successfully added copyable visual prompt text for slide: {section['title']}")
                        except Exception as e:
                            logger.error(f"Could not add visual prompt text for '{section['title']}': {e}")
                else:
                    logger.info(f"Skipping visual generation for slide {current_slide_number} (fast mode)")
                            
            else:
                # Other content slides
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                
                # Style title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(28)
                
                # Add content
                if section.get('content'):
                    content_shape = None
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            content_shape = placeholder
                            break
                    
                    if content_shape:
                        text_frame = content_shape.text_frame
                        text_frame.clear()
                        
                        for i, content_item in enumerate(section['content'][:6]):
                            p = text_frame.add_paragraph()
                            p.text = content_item
                            p.level = 0
                            p.font.size = Pt(20)
                            p.space_after = Pt(8)
        
        # Add a summary slide at the end
        summary_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Summary"
        
        # Find content placeholder safely for summary
        content_shape = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 1:
                content_shape = placeholder
                break
        
        if content_shape:
            summary_text = f"• {len(organized_slides)} main sections covered\n"
            summary_text += f"• Generated from {doc_structure.metadata['filename']}\n"
            summary_text += f"• Ready for presentation and editing\n\n"
            summary_text += "Next steps:\n"
            summary_text += "• Customize themes and layouts as needed\n"
            summary_text += "• Review and edit visual prompts\n"
            summary_text += "• Add additional images or formatting"
            
            content_shape.text = summary_text
        
        # Calculate total slides including optional title slide
        title_slide_count = 1 if has_h1_heading else 0
        total_slides = pptx_slides_created + title_slide_count + 1  # +1 for summary slide
        logger.info(f"POWERPOINT FINAL: Created {pptx_slides_created} content slides + {title_slide_count} title slide + 1 summary slide = {total_slides} total slides")
        
        # Save presentation
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(EXPORT_FOLDER, filename)
        os.makedirs(EXPORT_FOLDER, exist_ok=True)
        prs.save(filepath)
        
        return filepath
    
    def _organize_slides_by_hierarchy(self, slides: List[SlideContent]) -> List[Dict]:
        """Organize slides by heading hierarchy: H1=presentation title, H2=section, H3=subsection, H4=slide titles, cells=content"""
        organized = []
        presentation_title_added = False
        
        for slide in slides:
            title = slide.title
            content = slide.content
            slide_type = getattr(slide, 'slide_type', 'content')  # Get slide_type attribute
            
            # Check if this is explicitly a heading slide
            if slide_type == 'heading':
                # Use original heading level if available, otherwise determine from title
                level = getattr(slide, 'heading_level', None) or self._determine_heading_level(title)
                
                if level == 1 and not presentation_title_added:
                    # H1 = Title page for entire presentation (only use first H1)
                    organized.append({
                        'type': 'presentation_title',
                        'title': title,
                        'level': 1,
                        'overview': self._create_section_overview(content) if content else None
                    })
                    presentation_title_added = True
                    
                elif level == 1 or level == 2:
                    # H2 = Section title page
                    organized.append({
                        'type': 'section_title',
                        'title': title,
                        'level': 2,
                        'overview': self._create_section_overview(content) if content else None
                    })
                    
                elif level == 3:
                    # H3 = Subsection title page
                    organized.append({
                        'type': 'subsection_title',
                        'title': title,
                        'level': 3,
                        'overview': self._create_section_overview(content) if content else None
                    })
                    
                elif level == 4:
                    # H4 = Individual slide title (content slides with H4 titles)
                    organized.append({
                        'type': 'slide_title',
                        'title': title,
                        'content': content,
                        'level': 4
                    })

            elif slide_type == 'divider':
                # Section divider - visual break slide
                organized.append({
                    'type': 'divider',
                    'title': title,
                    'content': [],  # Dividers have no content
                    'level': getattr(slide, 'heading_level', 2)
                })

            elif slide_type == 'script':
                # Script content slides - each table row becomes one slide with bullet points
                organized.append({
                    'type': 'script',
                    'title': title,
                    'content': content,
                    'level': 4,
                    'subheader': getattr(slide, 'subheader', None),
                    'visual_cues': getattr(slide, 'visual_cues', None),
                    # AI visual generation fields
                    'visual_prompt': getattr(slide, 'visual_prompt', None),
                    'visual_image_url': getattr(slide, 'visual_image_url', None),
                    'visual_image_path': getattr(slide, 'visual_image_path', None),
                    'visual_type': getattr(slide, 'visual_type', None)
                })
                
            else:
                # Other content slides (fallback)
                if title.startswith('Table Cell'):
                    # Each table cell gets its own slide with blank title (only H4 should be titles)
                    organized.append({
                        'type': 'cell_content',
                        'title': "",  # Blank title - only H4 headings should be slide titles
                        'content': content,
                        'level': 4
                    })
                else:
                    # Other content slides with blank titles
                    organized.append({
                        'type': 'content',
                        'title': "",  # Blank title - only H4 headings should be slide titles
                        'content': content,
                        'level': 4
                    })
        
        return organized
    
    def _generate_sketch_suggestion(self, bullet_points: List[str], title: str) -> str:
        """Generate a contextual sketch suggestion based on specific bullet point content"""
        if not bullet_points:
            return "Simple concept illustration related to the topic"
        
        # Extract key concepts directly from bullet points
        key_concepts = self._extract_key_concepts_from_bullets(bullet_points)
        
        # Combine title and bullet points for analysis
        all_text = f"{title} {' '.join(bullet_points)}".lower()
        
        # Generate specific suggestion based on content
        if any(word in all_text for word in ['step', 'process', 'workflow', 'method', 'approach', 'first', 'then', 'next']):
            return f"Process flow showing: {key_concepts}. Draw connected steps with arrows and icons."
        
        elif any(word in all_text for word in ['compare', 'vs', 'versus', 'difference', 'traditional', 'new', 'before', 'after']):
            return f"Comparison diagram: {key_concepts}. Show side-by-side with contrasting elements."
        
        elif any(word in all_text for word in ['ai', 'genai', 'chatbot', 'model', 'prompt', 'artificial intelligence']):
            return f"AI concept visualization: {key_concepts}. Include neural network or brain diagram with data flows."
        
        elif any(word in all_text for word in ['ui', 'interface', 'design', 'user', 'screen', 'app', 'button', 'form']):
            return f"UI mockup showing: {key_concepts}. Create wireframe with interface elements and user interactions."
        
        elif any(word in all_text for word in ['data', 'chart', 'graph', 'metrics', 'analysis', 'dashboard', 'report']):
            return f"Data visualization: {key_concepts}. Show charts, graphs, and key metrics dashboard."
        
        elif any(word in all_text for word in ['system', 'architecture', 'component', 'api', 'database', 'server']):
            return f"System diagram: {key_concepts}. Draw components with connections and data flow arrows."
        
        elif any(word in all_text for word in ['learn', 'course', 'lesson', 'training', 'education', 'skill']):
            return f"Learning path: {key_concepts}. Show progression with milestones and knowledge building blocks."
        
        elif any(word in all_text for word in ['problem', 'solution', 'challenge', 'issue', 'fix', 'solve']):
            return f"Problem-solution flow: {key_concepts}. Use warning icons for problems, lightbulbs for solutions."
        
        elif any(word in all_text for word in ['benefit', 'advantage', 'improve', 'better', 'faster', 'efficiency']):
            return f"Benefits diagram: {key_concepts}. Use checkmarks, positive icons, and improvement arrows."
        
        elif any(word in all_text for word in ['team', 'collaboration', 'meeting', 'communication', 'stakeholder']):
            return f"Collaboration diagram: {key_concepts}. Show people icons, communication flows, and team interactions."
        
        elif any(word in all_text for word in ['timeline', 'schedule', 'roadmap', 'phases', 'milestone', 'time']):
            return f"Timeline visualization: {key_concepts}. Create timeline with key events and progress markers."
        
        else:
            # Default: create conceptual diagram based on content
            return f"Concept map: {key_concepts}. Show main ideas connected with relationships and supporting details."
    
    def _extract_key_concepts_from_bullets(self, bullet_points: List[str]) -> str:
        """Extract the most important concepts from bullet points for sketch suggestions"""
        if not bullet_points:
            return "main concepts"
        
        # Get first few key words from each bullet point
        concepts = []
        for bullet in bullet_points[:3]:  # Focus on first 3 bullets
            # Clean up bullet text and extract key terms
            clean_bullet = bullet.strip()
            if clean_bullet:
                # Take first significant phrase (up to first punctuation or first 4 words)
                words = clean_bullet.split()[:4]
                key_phrase = ' '.join(words)
                concepts.append(key_phrase)
        
        if concepts:
            return ' → '.join(concepts)
        else:
            return 'key concepts from content'
    
    def _create_content_specific_sketch(self, bullet_points: List[str], title: str, slide_number: int = 0) -> str:
        """Create a thought bubble with visual prompt on light blue background"""
        logger.info(f"_create_content_specific_sketch called for slide {slide_number}, title: '{title}', bullets: {len(bullet_points)}")
        if not bullet_points and not title:
            logger.warning(f"No content for sketch - bullets: {bullet_points}, title: '{title}'")
            return None
            
        try:
            # Generate a unique drawing prompt
            drawing_prompt = self._generate_drawing_prompt(bullet_points, title, slide_number)
            logger.info(f"Generated drawing prompt for slide {slide_number}: {drawing_prompt}")
            
            # Keep text file creation for backup
            self._create_prompt_text_file(drawing_prompt, slide_number)
            
            # Generate AI image using OpenAI DALL-E if API key available
            logger.info(f"Checking AI image generation for slide {slide_number}")
            logger.info(f"self.client is not None: {self.client is not None}")
            if self.client:
                logger.info(f"Attempting AI image generation for slide {slide_number}")
                return self._generate_ai_image(drawing_prompt, slide_number)
            else:
                logger.info(f"No OpenAI client available, falling back to thought bubble for slide {slide_number}")
                # Fallback to thought bubble image with the prompt
                return self._create_thought_bubble_image(drawing_prompt, slide_number)
            
            # Analyze the specific content to determine the best sketch approach
            all_content = f"{title} {' '.join(bullet_points)}".lower()
            logger.info(f"All content for analysis: '{all_content[:100]}...'")
            
            # Create image with the same quality as before
            width, height = 800, 600
            image = Image.new('RGB', (width, height), '#f8f9fa')
            draw = ImageDraw.Draw(image)
            
            # Load fonts
            try:
                font_paths = [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/System/Library/Fonts/Arial.ttf", 
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/Windows/Fonts/arial.ttf"
                ]
                font_loaded = False
                for font_path in font_paths:
                    try:
                        font = ImageFont.truetype(font_path, 18)
                        title_font = ImageFont.truetype(font_path, 24)
                        small_font = ImageFont.truetype(font_path, 14)
                        large_font = ImageFont.truetype(font_path, 20)
                        font_loaded = True
                        break
                    except:
                        continue
                
                if not font_loaded:
                    raise Exception("No fonts found")
                    
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
                large_font = ImageFont.load_default()
            
            # Determine sketch type based on actual bullet point content
            sketch_type = self._determine_content_specific_sketch_type(bullet_points, title)
            logger.info(f"Determined sketch type: '{sketch_type}' for slide {slide_number}")
            
            # Draw the appropriate sketch based on content analysis
            if sketch_type == 'process':
                self._draw_process_from_content(draw, width, height, font, bullet_points, title)
            elif sketch_type == 'comparison':
                self._draw_comparison_from_content(draw, width, height, font, bullet_points, title)
            elif sketch_type == 'ai_tech':
                self._draw_ai_concept(draw, width, height, font, title_font)
            elif sketch_type == 'ui_design':
                self._draw_ui_design_from_content(draw, width, height, font, bullet_points, title)
            elif sketch_type == 'data_metrics':
                self._draw_data_visualization(draw, width, height, font, title_font)
            elif sketch_type == 'system_arch':
                self._draw_system_architecture(draw, width, height, font, title_font)
            elif sketch_type == 'learning_path':
                self._draw_learning_path(draw, width, height, font, title_font)
            elif sketch_type == 'benefits':
                self._draw_benefits_diagram(draw, width, height, font, bullet_points)
            else:
                # Default: create concept map from bullet points
                self._draw_concept_map_from_bullets(draw, width, height, font, bullet_points, title)
            
            # Add slide title at the top (smaller and cleaner)
            if title:
                bbox = draw.textbbox((0, 0), title, font=title_font)
                title_width = bbox[2] - bbox[0]
                title_height = bbox[3] - bbox[1]
                
                # Simple title without heavy background
                draw.text((width//2 - title_width//2, 20), title, fill='#333333', font=title_font)
            
            # Save with unique naming
            os.makedirs('temp_sketches', exist_ok=True)
            import random
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            random_component = random.randint(1000, 9999)
            unique_id = abs(hash(f"{slide_number}_{title}_{str(bullet_points)}_{sketch_type}_{timestamp}_{random_component}"))
            image_path = f'temp_sketches/sketch_slide{slide_number:03d}_{sketch_type}_{unique_id}_{timestamp}_{random_component}.png'
            
            image.save(image_path)
            logger.info(f"Created content-specific sketch: {image_path}")
            return image_path
            
        except Exception as e:
            logger.error(f"Error creating content-specific sketch: {e}")
            return None
    
    def _create_llm_generated_diagram(self, bullet_points: List[str], title: str, slide_number: int) -> str:
        """Use LLM to generate a diagram based on slide content"""
        try:
            import openai
            
            # Prepare the content for LLM analysis
            content_text = f"Title: {title}\nBullet Points:\n" + "\n".join([f"• {bullet}" for bullet in bullet_points])
            
            # Create a prompt for diagram generation
            diagram_prompt = f"""
Analyze the following slide content and create a detailed description for a diagram that would best illustrate the concepts:

{content_text}

Please respond with a JSON object containing:
1. "diagram_type": The best type of diagram (flowchart, comparison, concept_map, process, timeline, hierarchy, etc.)
2. "description": A detailed description of what the diagram should show
3. "elements": A list of key visual elements to include
4. "layout": Suggested layout/structure
5. "svg_code": Simple SVG code (800x600) that represents this diagram with professional styling

Focus on creating a clean, professional diagram that clearly communicates the key concepts from the slide content.
"""

            # Get API key from environment or config
            api_key = os.getenv('OPENAI_API_KEY')
            if not api_key:
                logger.info("No OpenAI API key found, skipping LLM diagram generation")
                return None
            
            # Make API call to OpenAI
            try:
                client = openai.OpenAI(api_key=api_key)
            except Exception as e:
                logger.error(f"Failed to initialize OpenAI client: {e}")
                return None
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",  # Faster, cheaper model
                messages=[
                    {"role": "system", "content": "Create SVG diagrams. Respond with JSON: {'svg_content': '<svg>...</svg>', 'diagram_type': 'process|comparison|flowchart'}"},
                    {"role": "user", "content": diagram_prompt}
                ],
                max_tokens=1000,  # Shorter responses for speed
                temperature=0.2   # Less creativity for speed
            )
            
            # Parse the response
            import json
            response_content = response.choices[0].message.content
            logger.info(f"LLM response for slide {slide_number}: {response_content[:200]}...")
            
            try:
                diagram_data = json.loads(response_content)
            except json.JSONDecodeError:
                logger.error(f"Failed to parse LLM response as JSON for slide {slide_number}")
                return None
            
            if 'svg_code' in diagram_data:
                # Convert SVG to PNG
                svg_content = diagram_data['svg_code']
                logger.info(f"Generated SVG for slide {slide_number}: {len(svg_content)} characters")
                return self._convert_svg_to_image(svg_content, slide_number, diagram_data.get('diagram_type', 'llm'))
            
            logger.warning(f"No SVG code in LLM response for slide {slide_number}")
            return None
            
        except Exception as e:
            logger.error(f"Error in LLM diagram generation: {e}")
            return None
    
    def _convert_svg_to_image(self, svg_content: str, slide_number: int, diagram_type: str) -> str:
        """Convert SVG content to PNG image"""
        try:
            # Save the SVG content to a file for inspection
            svg_filename = f"temp_sketches/slide{slide_number:03d}_{diagram_type}_debug.svg"
            os.makedirs('temp_sketches', exist_ok=True)
            with open(svg_filename, 'w', encoding='utf-8') as f:
                f.write(svg_content)
            logger.info(f"Saved SVG content to {svg_filename}")
            
            # Try using cairosvg if available and working
            try:
                import cairosvg
                from io import BytesIO
                
                # Convert SVG to PNG
                png_data = cairosvg.svg2png(bytestring=svg_content.encode('utf-8'))
                
                # Create unique filename
                content_hash = abs(hash(svg_content)) % (10**16)
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
                random_id = random.randint(1000, 9999)
                filename = f"sketch_slide{slide_number:03d}_llm_{diagram_type}_{content_hash}_{timestamp}_{random_id}.png"
                
                # Ensure temp_sketches directory exists
                image_path = os.path.join('temp_sketches', filename)
                
                # Save the PNG
                with open(image_path, 'wb') as f:
                    f.write(png_data)
                
                logger.info(f"Successfully converted LLM SVG to PNG: {image_path}")
                return image_path
                
            except Exception as cairo_error:
                logger.warning(f"cairosvg conversion failed: {cairo_error}")
                
                # Second attempt: Use ImageMagick convert command
                try:
                    import subprocess
                    
                    # Create unique filename
                    content_hash = abs(hash(svg_content)) % (10**16)
                    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
                    random_id = random.randint(1000, 9999)
                    filename = f"sketch_slide{slide_number:03d}_llm_{diagram_type}_{content_hash}_{timestamp}_{random_id}.png"
                    
                    # Save SVG temporarily
                    temp_svg = os.path.join('temp_sketches', f'temp_{slide_number}_{diagram_type}.svg')
                    with open(temp_svg, 'w', encoding='utf-8') as f:
                        f.write(svg_content)
                    
                    # Convert using ImageMagick
                    image_path = os.path.join('temp_sketches', filename)
                    subprocess.run([
                        'convert', 
                        '-background', 'white',
                        '-density', '96',
                        temp_svg,
                        image_path
                    ], check=True, capture_output=True)
                    
                    # Clean up temp SVG
                    os.remove(temp_svg)
                    
                    logger.info(f"Successfully converted SVG using ImageMagick: {image_path}")
                    return image_path
                    
                except Exception as magick_error:
                    logger.warning(f"ImageMagick conversion failed: {magick_error}, using fallback")
                    return self._create_image_from_svg_description(svg_content, slide_number, diagram_type)
            
        except Exception as e:
            logger.error(f"Error converting SVG to image: {e}")
            return None

    def _generate_drawing_prompt(self, bullet_points: List[str], title: str, slide_number: int) -> str:
        """Generate simple, flat-design visual concepts for educational slides"""
        try:
            # Analyze content to determine main educational concepts
            all_text = f"{title} {' '.join(bullet_points)}".lower()
            
            # Identify the core educational theme
            visual_prompt = self._create_educational_visual_concept(title, all_text, bullet_points)
            
            return visual_prompt
            
        except Exception as e:
            logger.error(f"Error generating visual prompt: {e}")
            return self._get_fallback_visual_prompt()
    
    def _create_educational_visual_concept(self, title: str, all_text: str, bullet_points: List[str]) -> str:
        """Create educational flat-design visual concepts using the new template format"""
        
        # Use the new prompt template format
        visual_prompt = "Task: Create a visual concept for a course slide.\n\n"
        
        # Slide Purpose - extract learning goal from title and content
        purpose = self._extract_slide_purpose(title, all_text, bullet_points)
        visual_prompt += f"Slide Purpose: {purpose}\n\n"
        
        # Key Ideas / Bullet Points
        visual_prompt += "Key Ideas / Bullet Points:\n"
        if bullet_points:
            for bullet in bullet_points[:4]:  # Keep concise
                visual_prompt += f"    •    {bullet}\n"
        else:
            visual_prompt += f"    •    {title}\n"
        visual_prompt += "\n"
        
        # Visual Goal - determine appropriate visual narrative
        visual_goal = self._create_visual_goal(all_text, bullet_points)
        visual_prompt += f"Visual Goal: {visual_goal}\n\n"
        
        # Style Guidelines (consistent for all slides)
        visual_prompt += "Style Guidelines:\n"
        visual_prompt += "    •    Flat design or minimal 2D illustration\n"
        visual_prompt += "    •    Clean, modern, and easy to animate\n"
        visual_prompt += "    •    No text in the image\n"
        visual_prompt += "    •    No abstract symbols or meaningless geometry\n"
        visual_prompt += "    •    Prefer narrative visuals (e.g., person solving a problem, building an app, exploring data)\n"
        visual_prompt += "    •    Use metaphors only if they directly reinforce the concept\n\n"
        
        # Intended Use
        visual_prompt += "Intended Use: Slide visuals and animations in an online GenAI course. Learners are developers or data scientists prototyping with LLMs."
        
        return visual_prompt
    
    def _extract_slide_purpose(self, title: str, all_text: str, bullet_points: List[str]) -> str:
        """Extract the main learning purpose from slide content"""
        # Analyze content to determine educational purpose
        
        # Check for specific learning patterns
        if any(word in all_text for word in ['prototype', 'build', 'create', 'develop']):
            if 'genai' in all_text or 'ai' in all_text:
                return "Teach learners how to prototype and build applications using GenAI tools"
            else:
                return f"Show learners how to build and develop {title.lower()}"
        
        elif any(word in all_text for word in ['understand', 'learn', 'concept']):
            return f"Help learners understand key concepts in {title.lower()}"
        
        elif any(word in all_text for word in ['compare', 'versus', 'difference']):
            return f"Show learners the differences and trade-offs in {title.lower()}"
        
        elif any(word in all_text for word in ['test', 'evaluate', 'measure']):
            return f"Teach learners how to test and evaluate {title.lower()}"
        
        else:
            # Default based on title and first bullet point
            first_bullet = bullet_points[0] if bullet_points else title
            return f"Teach learners about {first_bullet.lower()}"
    
    def _create_visual_goal(self, all_text: str, bullet_points: List[str]) -> str:
        """Create content-specific visual goal based on slide content"""
        
        # Combine all content for analysis
        combined_content = f"{all_text} {' '.join(bullet_points)}".lower()
        
        # Content-specific visual scenarios
        visual_scenarios = []
        
        # Development/Coding scenarios
        if any(word in combined_content for word in ['code', 'programming', 'python', 'streamlit', 'api']):
            visual_scenarios.extend([
                "Show a developer at a computer writing code, with code snippets floating around them in a creative, organized way",
                "Illustrate a person building blocks of code, stacking them like LEGO pieces to create an application",
                "Show someone debugging code with a magnifying glass, finding and fixing errors in a visual code landscape",
                "Depict a developer testing their app on multiple devices, showing the app responding and working correctly"
            ])
        
        # GenAI/AI scenarios
        if any(word in combined_content for word in ['genai', 'ai', 'artificial intelligence', 'chatgpt', 'llm']):
            visual_scenarios.extend([
                "Show a person having a conversation with an AI assistant, with creative ideas flowing between them as visual elements",
                "Illustrate someone feeding data into an AI system and receiving creative solutions and insights in return",
                "Depict a human and AI working together to solve a complex problem, showing collaboration and synergy",
                "Show an AI helping a developer brainstorm and iterate on app ideas, with multiple concept bubbles around them"
            ])
        
        # Data/Analytics scenarios
        if any(word in combined_content for word in ['data', 'analysis', 'dashboard', 'chart', 'visualization']):
            visual_scenarios.extend([
                "Show someone exploring data like an archaeologist, uncovering insights and patterns from raw information",
                "Illustrate a person transforming messy data into clear, beautiful visualizations and charts",
                "Depict someone making data-driven decisions, with charts and graphs helping them choose the right path",
                "Show a data detective solving mysteries by analyzing trends and patterns in colorful data landscapes"
            ])
        
        # Learning/Education scenarios
        if any(word in combined_content for word in ['learn', 'teach', 'understand', 'concept', 'course', 'lesson']):
            visual_scenarios.extend([
                "Show a lightbulb moment - a person suddenly understanding a complex concept with clarity and excitement",
                "Illustrate someone climbing a mountain of knowledge, reaching new levels of understanding step by step",
                "Depict a person connecting dots between different concepts, creating a web of interconnected knowledge",
                "Show someone teaching others, sharing knowledge in an engaging, visual way with enthusiasm"
            ])
        
        # Process/Workflow scenarios
        if any(word in combined_content for word in ['process', 'step', 'workflow', 'method', 'procedure']):
            visual_scenarios.extend([
                "Show someone following a clear path through a process, with each step marked and progress clearly visible",
                "Illustrate a person organizing chaos into structured, efficient workflows like a conductor leading an orchestra",
                "Depict someone streamlining a complex process, removing unnecessary steps and optimizing efficiency",
                "Show a team working together in a well-orchestrated process, each person contributing their part"
            ])
        
        # Problem-solving scenarios
        if any(word in combined_content for word in ['problem', 'solution', 'fix', 'debug', 'troubleshoot']):
            visual_scenarios.extend([
                "Show a detective solving a mystery, using clues and evidence to find the solution to a complex problem",
                "Illustrate someone untangling a knot, gradually working through complexity to find a clear solution",
                "Depict a person building a bridge across a gap, connecting problems with creative solutions",
                "Show someone with a toolkit, methodically fixing and improving a broken system or process"
            ])
        
        # Testing/Evaluation scenarios
        if any(word in combined_content for word in ['test', 'evaluate', 'measure', 'validate', 'quality']):
            visual_scenarios.extend([
                "Show someone as a quality inspector, carefully examining work with attention to detail and precision",
                "Illustrate a person conducting experiments, testing different approaches to find what works best",
                "Depict someone using measuring tools to evaluate progress and ensure standards are met",
                "Show a feedback loop in action, with someone making improvements based on test results"
            ])
        
        # Collaboration scenarios
        if any(word in combined_content for word in ['team', 'collaborate', 'together', 'group', 'share']):
            visual_scenarios.extend([
                "Show people working together like puzzle pieces, each contributing unique skills to complete the picture",
                "Illustrate a diverse team brainstorming around a table, with creative ideas flowing between all participants",
                "Depict people building something together, each person adding their expertise to create something amazing",
                "Show team members passing ideas like a relay race, building on each other's contributions"
            ])
        
        # If we have content-specific scenarios, pick one randomly
        if visual_scenarios:
            selected_scenario = random.choice(visual_scenarios)
            return f"{selected_scenario}. Use bright, engaging colors and show clear emotions like curiosity, satisfaction, and achievement."
        
        # Fallback to generic but engaging visual
        return ("Show someone on a journey of discovery and growth, moving from challenge to success with determination and creativity. "
               "Use visual metaphors that represent transformation and achievement. Include clear emotional progression from effort to satisfaction.")
    
    def _determine_educational_visual_metaphor(self, all_text: str, bullet_points: List[str]) -> str:
        """Determine the best visual metaphor for educational content"""
        
        # Process/workflow/steps content
        if any(word in all_text for word in ['process', 'step', 'workflow', 'sequence', 'first', 'then', 'next', 'finally']):
            return "A series of connected geometric shapes forming a path, with icons representing each step in the process. Simple arrows guide the flow from start to finish."
        
        # Comparison/contrast content
        elif any(word in all_text for word in ['compare', 'versus', 'vs', 'difference', 'before', 'after', 'traditional', 'modern']):
            return "A split composition with two distinct sides - one representing the old/traditional approach (muted colors, simple shapes) and the other showing the new/modern approach (bright colors, dynamic shapes)."
        
        # AI/Technology content
        elif any(word in all_text for word in ['ai', 'artificial', 'intelligence', 'machine learning', 'algorithm', 'neural', 'model']):
            return "A stylized brain made of interconnected nodes and circuits, with data flowing through neural pathways. Geometric patterns suggest computational processing."
        
        # Speed/efficiency content
        elif any(word in all_text for word in ['fast', 'quick', 'rapid', 'speed', 'efficient', 'streamline']):
            return "A minimalist rocket or arrow moving through simplified obstacles, leaving a trail of progress markers. Motion lines suggest speed and efficiency."
        
        # Building/creating content
        elif any(word in all_text for word in ['build', 'create', 'develop', 'construct', 'design', 'prototype']):
            return "Building blocks or puzzle pieces coming together to form a complete structure. Each piece represents a component of the solution."
        
        # Learning/education content
        elif any(word in all_text for word in ['learn', 'understand', 'master', 'discover', 'explore', 'knowledge']):
            return "A lightbulb with rays emanating outward, surrounded by simple icons representing different concepts being illuminated and understood."
        
        # Problem/solution content
        elif any(word in all_text for word in ['problem', 'solution', 'challenge', 'solve', 'fix', 'resolve']):
            return "A maze or tangled path on the left transforming into a clear, straight path on the right. A key or tool bridges the transformation."
        
        # Tools/resources content
        elif any(word in all_text for word in ['tool', 'resource', 'platform', 'framework', 'library', 'api']):
            return "A toolbox with simplified tool icons floating above it, each tool connected to its application with dotted lines."
        
        # Data/analysis content
        elif any(word in all_text for word in ['data', 'analysis', 'metrics', 'measure', 'statistics', 'visualization']):
            return "Abstract data points transforming into organized charts and graphs. Flowing lines connect raw data to meaningful insights."
        
        # Collaboration/team content
        elif any(word in all_text for word in ['team', 'collaborate', 'together', 'share', 'communicate', 'feedback']):
            return "Simplified human figures arranged in a circle, with connecting lines showing information flow. Speech bubbles or thought clouds overlap to show shared ideas."
        
        # Growth/improvement content
        elif any(word in all_text for word in ['grow', 'improve', 'enhance', 'evolve', 'progress', 'advance']):
            return "A plant or tree growing from seed to full bloom, with each stage clearly defined. Growth stages align with concept progression."
        
        # Testing/validation content
        elif any(word in all_text for word in ['test', 'validate', 'verify', 'check', 'quality', 'debug']):
            return "A magnifying glass examining different geometric shapes, with checkmarks appearing on validated items and X marks on issues to fix."
        
        # Default educational visual
        else:
            return "Interconnected circles representing concepts, with the largest circle in the center and smaller related concepts orbiting around it. Lines show relationships between ideas."
    
    def _get_fallback_visual_prompt(self) -> str:
        """Provide a generic educational visual prompt as fallback"""
        return """Create a visual concept for this course slide:

Image idea: Abstract geometric shapes arranged to suggest learning and progress. Use circles, triangles, and squares in a balanced composition with arrows showing flow and connection.

Style: clean flat illustration, ideal for a presentation."""
    
    def _create_prompt_text_file(self, drawing_prompt: str, slide_number: int) -> str:
        """Create a text file containing the drawing prompt for easy copy-paste into LLMs"""
        try:
            # Create unique filename
            import hashlib
            prompt_hash = hashlib.md5(drawing_prompt.encode()).hexdigest()[:8]
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
            random_id = random.randint(1000, 9999)
            filename = f"prompt_slide{slide_number:03d}_{prompt_hash}_{timestamp}_{random_id}.txt"
            
            # Create the text content
            text_content = f"""RICH VISUAL SUGGESTION FOR SLIDE {slide_number}
==============================================

{drawing_prompt}

==============================================
This creative visual suggestion is designed to help learners understand and remember the concepts through engaging imagery, storytelling, and memorable metaphors.

Copy the text above and paste it into any image generation AI (DALL-E, Midjourney, Stable Diffusion, etc.) or describe it to an illustrator to create a compelling visual that will animate your slide and captivate your audience.

Generated: {timestamp}
"""
            
            # Save text file
            os.makedirs('temp_sketches', exist_ok=True)
            file_path = os.path.join('temp_sketches', filename)
            
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(text_content)
            
            logger.info(f"Created drawing prompt text file: {file_path}")
            return file_path
            
        except Exception as e:
            logger.error(f"Error creating prompt text file: {e}")
            return None
    
    def _create_image_from_svg_description(self, svg_content: str, slide_number: int, diagram_type: str) -> str:
        """Create a simple diagram image based on SVG content description (fallback)"""
        try:
            # Create a basic image with SVG-inspired content
            width, height = 800, 600
            image = Image.new('RGB', (width, height), '#ffffff')
            draw = ImageDraw.Draw(image)
            
            # Add a title indicating this is an LLM-generated concept
            try:
                font = ImageFont.truetype('/System/Library/Fonts/Arial.ttf', 16)
                title_font = ImageFont.truetype('/System/Library/Fonts/Arial Bold.ttf', 20)
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
            
            # Add header
            draw.text((20, 20), f"🤖 LLM-Generated {diagram_type.replace('_', ' ').title()}", fill='#2c3e50', font=title_font)
            
            # Try to extract and display some text content from the SVG
            import re
            
            # Look for text elements in SVG
            text_matches = re.findall(r'<text[^>]*>([^<]+)</text>', svg_content, re.IGNORECASE)
            if text_matches:
                y_pos = 60
                draw.text((20, y_pos), "Key Elements:", fill='#7f8c8d', font=title_font)
                y_pos += 30
                
                for i, text in enumerate(text_matches[:8]):  # Max 8 elements
                    draw.text((30, y_pos), f"• {text.strip()}", fill='#2c3e50', font=font)
                    y_pos += 25
            
            # Look for rectangles and circles to show structure
            rect_matches = re.findall(r'<rect[^>]*>', svg_content, re.IGNORECASE)
            circle_matches = re.findall(r'<circle[^>]*>', svg_content, re.IGNORECASE)
            
            if rect_matches or circle_matches:
                y_pos = height - 150
                draw.text((20, y_pos), f"Structure: {len(rect_matches)} boxes, {len(circle_matches)} circles", 
                         fill='#7f8c8d', font=font)
            
            # Add a visual indicator that this is AI-generated
            draw.rectangle([width-200, height-100, width-20, height-20], 
                          fill='#e3f2fd', outline='#2196f3', width=2)
            draw.text((width-180, height-80), "Generated by AI", fill='#1976d2', font=font)
            draw.text((width-180, height-60), "from slide content", fill='#1976d2', font=font)
            
            # Create unique filename
            content_hash = abs(hash(svg_content)) % (10**16)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
            random_id = random.randint(1000, 9999)
            filename = f"sketch_slide{slide_number:03d}_llm_{diagram_type}_{content_hash}_{timestamp}_{random_id}.png"
            
            # Save the image
            os.makedirs('temp_sketches', exist_ok=True)
            image_path = os.path.join('temp_sketches', filename)
            image.save(image_path, 'PNG')
            
            logger.info(f"Created fallback LLM diagram with {len(text_matches)} text elements: {image_path}")
            return image_path
            
        except Exception as e:
            logger.error(f"Error creating fallback diagram: {e}")
            return None

    def _determine_content_specific_sketch_type(self, bullet_points: List[str], title: str) -> str:
        """Analyze bullet points and title to determine the best sketch type"""
        all_text = f"{title} {' '.join(bullet_points)}".lower()
        
        # Process/workflow indicators
        if any(word in all_text for word in ['process', 'step', 'workflow', 'procedure', 'method', 'approach', 'how to']):
            return 'process'
            
        # Comparison indicators
        if any(word in all_text for word in ['vs', 'versus', 'compare', 'comparison', 'difference', 'traditional', 'modern']):
            return 'comparison'
            
        # AI/Technology indicators
        if any(word in all_text for word in ['ai', 'genai', 'artificial intelligence', 'machine learning', 'algorithm', 'model']):
            return 'ai_tech'
            
        # UI/Design indicators  
        if any(word in all_text for word in ['interface', 'design', 'user', 'ui', 'ux', 'mockup', 'prototype']):
            return 'ui_design'
            
        # Data/Metrics indicators
        if any(word in all_text for word in ['data', 'metrics', 'analytics', 'chart', 'graph', 'statistics', 'performance']):
            return 'data_metrics'
            
        # System/Architecture indicators
        if any(word in all_text for word in ['system', 'architecture', 'component', 'service', 'api', 'database']):
            return 'system_arch'
            
        # Learning/Education indicators
        if any(word in all_text for word in ['learn', 'course', 'lesson', 'tutorial', 'training', 'education', 'skill']):
            return 'learning_path'
            
        # Benefits/Advantages indicators
        if any(word in all_text for word in ['benefit', 'advantage', 'improve', 'better', 'faster', 'efficient']):
            return 'benefits'
            
        # Default to concept map
        return 'concept_map'
    
    def _draw_process_from_content(self, draw, width, height, font, bullet_points, title):
        """Draw a process flow based on actual bullet point content"""
        num_steps = min(len(bullet_points), 5)  # Max 5 steps for readability
        if num_steps == 0:
            return
            
        # Calculate step positions
        step_width = 120
        step_height = 60
        total_width = num_steps * step_width + (num_steps - 1) * 40
        start_x = (width - total_width) // 2
        y_pos = height // 2 - step_height // 2
        
        colors = ['#3498db', '#2ecc71', '#f39c12', '#e74c3c', '#9b59b6']
        
        for i, bullet in enumerate(bullet_points[:num_steps]):
            x_pos = start_x + i * (step_width + 40)
            
            # Draw step box with shadow (ensure valid dimensions)
            shadow_offset = 3
            if step_width > 0 and step_height > 0:
                draw.rounded_rectangle([x_pos + shadow_offset, y_pos + shadow_offset, 
                                      x_pos + step_width + shadow_offset, y_pos + step_height + shadow_offset],
                                     radius=8, fill='#00000020')
                draw.rounded_rectangle([x_pos, y_pos, x_pos + step_width, y_pos + step_height],
                                     radius=8, fill=colors[i % len(colors)], outline='#2c3e50', width=2)
            
            # Add step number
            step_num = f"{i+1}"
            draw.text((x_pos + 10, y_pos + 10), step_num, fill='white', font=font)
            
            # Add abbreviated bullet text
            text_lines = self._wrap_text(bullet[:30] + "...", 14, font)
            for j, line in enumerate(text_lines[:2]):  # Max 2 lines
                draw.text((x_pos + 10, y_pos + 25 + j * 15), line, fill='white', font=font)
            
            # Draw arrow to next step
            if i < num_steps - 1:
                arrow_start_x = x_pos + step_width + 5
                arrow_end_x = arrow_start_x + 30
                arrow_y = y_pos + step_height // 2
                
                # Arrow line
                draw.line([arrow_start_x, arrow_y, arrow_end_x - 8, arrow_y], fill='#2c3e50', width=3)
                # Arrow head
                draw.polygon([(arrow_end_x, arrow_y), (arrow_end_x - 8, arrow_y - 4), (arrow_end_x - 8, arrow_y + 4)], 
                           fill='#2c3e50')
    
    def _draw_comparison_from_content(self, draw, width, height, font, bullet_points, title):
        """Draw a comparison diagram based on bullet points"""
        if len(bullet_points) < 2:
            return
            
        # Split bullets into two groups
        mid_point = len(bullet_points) // 2
        left_items = bullet_points[:mid_point] if mid_point > 0 else bullet_points[:1]
        right_items = bullet_points[mid_point:] if mid_point > 0 else bullet_points[1:]
        
        # Draw two columns
        col_width = width // 2 - 60
        left_x = 30
        right_x = width // 2 + 30
        
        # Left column (Traditional/Old)
        left_rect = [left_x, 100, left_x + max(col_width, 50), height - 100]
        if left_rect[3] > left_rect[1]:  # Ensure y2 > y1
            draw.rounded_rectangle(left_rect, radius=10, fill='#ecf0f1', outline='#bdc3c7', width=2)
        draw.text((left_x + 10, 110), "Before/Traditional", fill='#2c3e50', font=font)
        
        for i, item in enumerate(left_items[:4]):  # Max 4 items
            y_pos = 140 + i * 50
            # Bullet point
            draw.ellipse([left_x + 15, y_pos, left_x + 25, y_pos + 10], fill='#e74c3c')
            # Text
            text_lines = self._wrap_text(item[:40] + "...", col_width - 50, font)
            draw.text((left_x + 35, y_pos), text_lines[0] if text_lines else "", fill='#2c3e50', font=font)
        
        # Right column (Modern/New)
        right_rect = [right_x, 100, right_x + max(col_width, 50), height - 100]
        if right_rect[3] > right_rect[1]:  # Ensure y2 > y1
            draw.rounded_rectangle(right_rect, radius=10, fill='#e8f5e8', outline='#27ae60', width=2)
        draw.text((right_x + 10, 110), "After/Modern", fill='#27ae60', font=font)
        
        for i, item in enumerate(right_items[:4]):  # Max 4 items
            y_pos = 140 + i * 50
            # Bullet point
            draw.ellipse([right_x + 15, y_pos, right_x + 25, y_pos + 10], fill='#27ae60')
            # Text
            text_lines = self._wrap_text(item[:40] + "...", col_width - 50, font)
            draw.text((right_x + 35, y_pos), text_lines[0] if text_lines else "", fill='#27ae60', font=font)
        
        # Arrow between columns
        arrow_y = height // 2
        draw.line([width // 2 - 20, arrow_y, width // 2 + 20, arrow_y], fill='#3498db', width=4)
        draw.polygon([(width // 2 + 20, arrow_y), (width // 2 + 12, arrow_y - 6), (width // 2 + 12, arrow_y + 6)], 
                   fill='#3498db')
    
    def _draw_benefits_diagram(self, draw, width, height, font, bullet_points):
        """Draw a benefits/advantages diagram"""
        if not bullet_points:
            return
            
        # Central hub
        center_x, center_y = width // 2, height // 2
        hub_radius = 60
        
        # Draw central hub
        draw.ellipse([center_x - hub_radius, center_y - hub_radius, 
                     center_x + hub_radius, center_y + hub_radius],
                    fill='#3498db', outline='#2980b9', width=3)
        draw.text((center_x - 30, center_y - 10), "Benefits", fill='white', font=font)
        
        # Draw benefit nodes around the hub
        num_benefits = min(len(bullet_points), 6)  # Max 6 for readability
        angle_step = 2 * 3.14159 / num_benefits
        
        for i, benefit in enumerate(bullet_points[:num_benefits]):
            angle = i * angle_step
            node_x = center_x + int(140 * cos(angle))
            node_y = center_y + int(140 * sin(angle))
            
            # Draw connecting line
            draw.line([center_x + int(hub_radius * cos(angle)), 
                      center_y + int(hub_radius * sin(angle)),
                      node_x - int(30 * cos(angle)), 
                      node_y - int(30 * sin(angle))], 
                     fill='#34495e', width=2)
            
            # Draw benefit node
            node_width, node_height = 120, 40
            draw.rounded_rectangle([node_x - node_width//2, node_y - node_height//2,
                                  node_x + node_width//2, node_y + node_height//2],
                                 radius=8, fill='#2ecc71', outline='#27ae60', width=2)
            
            # Add benefit text
            benefit_text = benefit[:25] + "..." if len(benefit) > 25 else benefit
            text_lines = self._wrap_text(benefit_text, node_width - 10, font)
            for j, line in enumerate(text_lines[:2]):
                draw.text((node_x - len(line) * 4, node_y - 8 + j * 12), line, fill='white', font=font)
    
    def _draw_concept_map_from_bullets(self, draw, width, height, font, bullet_points, title):
        """Draw a concept map directly from bullet points"""
        if not bullet_points:
            return
            
        # Main concept in center
        center_x, center_y = width // 2, height // 2 + 20
        main_width, main_height = 160, 50
        
        draw.rounded_rectangle([center_x - main_width//2, center_y - main_height//2,
                              center_x + main_width//2, center_y + main_height//2],
                             radius=10, fill='#3498db', outline='#2980b9', width=3)
        
        # Truncate title for center
        center_title = title[:20] + "..." if len(title) > 20 else title
        title_lines = self._wrap_text(center_title, main_width - 10, font)
        for i, line in enumerate(title_lines[:2]):
            draw.text((center_x - len(line) * 4, center_y - 10 + i * 12), line, fill='white', font=font)
        
        # Surrounding concepts from bullet points
        num_concepts = min(len(bullet_points), 5)
        colors = ['#e74c3c', '#2ecc71', '#f39c12', '#9b59b6', '#1abc9c']
        
        for i, bullet in enumerate(bullet_points[:num_concepts]):
            angle = (i * 2 * 3.14159) / num_concepts - 3.14159/2  # Start from top
            radius = 120
            
            concept_x = center_x + int(radius * cos(angle))
            concept_y = center_y + int(radius * sin(angle))
            
            # Connection line
            draw.line([center_x + int((main_width//2) * cos(angle)),
                      center_y + int((main_height//2) * sin(angle)),
                      concept_x - int(40 * cos(angle)),
                      concept_y - int(20 * sin(angle))],
                     fill='#7f8c8d', width=2)
            
            # Concept box
            concept_width, concept_height = 100, 40
            draw.rounded_rectangle([concept_x - concept_width//2, concept_y - concept_height//2,
                                  concept_x + concept_width//2, concept_y + concept_height//2],
                                 radius=8, fill=colors[i], outline='#2c3e50', width=2)
            
            # Concept text (first few words)
            concept_text = bullet.split()[:3]  # First 3 words
            concept_text = ' '.join(concept_text)[:15] + "..." if len(' '.join(concept_text)) > 15 else ' '.join(concept_text)
            
            text_lines = self._wrap_text(concept_text, concept_width - 10, font)
            for j, line in enumerate(text_lines[:2]):
                draw.text((concept_x - len(line) * 3, concept_y - 8 + j * 12), line, fill='white', font=font)

    def _draw_ui_design_from_content(self, draw, width, height, font, bullet_points, title):
        """Draw a UI mockup based on bullet points"""
        if not bullet_points:
            return
            
        # Draw a simple mockup with header, content areas, and buttons based on content
        # Header area
        header_height = 60
        if header_height > 0:
            draw.rounded_rectangle([50, 50, width - 50, 50 + header_height],
                                 radius=8, fill='#3498db', outline='#2980b9', width=2)
            draw.text((60, 70), title[:30] + "..." if len(title) > 30 else title, fill='white', font=font)
        
        # Main content areas based on bullet points
        content_start_y = 120
        content_height = (height - content_start_y - 80) // min(len(bullet_points), 4)
        
        for i, bullet in enumerate(bullet_points[:4]):
            y_pos = content_start_y + i * (content_height + 10)
            
            # Content box
            if content_height > 20:  # Ensure valid dimensions
                draw.rounded_rectangle([60, y_pos, width - 60, y_pos + max(content_height - 10, 20)],
                                     radius=5, fill='#ecf0f1', outline='#bdc3c7', width=1)
                
                # Content text
                content_text = bullet[:25] + "..." if len(bullet) > 25 else bullet
                draw.text((70, y_pos + 10), content_text, fill='#2c3e50', font=font)
        
        # Action buttons at bottom
        button_y = height - 60
        button_width = 100
        if button_y > content_start_y + 50:  # Ensure buttons don't overlap content
            # Primary button
            draw.rounded_rectangle([60, button_y, 60 + button_width, button_y + 30],
                                 radius=5, fill='#27ae60', outline='#229954', width=1)
            draw.text((75, button_y + 8), "Submit", fill='white', font=font)
            
            # Secondary button
            draw.rounded_rectangle([180, button_y, 180 + button_width, button_y + 30],
                                 radius=5, fill='#95a5a6', outline='#7f8c8d', width=1)
            draw.text((205, button_y + 8), "Cancel", fill='white', font=font)

    def _wrap_text(self, text: str, max_width: int, font) -> List[str]:
        """Wrap text to fit within specified width"""
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            # Test if adding this word would exceed the width
            test_line = ' '.join(current_line + [word])
            try:
                bbox = ImageDraw.Draw(Image.new('RGB', (1, 1))).textbbox((0, 0), test_line, font=font)
                text_width = bbox[2] - bbox[0]
            except:
                # Fallback for simple width estimation
                text_width = len(test_line) * 8
            
            if text_width <= max_width or not current_line:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return lines if lines else ['']

    def _create_sketch_image(self, bullet_points: List[str], title: str, suggestion: str, slide_number: int = 0) -> str:
        """Create a high-quality, professional sketch based on content and suggestion"""
        try:
            # Create high-resolution image for better quality
            width, height = 800, 600
            image = Image.new('RGB', (width, height), '#f8f9fa')  # Light gray background
            draw = ImageDraw.Draw(image)
            
            # Try to load better fonts with multiple fallbacks
            try:
                # Try to find the best available font
                font_paths = [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/System/Library/Fonts/Arial.ttf", 
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/Windows/Fonts/arial.ttf"
                ]
                font_loaded = False
                for font_path in font_paths:
                    try:
                        font = ImageFont.truetype(font_path, 18)
                        title_font = ImageFont.truetype(font_path, 24)
                        small_font = ImageFont.truetype(font_path, 14)
                        large_font = ImageFont.truetype(font_path, 20)
                        font_loaded = True
                        break
                    except:
                        continue
                
                if not font_loaded:
                    raise Exception("No fonts found")
                    
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
                large_font = ImageFont.load_default()
            
            # Combine content for analysis - include suggestion text for better matching
            all_content = f"{title} {' '.join(bullet_points)} {suggestion}".lower()
            
            # More specific keyword matching based on suggestion text
            sketch_type = self._determine_sketch_type(all_content, suggestion)
            
            # Draw sketch based on determined type
            if sketch_type == 'flowchart':
                self._draw_process_flowchart(draw, width, height, font, title_font)
            elif sketch_type == 'ai':
                self._draw_ai_concept(draw, width, height, font, title_font)
            elif sketch_type == 'comparison':
                self._draw_comparison_diagram(draw, width, height, font, title_font)
            elif sketch_type == 'data':
                self._draw_data_visualization(draw, width, height, font, title_font)
            elif sketch_type == 'ui':
                self._draw_ui_mockup(draw, width, height, font, title_font)
            elif sketch_type == 'system':
                self._draw_system_architecture(draw, width, height, font, title_font)
            elif sketch_type == 'learning':
                self._draw_learning_path(draw, width, height, font, title_font)
            elif sketch_type == 'collaboration':
                self._draw_collaboration_diagram(draw, width, height, font, title_font)
            elif sketch_type == 'problem_solution':
                self._draw_problem_solution(draw, width, height, font, title_font)
            elif sketch_type == 'timeline':
                self._draw_timeline(draw, width, height, font, title_font)
            elif sketch_type == 'business':
                self._draw_business_diagram(draw, width, height, font, title_font)
            else:
                # Generic concept diagram
                self._draw_generic_concept(draw, width, height, font, title_font, bullet_points)
            
            # Add a professional title with modern styling
            sketch_title = self._get_sketch_title(all_content, sketch_type)
            bbox = draw.textbbox((0, 0), sketch_title, font=title_font)
            title_width = bbox[2] - bbox[0]
            title_height = bbox[3] - bbox[1]
            
            # Title background with rounded corners
            title_x = width//2 - title_width//2 - 15
            title_y = 15
            draw.rounded_rectangle([title_x, title_y, title_x + title_width + 30, title_y + title_height + 15], 
                                 radius=8, fill='#37474F', outline='#263238', width=2)
            
            # Title text with white color for contrast
            draw.text((width//2 - title_width//2, 22), sketch_title, fill='white', font=title_font)
            
            # Save image to temporary file with unique naming
            os.makedirs('temp_sketches', exist_ok=True)
            # Use slide number, timestamp, and random component to ensure every slide gets a unique image
            import random
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            random_component = random.randint(1000, 9999)
            unique_id = abs(hash(f"{slide_number}_{title}_{suggestion}_{str(bullet_points)}_{sketch_type}_{timestamp}_{random_component}"))
            image_path = f'temp_sketches/sketch_slide{slide_number:03d}_{sketch_type}_{unique_id}_{timestamp}_{random_component}.png'
            image.save(image_path)
            
            logger.info(f"Created {sketch_type} sketch for slide {slide_number}: {image_path}")
            return image_path
            
        except Exception as e:
            logger.error(f"Error creating sketch image: {e}")
            return None
    
    def _determine_sketch_type(self, all_content: str, suggestion: str) -> str:
        """Determine the most appropriate sketch type based on content and suggestion"""
        suggestion_lower = suggestion.lower()
        
        # Check suggestion text first for most accurate matching
        if 'flowchart' in suggestion_lower or 'step-by-step' in suggestion_lower or 'arrows connecting' in suggestion_lower:
            return 'flowchart'
        elif 'side-by-side' in suggestion_lower or 'comparison' in suggestion_lower or 'two columns' in suggestion_lower:
            return 'comparison'
        elif 'brain' in suggestion_lower or 'neural network' in suggestion_lower or 'ai concept' in suggestion_lower:
            return 'ai'
        elif 'charts and graphs' in suggestion_lower or 'bar chart' in suggestion_lower or 'dashboard' in suggestion_lower:
            return 'data'
        elif 'wireframe' in suggestion_lower or 'buttons' in suggestion_lower or 'user interface' in suggestion_lower:
            return 'ui'
        elif 'system architecture' in suggestion_lower or 'components' in suggestion_lower or 'data flow' in suggestion_lower:
            return 'system'
        elif 'learning path' in suggestion_lower or 'milestone markers' in suggestion_lower or 'progression' in suggestion_lower:
            return 'learning'
        elif 'collaboration' in suggestion_lower or 'meeting table' in suggestion_lower or 'team members' in suggestion_lower:
            return 'collaboration'
        elif 'problem-solution' in suggestion_lower or 'warning icons' in suggestion_lower or 'checkmarks' in suggestion_lower:
            return 'problem_solution'
        elif 'timeline' in suggestion_lower or 'milestone flags' in suggestion_lower or 'progression' in suggestion_lower:
            return 'timeline'
        elif 'business diagrams' in suggestion_lower or 'growth arrows' in suggestion_lower or 'market segments' in suggestion_lower:
            return 'business'
        
        # Fallback to content analysis
        if any(keyword in all_content for keyword in ['process', 'workflow', 'steps', 'cycle', 'development', 'build', 'deploy']):
            return 'flowchart'
        elif any(keyword in all_content for keyword in ['ai', 'genai', 'chatbot', 'gpt', 'model', 'prompt']):
            return 'ai'
        elif any(keyword in all_content for keyword in ['compare', 'vs', 'versus', 'difference', 'before', 'after']):
            return 'comparison'
        elif any(keyword in all_content for keyword in ['data', 'analytics', 'chart', 'graph', 'metrics', 'visualization']):
            return 'data'
        elif any(keyword in all_content for keyword in ['interface', 'ui', 'app', 'website', 'screen', 'button']):
            return 'ui'
        elif any(keyword in all_content for keyword in ['system', 'architecture', 'components', 'api', 'database']):
            return 'system'
        elif any(keyword in all_content for keyword in ['learn', 'course', 'lesson', 'education', 'skill', 'training']):
            return 'learning'
        elif any(keyword in all_content for keyword in ['team', 'collaboration', 'meeting', 'communication']):
            return 'collaboration'
        elif any(keyword in all_content for keyword in ['problem', 'solution', 'issue', 'challenge', 'fix']):
            return 'problem_solution'
        elif any(keyword in all_content for keyword in ['timeline', 'schedule', 'roadmap', 'phases', 'milestone']):
            return 'timeline'
        elif any(keyword in all_content for keyword in ['business', 'strategy', 'market', 'growth', 'revenue']):
            return 'business'
        else:
            return 'generic'
    
    def _get_sketch_title(self, content: str, sketch_type: str) -> str:
        """Generate appropriate title for the sketch"""
        titles = {
            'flowchart': 'Process Flow',
            'ai': 'AI System',
            'comparison': 'Comparison',
            'data': 'Data Analysis',
            'ui': 'User Interface',
            'system': 'System Architecture',
            'learning': 'Learning Path',
            'collaboration': 'Team Collaboration',
            'problem_solution': 'Problem → Solution',
            'timeline': 'Timeline',
            'business': 'Business Strategy',
            'generic': 'Concept Overview'
        }
        return titles.get(sketch_type, 'Diagram')
    
    def _draw_process_flowchart(self, draw, width, height, font, title_font):
        """Draw a professional process flowchart with modern styling"""
        steps = ["Start", "Analyze", "Process", "Review", "Complete"]
        box_width, box_height = 130, 60
        spacing_x = 150
        start_x = (width - (len(steps) * box_width + (len(steps)-1) * 20)) // 2
        y = height // 2 - box_height // 2
        
        # Color scheme for different step types
        colors = {
            'start': '#4CAF50',      # Green
            'process': '#2196F3',    # Blue  
            'review': '#FF9800',     # Orange
            'end': '#9C27B0'         # Purple
        }
        
        for i, step in enumerate(steps):
            if i == 0:  # Vertical layout for better fit
                x = width // 2 - box_width // 2
                current_y = 80 + i * 90
            else:
                x = width // 2 - box_width // 2
                current_y = 80 + i * 90
            
            # Choose color based on step type
            if i == 0:
                color = colors['start']
            elif i == len(steps) - 1:
                color = colors['end']
            elif 'review' in step.lower() or 'analyze' in step.lower():
                color = colors['review']
            else:
                color = colors['process']
            
            # Draw shadow for depth
            shadow_offset = 4
            draw.rounded_rectangle([x + shadow_offset, current_y + shadow_offset, 
                                  x + box_width + shadow_offset, current_y + box_height + shadow_offset], 
                                 radius=12, fill='#00000030')
            
            # Draw main box with gradient-like effect
            draw.rounded_rectangle([x, current_y, x + box_width, current_y + box_height], 
                                 radius=12, fill=color, outline='#333333', width=3)
            
            # Add inner highlight for 3D effect
            draw.rounded_rectangle([x + 3, current_y + 3, x + box_width - 3, current_y + box_height - 3], 
                                 radius=9, outline='#ffffff40', width=1)
            
            # Add step text with better positioning
            bbox = draw.textbbox((0, 0), step, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            text_x = x + box_width//2 - text_width//2
            text_y = current_y + box_height//2 - text_height//2
            
            # Add text shadow for better readability
            draw.text((text_x + 1, text_y + 1), step, fill='#00000080', font=font)
            draw.text((text_x, text_y), step, fill='white', font=font)
            
            # Draw modern arrow to next step
            if i < len(steps) - 1:
                arrow_start_y = current_y + box_height + 5
                arrow_end_y = 80 + (i+1) * 90 - 5
                arrow_x = width // 2
                
                # Draw arrow line with gradient effect
                draw.line([arrow_x, arrow_start_y, arrow_x, arrow_end_y], 
                         fill='#666666', width=4)
                
                # Draw modern arrow head
                arrow_size = 10
                draw.polygon([arrow_x - arrow_size, arrow_end_y - arrow_size,
                             arrow_x + arrow_size, arrow_end_y - arrow_size,
                             arrow_x, arrow_end_y], 
                            fill='#666666')
    
    def _draw_ai_concept(self, draw, width, height, font, title_font):
        """Draw a modern AI brain/neural network concept"""
        center_x, center_y = width // 2, height // 2
        
        # Draw stylized brain with gradient effect
        brain_width, brain_height = 200, 140
        
        # Brain shadow for depth
        shadow_offset = 6
        draw.ellipse([center_x - brain_width//2 + shadow_offset, center_y - brain_height//2 + shadow_offset,
                     center_x + brain_width//2 + shadow_offset, center_y + brain_height//2 + shadow_offset], 
                    fill='#00000030')
        
        # Main brain with layered gradient effect
        colors = ['#E8F5E8', '#C8E6C9', '#A5D6A7', '#81C784']
        for i, color in enumerate(colors):
            offset = i * 10
            draw.ellipse([center_x - brain_width//2 + offset, center_y - brain_height//2 + offset,
                         center_x + brain_width//2 - offset, center_y + brain_height//2 - offset], 
                        fill=color, outline='#4CAF50', width=3)
        
        # Draw professional neural network layers
        layer_positions = [
            [(center_x-80, center_y-50), (center_x-80, center_y-15), (center_x-80, center_y+15), (center_x-80, center_y+50)],
            [(center_x-30, center_y-60), (center_x-30, center_y-20), (center_x-30, center_y+20), (center_x-30, center_y+60)],
            [(center_x+20, center_y-40), (center_x+20, center_y-10), (center_x+20, center_y+25), (center_x+20, center_y+55)],
            [(center_x+70, center_y-25), (center_x+70, center_y+25)]
        ]
        
        # Draw connections with varying thickness for importance
        for i in range(len(layer_positions)-1):
            for j, node1 in enumerate(layer_positions[i]):
                for k, node2 in enumerate(layer_positions[i+1]):
                    # Vary connection strength visually
                    thickness = 2 if (j + k) % 2 == 0 else 1
                    alpha = 'AA' if thickness == 2 else '66'
                    draw.line([node1[0], node1[1], node2[0], node2[1]], 
                             fill=f'#2196F3{alpha}', width=thickness)
        
        # Draw nodes with modern styling
        for layer_idx, layer in enumerate(layer_positions):
            node_color = ['#4CAF50', '#2196F3', '#FF9800', '#9C27B0'][layer_idx % 4]
            for x, y in layer:
                # Node shadow
                draw.ellipse([x-7, y-7, x+7, y+7], fill='#00000020')
                # Main node
                draw.ellipse([x-6, y-6, x+6, y+6], fill=node_color, outline='#333333', width=2)
                # Inner highlight
                draw.ellipse([x-3, y-3, x+3, y+3], fill='#ffffff60')
        
        # Add professional labels with backgrounds
        # Input label
        input_text = "INPUT"
        bbox = draw.textbbox((0, 0), input_text, font=font)
        text_width = bbox[2] - bbox[0]
        draw.rounded_rectangle([center_x-80-text_width//2-5, center_y-75, 
                              center_x-80+text_width//2+5, center_y-60], 
                             radius=5, fill='#4CAF50', outline='#333333', width=1)
        draw.text((center_x-80-text_width//2, center_y-72), input_text, fill='white', font=font)
        
        # Output label  
        output_text = "OUTPUT"
        bbox = draw.textbbox((0, 0), output_text, font=font)
        text_width = bbox[2] - bbox[0]
        draw.rounded_rectangle([center_x+70-text_width//2-5, center_y-40, 
                              center_x+70+text_width//2+5, center_y-25], 
                             radius=5, fill='#9C27B0', outline='#333333', width=1)
        draw.text((center_x+70-text_width//2, center_y-37), output_text, fill='white', font=font)
        
        # Main AI label with background
        ai_text = "NEURAL NETWORK"
        bbox = draw.textbbox((0, 0), ai_text, font=title_font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        draw.rounded_rectangle([center_x-text_width//2-10, center_y+80, 
                              center_x+text_width//2+10, center_y+80+text_height+10], 
                             radius=8, fill='#2196F3', outline='#0D47A1', width=2)
        draw.text((center_x-text_width//2, center_y+85), ai_text, fill='white', font=title_font)
    
    def _draw_comparison_diagram(self, draw, width, height, font, title_font):
        """Draw side-by-side comparison diagram"""
        center_x = width // 2
        y_start = 80
        
        # Draw dividing line
        draw.line([center_x, y_start, center_x, height-40], fill='black', width=2)
        
        # Left side - "Before"
        left_x = center_x // 2
        draw.text((left_x-20, y_start-20), "Before", fill='black', font=title_font)
        
        # Draw negative indicators
        draw.rectangle([left_x-40, y_start+20, left_x+40, y_start+50], outline='red', width=2)
        draw.text((left_x-25, y_start+30), "Problem", fill='black', font=font)
        
        draw.rectangle([left_x-40, y_start+70, left_x+40, y_start+100], outline='red', width=2)
        draw.text((left_x-20, y_start+80), "Issues", fill='black', font=font)
        
        # Right side - "After"
        right_x = center_x + center_x // 2
        draw.text((right_x-15, y_start-20), "After", fill='black', font=title_font)
        
        # Draw positive indicators
        draw.rectangle([right_x-40, y_start+20, right_x+40, y_start+50], outline='green', width=2)
        draw.text((right_x-25, y_start+30), "Solution", fill='black', font=font)
        
        draw.rectangle([right_x-40, y_start+70, right_x+40, y_start+100], outline='green', width=2)
        draw.text((right_x-25, y_start+80), "Benefits", fill='black', font=font)
        
        # Add arrow
        arrow_y = y_start + 60
        draw.polygon([center_x-20, arrow_y, center_x+20, arrow_y, center_x+15, arrow_y-5, center_x+15, arrow_y+5], fill='black')
    
    def _draw_timeline(self, draw, width, height, font, title_font):
        """Draw timeline with milestones"""
        y = height // 2 + 20
        start_x = 60
        end_x = width - 60
        
        # Draw main timeline
        draw.line([start_x, y, end_x, y], fill='black', width=4)
        
        # Draw milestones
        milestones = ["Start", "Phase 1", "Phase 2", "Complete"]
        milestone_x = [start_x + i * (end_x - start_x) // 3 for i in range(4)]
        
        for i, (x, milestone) in enumerate(zip(milestone_x, milestones)):
            # Draw milestone marker
            draw.ellipse([x-8, y-8, x+8, y+8], fill='black', outline='black')
            draw.ellipse([x-5, y-5, x+5, y+5], fill='white', outline='black')
            
            # Add milestone text
            bbox = draw.textbbox((0, 0), milestone, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((x - text_width//2, y + 20), milestone, fill='black', font=font)
            
            # Add date placeholders
            date = f"Week {i+1}"
            bbox = draw.textbbox((0, 0), date, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((x - text_width//2, y - 25), date, fill='gray', font=font)
    
    def _draw_business_diagram(self, draw, width, height, font, title_font):
        """Draw business strategy diagram"""
        center_x, center_y = width // 2, height // 2 + 20
        
        # Draw growth arrow
        arrow_start_x = 80
        arrow_end_x = width - 80
        arrow_y = center_y
        
        draw.line([arrow_start_x, arrow_y, arrow_end_x, arrow_y], fill='black', width=4)
        draw.polygon([arrow_end_x-15, arrow_y-8, arrow_end_x-15, arrow_y+8, arrow_end_x, arrow_y], fill='black')
        
        # Add growth indicators
        growth_points = [(120, center_y-20), (200, center_y-35), (280, center_y-50), (360, center_y-65)]
        
        for i, (x, y) in enumerate(growth_points):
            if x < width - 100:
                # Draw upward bar
                bar_height = 20 + i * 15
                draw.rectangle([x-10, y, x+10, y+bar_height], fill='green', outline='black')
                
                # Add percentage
                draw.text((x-15, y-15), f"{(i+1)*25}%", fill='black', font=font)
        
        # Add labels
        draw.text((arrow_start_x-30, arrow_y+30), "Growth", fill='black', font=title_font)
        draw.text((arrow_end_x-40, arrow_y+30), "Success", fill='black', font=title_font)
    
    def _draw_data_visualization(self, draw, width, height, font, title_font):
        """Draw simple charts and graphs"""
        # Bar chart with better proportions
        bar_x = 60
        bar_y = height - 100
        bar_heights = [30, 50, 40, 60, 35]
        bar_colors = ['black', 'gray', 'black', 'gray', 'black']
        
        # Draw bars
        for i, h in enumerate(bar_heights):
            x = bar_x + i * 35
            draw.rectangle([x, bar_y - h, x + 25, bar_y], fill=bar_colors[i], outline='black', width=2)
            draw.text((x + 8, bar_y + 5), f"Q{i+1}", fill='black', font=font)
        
        # Draw axes
        draw.line([bar_x-10, bar_y, bar_x + len(bar_heights)*35, bar_y], fill='black', width=2)  # X-axis
        draw.line([bar_x-10, bar_y, bar_x-10, bar_y-80], fill='black', width=2)  # Y-axis
        
        # Simple line graph on the right
        line_start_x = 280
        line_points = [(line_start_x, 200), (line_start_x+40, 180), (line_start_x+80, 160), (line_start_x+120, 140)]
        
        # Draw line graph
        for i in range(len(line_points) - 1):
            draw.line([line_points[i], line_points[i+1]], fill='black', width=3)
        
        # Draw points with values
        for i, point in enumerate(line_points):
            draw.ellipse([point[0]-4, point[1]-4, point[0]+4, point[1]+4], fill='white', outline='black', width=2)
            draw.text((point[0]-5, point[1]-20), f"{100-i*10}", fill='black', font=font)
        
        # Add chart labels
        draw.text((bar_x, 50), "Sales Data", fill='black', font=title_font)
        draw.text((line_start_x, 50), "Growth Trend", fill='black', font=title_font)
    
    def _draw_ui_mockup(self, draw, width, height, font, title_font):
        """Draw a modern, professional UI wireframe"""
        # Main container with shadow and modern styling
        container_shadow = 4
        draw.rounded_rectangle([50 + container_shadow, 80 + container_shadow, 
                              width-50 + container_shadow, height-60 + container_shadow], 
                             radius=12, fill='#00000020')
        
        # Main container
        draw.rounded_rectangle([50, 80, width-50, height-60], 
                             radius=12, fill='white', outline='#E0E0E0', width=3)
        
        # Modern header with gradient-like effect
        draw.rounded_rectangle([60, 90, width-60, 130], 
                             radius=8, fill='#2196F3', outline='#1976D2', width=2)
        draw.rounded_rectangle([62, 92, width-62, 128], 
                             radius=6, outline='#64B5F6', width=1)
        
        # Header text
        header_text = "App Header / Navigation"
        bbox = draw.textbbox((0, 0), header_text, font=font)
        text_width = bbox[2] - bbox[0]
        draw.text((70, 105), header_text, fill='white', font=font)
        
        # Menu icon (hamburger)
        menu_x = width - 100
        for i in range(3):
            draw.rectangle([menu_x, 100 + i*6, menu_x + 20, 102 + i*6], fill='white')
        
        # Modern buttons with different styles
        button_colors = ['#4CAF50', '#FF9800', '#9C27B0']
        button_texts = ['Primary', 'Secondary', 'Action']
        
        for i, (color, text) in enumerate(zip(button_colors, button_texts)):
            x = 70 + i * 120
            y = 150
            
            # Button shadow
            draw.rounded_rectangle([x + 2, y + 2, x + 102, y + 37], 
                                 radius=8, fill='#00000030')
            # Main button
            draw.rounded_rectangle([x, y, x + 100, y + 35], 
                                 radius=8, fill=color, outline='#333333', width=2)
            # Button highlight
            draw.rounded_rectangle([x + 2, y + 2, x + 98, y + 15], 
                                 radius=6, fill='#ffffff40')
            
            # Button text
            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((x + 50 - text_width//2, y + 12), text, fill='white', font=font)
        
        # Modern form fields
        field_labels = ['Email Address', 'Password', 'Confirm Password']
        for i, label in enumerate(field_labels):
            y = 210 + i * 50
            
            # Field label
            draw.text((70, y), label, fill='#333333', font=font)
            
            # Field container with modern styling
            draw.rounded_rectangle([70, y + 18, 350, y + 43], 
                                 radius=6, fill='white', outline='#BDBDBD', width=2)
            draw.rounded_rectangle([72, y + 20, 348, y + 41], 
                                 radius=4, outline='#E0E0E0', width=1)
            
            # Placeholder text
            placeholder = "Enter " + label.lower()
            draw.text((80, y + 26), placeholder, fill='#9E9E9E', font=font)
        
        # Modern mobile frame
        mobile_x = width - 150
        mobile_y = 140
        
        # Phone shadow
        draw.rounded_rectangle([mobile_x + 3, mobile_y + 3, mobile_x + 83, mobile_y + 163], 
                             radius=20, fill='#00000030')
        
        # Phone body
        draw.rounded_rectangle([mobile_x, mobile_y, mobile_x + 80, mobile_y + 160], 
                             radius=20, fill='#263238', outline='#37474F', width=3)
        
        # Screen
        draw.rounded_rectangle([mobile_x + 8, mobile_y + 20, mobile_x + 72, mobile_y + 140], 
                             radius=8, fill='white', outline='#90A4AE', width=1)
        
        # Mobile header
        draw.rounded_rectangle([mobile_x + 10, mobile_y + 25, mobile_x + 70, mobile_y + 45], 
                             radius=4, fill='#2196F3')
        draw.text((mobile_x + 25, mobile_y + 30), "Mobile", fill='white', font=font)
        
        # Mobile content blocks
        for i in range(3):
            block_y = mobile_y + 55 + i * 20
            draw.rounded_rectangle([mobile_x + 12, block_y, mobile_x + 68, block_y + 15], 
                                 radius=3, fill='#F5F5F5', outline='#E0E0E0', width=1)
    
    def _draw_system_architecture(self, draw, width, height, font, title_font):
        """Draw system architecture diagram"""
        # Database
        db_x, db_y = 60, 180
        draw.ellipse([db_x, db_y, db_x + 60, db_y + 40], outline='black', width=2)
        draw.text((db_x + 15, db_y + 15), "DB", fill='black', font=font)
        
        # API/Server
        api_x, api_y = 180, 140
        draw.rectangle([api_x, api_y, api_x + 80, api_y + 60], outline='black', width=2)
        draw.text((api_x + 25, api_y + 25), "API", fill='black', font=font)
        
        # Frontend
        ui_x, ui_y = 320, 100
        draw.rectangle([ui_x, ui_y, ui_x + 60, ui_y + 80], outline='black', width=2)
        draw.text((ui_x + 20, ui_y + 35), "UI", fill='black', font=font)
        
        # Connections
        draw.line([db_x + 60, db_y + 20, api_x, api_y + 30], fill='black', width=2)
        draw.line([api_x + 80, api_y + 30, ui_x, ui_y + 40], fill='black', width=2)
        
        # Labels
        draw.text((50, 30), "System Architecture", fill='black', font=font)
    
    def _draw_learning_path(self, draw, width, height, font, title_font):
        """Draw learning progression diagram"""
        # Learning steps as ascending stairs
        steps = ["Start", "Learn", "Practice", "Master"]
        for i, step in enumerate(steps):
            x = 50 + i * 80
            y = height - 60 - i * 30
            w, h = 60, 30
            
            draw.rectangle([x, y, x + w, y + h], outline='black', width=2)
            bbox = draw.textbbox((0, 0), step, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((x + w//2 - text_width//2, y + h//2 - 6), step, fill='black', font=font)
            
            # Arrow to next step
            if i < len(steps) - 1:
                draw.line([x + w + 5, y + h//2, x + w + 15, y + h//2 - 15], fill='black', width=2)
        
        # Add lightbulb icon
        bulb_x, bulb_y = width - 80, 50
        draw.ellipse([bulb_x, bulb_y, bulb_x + 30, bulb_y + 30], outline='black', width=2)
        draw.text((bulb_x + 8, bulb_y + 8), "💡", fill='black', font=font)
    
    def _draw_collaboration_diagram(self, draw, width, height, font, title_font):
        """Draw team collaboration diagram"""
        # Central meeting table (ellipse)
        center_x, center_y = width // 2, height // 2
        draw.ellipse([center_x - 60, center_y - 30, center_x + 60, center_y + 30], 
                    outline='black', width=2)
        
        # People around table (circles)
        positions = [(center_x - 80, center_y - 60), (center_x + 80, center_y - 60),
                    (center_x - 80, center_y + 60), (center_x + 80, center_y + 60)]
        
        for i, (x, y) in enumerate(positions):
            # Person (circle)
            draw.ellipse([x - 15, y - 15, x + 15, y + 15], outline='black', width=2)
            draw.text((x - 5, y - 5), str(i+1), fill='black', font=font)
            
            # Speech bubble
            bubble_x = x + (20 if i % 2 == 0 else -40)
            bubble_y = y - 30
            draw.ellipse([bubble_x, bubble_y, bubble_x + 20, bubble_y + 15], outline='black', width=1)
        
        draw.text((center_x - 30, 30), "Team Meeting", fill='black', font=font)
    
    def _draw_problem_solution(self, draw, width, height, font, title_font):
        """Draw problem-solution flow"""
        # Problem box (left)
        prob_x, prob_y = 40, height//2 - 40
        draw.rectangle([prob_x, prob_y, prob_x + 80, prob_y + 80], outline='black', width=2)
        draw.text((prob_x + 15, prob_y + 35), "Problem", fill='black', font=font)
        
        # Warning icon
        draw.text((prob_x + 35, prob_y + 10), "⚠", fill='black', font=font)
        
        # Arrow
        arrow_y = height // 2
        draw.line([prob_x + 90, arrow_y, prob_x + 150, arrow_y], fill='black', width=3)
        draw.line([prob_x + 145, arrow_y - 5, prob_x + 150, arrow_y], fill='black', width=3)
        draw.line([prob_x + 145, arrow_y + 5, prob_x + 150, arrow_y], fill='black', width=3)
        
        # Solution box (right)
        sol_x = prob_x + 160
        draw.rectangle([sol_x, prob_y, sol_x + 80, prob_y + 80], outline='black', width=2)
        draw.text((sol_x + 15, prob_y + 35), "Solution", fill='black', font=font)
        
        # Checkmark
        draw.text((sol_x + 35, prob_y + 10), "✓", fill='black', font=font)
    
    def _draw_generic_concept(self, draw, width, height, font, title_font, bullet_points):
        """Draw generic concept map with key terms"""
        center_x, center_y = width // 2, height // 2
        
        # Central concept circle
        draw.ellipse([center_x - 40, center_y - 30, center_x + 40, center_y + 30], 
                    outline='black', width=2)
        draw.text((center_x - 20, center_y - 10), "Concept", fill='black', font=font)
        
        # Surrounding concept bubbles
        if bullet_points:
            positions = [(center_x - 100, center_y - 80), (center_x + 100, center_y - 80),
                        (center_x - 100, center_y + 80), (center_x + 100, center_y + 80)]
            
            for i, (x, y) in enumerate(positions[:min(4, len(bullet_points))]):
                # Extract first word from bullet point
                words = bullet_points[i].split()[:2] if i < len(bullet_points) else ["Item"]
                text = " ".join(words)[:10]  # Limit length
                
                # Draw bubble
                bubble_width = max(40, len(text) * 6)
                draw.ellipse([x - bubble_width//2, y - 15, x + bubble_width//2, y + 15], 
                           outline='black', width=1)
                
                # Add text
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                draw.text((x - text_width//2, y - 5), text, fill='black', font=font)
                
                # Connect to center
                draw.line([center_x, center_y, x, y], fill='black', width=1)
    
    def _calculate_heading_score(self, title: str, position_ratio: float = 0.5) -> dict:
        """
        Calculate multi-factor scores for each heading level (1-4).
        Returns dict with scores for each level and recommended level.

        Args:
            title: The heading text
            position_ratio: Where in document (0.0 = start, 1.0 = end)

        Returns:
            {'level_1_score': float, 'level_2_score': float, ..., 'recommended_level': int}
        """
        title_lower = title.lower()
        word_count = len(title.split())
        char_count = len(title)

        # Initialize scores
        scores = {1: 0.0, 2: 0.0, 3: 0.0, 4: 0.0}

        # LEVEL 1 SCORING (Major sections: Introduction, Overview, Conclusion)
        # Strong indicators (weight: 3.0)
        level_1_strong = ['introduction', 'overview', 'conclusion', 'summary', 'abstract',
                          'table of contents', 'appendix', 'references', 'about']
        # Medium indicators (weight: 2.0)
        level_1_medium = ['what is', 'getting started', 'fundamentals', 'key concepts',
                          'background', 'motivation', 'objectives', 'goals']
        # Weak indicators (weight: 1.0)
        level_1_weak = ['welcome', 'agenda', 'outline', 'roadmap']

        for keyword in level_1_strong:
            if keyword in title_lower:
                scores[1] += 3.0
        for keyword in level_1_medium:
            if keyword in title_lower:
                scores[1] += 2.0
        for keyword in level_1_weak:
            if keyword in title_lower:
                scores[1] += 1.0

        # Position bonus for H1 (early in document = more likely H1)
        if position_ratio < 0.2:  # First 20% of document
            scores[1] += 2.0
        elif position_ratio > 0.8:  # Last 20% of document (conclusions)
            scores[1] += 1.5

        # LEVEL 2 SCORING (Major subsections)
        # Strong indicators
        level_2_strong = ['types of', 'categories', 'classification', 'approaches',
                          'methodology', 'architecture', 'components', 'applications']
        # Medium indicators
        level_2_medium = ['supervised', 'unsupervised', 'reinforcement', 'algorithm',
                          'healthcare', 'finance', 'technology', 'prerequisites',
                          'requirements', 'setup', 'installation', 'configuration']
        # Weak indicators
        level_2_weak = ['first', 'second', 'third', 'next', 'then', 'finally']

        for keyword in level_2_strong:
            if keyword in title_lower:
                scores[2] += 3.0
        for keyword in level_2_medium:
            if keyword in title_lower:
                scores[2] += 2.0
        for keyword in level_2_weak:
            if keyword in title_lower:
                scores[2] += 1.0

        # LEVEL 3 SCORING (Subsections)
        level_3_strong = ['step', 'phase', 'stage', 'part', 'section', 'module']
        level_3_medium = ['lesson', 'chapter', 'unit', 'topic', 'case study',
                          'example', 'scenario', 'workflow']
        level_3_weak = ['process', 'procedure', 'method', 'technique']

        for keyword in level_3_strong:
            if keyword in title_lower:
                scores[3] += 3.0
        for keyword in level_3_medium:
            if keyword in title_lower:
                scores[3] += 2.0
        for keyword in level_3_weak:
            if keyword in title_lower:
                scores[3] += 1.0

        # LEVEL 4 SCORING (Individual slide titles - specific topics)
        level_4_strong = ['how to', 'why', 'when to', 'best practices', 'tips', 'tricks']
        level_4_medium = ['implementation', 'details', 'features', 'benefits',
                          'challenges', 'limitations', 'advantages', 'disadvantages']
        level_4_weak = ['what', 'which', 'who', 'where', 'common', 'key points',
                        'important', 'note', 'remember']

        for keyword in level_4_strong:
            if keyword in title_lower:
                scores[4] += 3.0
        for keyword in level_4_medium:
            if keyword in title_lower:
                scores[4] += 2.0
        for keyword in level_4_weak:
            if keyword in title_lower:
                scores[4] += 1.0

        # LENGTH-BASED SCORING (optimal ranges per level)
        # H1: 15-40 chars, 2-5 words
        # H2: 20-50 chars, 3-7 words
        # H3: 15-45 chars, 2-6 words
        # H4: 10-60 chars, 2-10 words (most flexible)

        if 15 <= char_count <= 40 and 2 <= word_count <= 5:
            scores[1] += 1.5
        if 20 <= char_count <= 50 and 3 <= word_count <= 7:
            scores[2] += 1.5
        if 15 <= char_count <= 45 and 2 <= word_count <= 6:
            scores[3] += 1.5
        if 10 <= char_count <= 60 and 2 <= word_count <= 10:
            scores[4] += 1.5

        # Very short titles lean toward H4 (specific topics)
        if word_count <= 3:
            scores[4] += 1.0

        # Very long titles lean toward H2/H3 (descriptive sections)
        if word_count >= 8:
            scores[2] += 0.5
            scores[3] += 0.5

        # Determine recommended level (highest score)
        recommended_level = max(scores.keys(), key=lambda k: scores[k])

        # If all scores are zero or very low, use length-based fallback
        if max(scores.values()) < 1.0:
            if word_count <= 3:
                recommended_level = 4
            elif word_count <= 6:
                recommended_level = 3
            else:
                recommended_level = 2

        result = {
            'level_1_score': scores[1],
            'level_2_score': scores[2],
            'level_3_score': scores[3],
            'level_4_score': scores[4],
            'recommended_level': recommended_level
        }

        logger.debug(f"Heading score for '{title}': {result}")
        return result

    def _validate_heading_hierarchy(self, slides: List[SlideContent]) -> List[SlideContent]:
        """
        Validate and correct heading hierarchy to ensure logical progression.
        Rules:
        - H3 cannot appear without parent H2
        - H4 cannot appear without parent H3
        - Large jumps (H1 → H4) are corrected
        - First heading slide should be H1 or H2

        Args:
            slides: List of SlideContent objects

        Returns:
            Updated list with corrected heading levels
        """
        last_heading_level = 0
        corrections_made = 0

        for i, slide in enumerate(slides):
            if slide.slide_type == 'heading' and slide.heading_level:
                current_level = slide.heading_level
                original_level = current_level

                # First heading should be H1 or H2
                if last_heading_level == 0:
                    if current_level > 2:
                        current_level = 2
                        logger.info(f"📐 Corrected first heading from H{original_level} to H{current_level}: '{slide.title}'")
                        corrections_made += 1

                # Cannot skip levels downward (H1 → H3 or H1 → H4)
                elif current_level > last_heading_level + 1:
                    # Allow max 1 level jump
                    current_level = last_heading_level + 1
                    logger.info(f"📐 Corrected heading jump from H{last_heading_level} → H{original_level} to H{last_heading_level} → H{current_level}: '{slide.title}'")
                    corrections_made += 1

                # Cannot have H3 without H2 parent
                elif current_level == 3 and last_heading_level < 2:
                    current_level = 2
                    logger.info(f"📐 Promoted orphaned H3 to H2: '{slide.title}'")
                    corrections_made += 1

                # Cannot have H4 without H3 parent
                elif current_level == 4 and last_heading_level < 3:
                    current_level = 3
                    logger.info(f"📐 Promoted orphaned H4 to H3: '{slide.title}'")
                    corrections_made += 1

                # Update slide if level changed
                if current_level != original_level:
                    slide.heading_level = current_level

                last_heading_level = current_level

        if corrections_made > 0:
            logger.info(f"📐 Heading hierarchy validation: {corrections_made} corrections applied")

        return slides

    def _optimize_slide_density(self, slides: List[SlideContent]) -> List[SlideContent]:
        """
        Optimize slide density by merging sparse slides and splitting dense slides.

        Rules:
        - Sparse slides (1-2 bullets): Merge consecutive content slides under same section
        - Dense slides (7+ bullets): Split into multiple slides
        - Optimal range: 3-6 bullets per slide
        - Don't merge across heading boundaries (H1/H2)
        - Preserve bullet order and slide metadata

        Args:
            slides: List of SlideContent objects

        Returns:
            Optimized list with better slide density
        """
        optimized_slides = []
        merges_made = 0
        splits_made = 0

        i = 0
        while i < len(slides):
            slide = slides[i]

            # Skip heading slides (never merge/split these)
            if slide.slide_type == 'heading':
                optimized_slides.append(slide)
                i += 1
                continue

            # DENSE SLIDE SPLITTING (7+ bullets)
            if slide.content and len(slide.content) >= 7:
                # Split into chunks of 4-5 bullets
                bullets = slide.content
                chunk_size = 5
                num_chunks = (len(bullets) + chunk_size - 1) // chunk_size

                for chunk_idx in range(num_chunks):
                    start_idx = chunk_idx * chunk_size
                    end_idx = min(start_idx + chunk_size, len(bullets))
                    chunk_bullets = bullets[start_idx:end_idx]

                    # Add continuation marker to title if not first chunk
                    chunk_title = slide.title
                    if chunk_idx > 0:
                        chunk_title = f"{slide.title} (cont.)"

                    optimized_slides.append(SlideContent(
                        title=chunk_title,
                        content=chunk_bullets,
                        slide_type=slide.slide_type,
                        heading_level=slide.heading_level,
                        subheader=slide.subheader if chunk_idx == 0 else None,
                        visual_cues=slide.visual_cues if chunk_idx == 0 else None
                    ))

                splits_made += 1
                logger.info(f"📊 Split dense slide '{slide.title}' ({len(bullets)} bullets) into {num_chunks} slides")
                i += 1
                continue

            # SPARSE SLIDE MERGING (1-2 bullets)
            if slide.content and len(slide.content) <= 2:
                # Look ahead to find consecutive sparse content slides under same section
                merge_candidates = [slide]
                j = i + 1

                # Track the last major heading we've seen
                last_major_heading = None
                for k in range(i - 1, -1, -1):
                    if slides[k].slide_type == 'heading' and slides[k].heading_level and slides[k].heading_level <= 2:
                        last_major_heading = slides[k]
                        break

                # Collect consecutive sparse slides (up to 6 bullets total)
                while j < len(slides):
                    next_slide = slides[j]

                    # Stop if we hit a heading slide
                    if next_slide.slide_type == 'heading':
                        # If it's H1/H2, definitely stop (section boundary)
                        if next_slide.heading_level and next_slide.heading_level <= 2:
                            break
                        # If it's H3/H4, still stop (different subsection)
                        break

                    # Stop if next slide has good density (3+ bullets)
                    if next_slide.content and len(next_slide.content) >= 3:
                        break

                    # Stop if adding would exceed optimal range
                    total_bullets = sum(len(c.content) for c in merge_candidates)
                    if next_slide.content and total_bullets + len(next_slide.content) > 6:
                        break

                    # Add to merge candidates
                    if next_slide.content and len(next_slide.content) > 0:
                        merge_candidates.append(next_slide)

                    j += 1

                # Only merge if we found at least 2 slides to combine
                if len(merge_candidates) >= 2:
                    # Merge bullets from all candidates
                    merged_bullets = []
                    for candidate in merge_candidates:
                        merged_bullets.extend(candidate.content)

                    # Use first slide's title
                    merged_slide = SlideContent(
                        title=merge_candidates[0].title,
                        content=merged_bullets,
                        slide_type=merge_candidates[0].slide_type,
                        heading_level=merge_candidates[0].heading_level,
                        subheader=merge_candidates[0].subheader,
                        visual_cues=merge_candidates[0].visual_cues
                    )

                    optimized_slides.append(merged_slide)
                    merges_made += 1
                    logger.info(f"📊 Merged {len(merge_candidates)} sparse slides into '{merged_slide.title}' ({len(merged_bullets)} bullets)")

                    # Skip all merged slides
                    i = j
                    continue

            # No optimization needed - keep slide as-is
            optimized_slides.append(slide)
            i += 1

        if merges_made > 0 or splits_made > 0:
            logger.info(f"📊 Slide density optimization: {merges_made} merges, {splits_made} splits applied")

        return optimized_slides

    def _insert_section_dividers(self, slides: List[SlideContent]) -> List[SlideContent]:
        """
        Insert visual divider slides before major section headings (H1 and H2).

        Divider slides create visual breaks in presentations, helping presenters
        transition between major topics and providing natural pause points.

        Rules:
        - Insert divider before H1 headings (except the very first heading)
        - Insert divider before H2 headings (major section transitions)
        - Dividers have type 'divider' for special rendering
        - Dividers preserve the section title and heading level

        Args:
            slides: List of slide content objects

        Returns:
            List of slides with dividers inserted
        """
        enhanced_slides = []
        dividers_added = 0
        first_heading_seen = False

        for i, slide in enumerate(slides):
            # Check if this is a major heading (H1 or H2)
            if slide.slide_type == 'heading' and slide.heading_level in [1, 2]:
                # Skip divider for the very first heading (that's the title/intro)
                if not first_heading_seen:
                    first_heading_seen = True
                    enhanced_slides.append(slide)
                    continue

                # Insert divider slide before this heading
                divider_slide = SlideContent(
                    title=slide.title,
                    content=[],
                    slide_type='divider',
                    heading_level=slide.heading_level,
                    subheader=None,
                    visual_cues=None
                )

                enhanced_slides.append(divider_slide)
                dividers_added += 1

                logger.info(f"📑 Inserted section divider before H{slide.heading_level}: '{slide.title}'")

            # Add the original slide
            enhanced_slides.append(slide)

        if dividers_added > 0:
            logger.info(f"📑 Section divider insertion: {dividers_added} divider slides added")
        else:
            logger.info("📑 No section dividers needed (no major section transitions)")

        return enhanced_slides

    def _determine_heading_level(self, title: str, position_ratio: float = 0.5) -> int:
        """
        Determine the heading level of a title using multi-factor scoring.

        Args:
            title: The heading text
            position_ratio: Position in document (0.0 = start, 1.0 = end)

        Returns:
            Heading level (1-4)
        """
        scoring_result = self._calculate_heading_score(title, position_ratio)
        recommended_level = scoring_result['recommended_level']

        logger.debug(f"📊 Determined '{title}' as H{recommended_level} (scores: " +
                    f"H1={scoring_result['level_1_score']:.1f}, " +
                    f"H2={scoring_result['level_2_score']:.1f}, " +
                    f"H3={scoring_result['level_3_score']:.1f}, " +
                    f"H4={scoring_result['level_4_score']:.1f})")

        return recommended_level
    
    def _create_section_overview(self, content: List[str]) -> str:
        """Create a brief overview for section intro slides"""
        if not content:
            return None
        
        # Take first sentence or two as overview
        overview_parts = []
        char_count = 0
        
        for item in content[:3]:  # Max 3 items
            if char_count + len(item) > 150:  # Keep overview concise
                break
            overview_parts.append(item)
            char_count += len(item)
        
        return ' '.join(overview_parts) if overview_parts else None
    
    def create_html_slides(self, doc_structure: DocumentStructure) -> str:
        """Generate HTML presentation from document structure"""
        slides_html = []
        
        # Title slide
        title_slide = f"""
        <section class="slide title-slide">
            <h1>{doc_structure.title}</h1>
            <p class="subtitle">Generated from {doc_structure.metadata['filename']}</p>
            <p class="date">{datetime.now().strftime('%B %d, %Y')}</p>
        </section>
        """
        slides_html.append(title_slide)
        
        # Content slides
        for slide_content in doc_structure.slides:
            content_items = '</li><li>'.join(slide_content.content[:8]) if slide_content.content else ''
            
            slide_html = f"""
            <section class="slide content-slide">
                <h2>{slide_content.title}</h2>
                {f'<ul><li>{content_items}</li></ul>' if content_items else ''}
            </section>
            """
            slides_html.append(slide_html)
        
        # Use the same HTML template from the original app
        html_template = self._get_html_template()
        full_html = html_template.replace('{{SLIDES}}', '\n'.join(slides_html))
        full_html = full_html.replace('{{TITLE}}', doc_structure.title)
        
        # Save HTML file
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        filepath = os.path.join(EXPORT_FOLDER, filename)
        os.makedirs(EXPORT_FOLDER, exist_ok=True)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(full_html)
        
        return filepath
    
    def _generate_ai_image(self, drawing_prompt: str, slide_number: int) -> str:
        """Generate an AI image using OpenAI DALL-E based on the drawing prompt"""
        try:
            logger.info(f"Generating AI image for slide {slide_number} with DALL-E")
            
            # Convert the educational prompt to a DALL-E-optimized prompt
            dalle_prompt = self._convert_to_dalle_prompt(drawing_prompt)
            logger.info(f"DALL-E prompt: {dalle_prompt[:100]}...")
            
            # Generate image using OpenAI DALL-E with timeout
            logger.info(f"Making DALL-E API call with prompt length: {len(dalle_prompt)}")
            
            # Set a shorter timeout for individual image generation
            import signal
            def timeout_handler(signum, frame):
                raise TimeoutError("DALL-E API call timed out")
            
            signal.signal(signal.SIGALRM, timeout_handler)
            signal.alarm(30)  # 30 second timeout per image
            
            try:
                response = self.client.images.generate(
                    model="dall-e-2",  # Use DALL-E 2 for better compatibility
                    prompt=dalle_prompt,
                    size="512x512",    # Use smaller size for faster generation
                    n=1
                )
            finally:
                signal.alarm(0)  # Cancel the alarm
            
            # Download and save the generated image
            image_url = response.data[0].url
            
            # Download the image
            image_response = requests.get(image_url, timeout=30)
            image_response.raise_for_status()
            
            # Save the image
            sketch_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'temp_sketches')
            os.makedirs(sketch_dir, exist_ok=True)
            
            # Create unique filename
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            image_filename = f"ai_generated_slide{slide_number:03d}_{timestamp}_{random.randint(1000000, 9999999)}.png"
            image_path = os.path.join(sketch_dir, image_filename)
            
            # Save the downloaded image
            with open(image_path, 'wb') as f:
                f.write(image_response.content)
            
            logger.info(f"✅ AI image generated successfully: {image_filename}")
            return image_path
            
        except Exception as e:
            logger.error(f"Failed to generate AI image for slide {slide_number}: {type(e).__name__}: {e}")
            # Log more details for debugging
            import traceback
            logger.error(f"Full traceback: {traceback.format_exc()}")
            # Fallback to thought bubble image
            return self._create_thought_bubble_image(drawing_prompt, slide_number)
    
    def _convert_to_dalle_prompt(self, educational_prompt: str) -> str:
        """Convert educational prompt to DALL-E optimized prompt"""
        # Extract the visual goal from the educational prompt
        if "Visual Goal:" in educational_prompt:
            visual_goal_start = educational_prompt.find("Visual Goal:") + len("Visual Goal:")
            style_start = educational_prompt.find("Style Guidelines:")
            if style_start != -1:
                visual_goal = educational_prompt[visual_goal_start:style_start].strip()
            else:
                visual_goal = educational_prompt[visual_goal_start:].strip()
        else:
            # Fallback - use the whole prompt
            visual_goal = educational_prompt
        
        # Add DALL-E specific optimizations
        dalle_prompt = visual_goal
        
        # Ensure it ends with style specifications
        if "flat design" not in dalle_prompt.lower():
            dalle_prompt += " Create this as a flat design illustration with clean lines and simple colors."
        
        # Ensure no text constraint
        if "no text" not in dalle_prompt.lower():
            dalle_prompt += " No text or labels in the image."
        
        # Limit length to DALL-E's requirements (1000 chars max)
        if len(dalle_prompt) > 950:
            dalle_prompt = dalle_prompt[:950] + "..."
        
        return dalle_prompt

    def _create_thought_bubble_image(self, drawing_prompt: str, slide_number: int) -> str:
        """Create a thought bubble with visual prompt on light blue background"""
        try:
            # Create image dimensions
            width, height = 800, 600
            
            # Light blue background (#E6F3FF)
            image = Image.new('RGB', (width, height), '#E6F3FF')
            draw = ImageDraw.Draw(image)
            
            # Load fonts
            try:
                font_paths = [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/System/Library/Fonts/Arial.ttf", 
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/Windows/Fonts/arial.ttf"
                ]
                font_loaded = False
                for font_path in font_paths:
                    try:
                        font = ImageFont.truetype(font_path, 16)
                        title_font = ImageFont.truetype(font_path, 20)
                        font_loaded = True
                        break
                    except:
                        continue
                
                if not font_loaded:
                    raise Exception("No fonts found")
                    
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
            
            # Draw thought bubble outline
            bubble_margin = 40
            bubble_x1 = bubble_margin
            bubble_y1 = bubble_margin
            bubble_x2 = width - bubble_margin
            bubble_y2 = height - bubble_margin
            
            # Main thought bubble (white with light blue border)
            draw.ellipse([bubble_x1, bubble_y1, bubble_x2, bubble_y2], 
                        fill='white', outline='#4A90E2', width=3)
            
            # Small thought bubbles leading to main bubble
            small_bubble_1_x = bubble_x1 + 30
            small_bubble_1_y = bubble_y2 + 10
            draw.ellipse([small_bubble_1_x, small_bubble_1_y, 
                         small_bubble_1_x + 20, small_bubble_1_y + 20],
                        fill='white', outline='#4A90E2', width=2)
            
            small_bubble_2_x = bubble_x1 + 60
            small_bubble_2_y = bubble_y2 + 25
            draw.ellipse([small_bubble_2_x, small_bubble_2_y,
                         small_bubble_2_x + 15, small_bubble_2_y + 15],
                        fill='white', outline='#4A90E2', width=2)
            
            # Add title at the top of bubble
            title_text = "💡 Visual Inspiration"
            title_bbox = draw.textbbox((0, 0), title_text, font=title_font)
            title_width = title_bbox[2] - title_bbox[0]
            title_x = width//2 - title_width//2
            title_y = bubble_y1 + 20
            
            draw.text((title_x, title_y), title_text, font=title_font, 
                     fill='#2C5282', anchor='mm')
            
            # Wrap and draw the prompt text
            max_width = bubble_x2 - bubble_x1 - 80  # Leave margins
            wrapped_lines = self._wrap_text(drawing_prompt, max_width, font)
            
            # Calculate total text height
            line_height = 24
            total_text_height = len(wrapped_lines) * line_height
            start_y = bubble_y1 + 80  # Start below title
            
            # Draw each line of wrapped text
            for i, line in enumerate(wrapped_lines):
                text_bbox = draw.textbbox((0, 0), line, font=font)
                text_width = text_bbox[2] - text_bbox[0]
                text_x = width//2 - text_width//2
                text_y = start_y + (i * line_height)
                
                draw.text((text_x, text_y), line, font=font, fill='#2D3748')
            
            # Add copy instruction at bottom
            copy_text = "💻 Copy & paste this into any AI image generator"
            copy_bbox = draw.textbbox((0, 0), copy_text, font=font)
            copy_width = copy_bbox[2] - copy_bbox[0]
            copy_x = width//2 - copy_width//2
            copy_y = bubble_y2 - 40
            
            draw.text((copy_x, copy_y), copy_text, font=font, 
                     fill='#4A90E2', anchor='mm')
            
            # Save the image
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
            random_suffix = random.randint(1000, 9999)
            filename = f"thought_bubble_slide{slide_number:03d}_{timestamp}_{random_suffix}.png"
            
            # Ensure temp_sketches directory exists
            os.makedirs('temp_sketches', exist_ok=True)
            filepath = os.path.join('temp_sketches', filename)
            
            image.save(filepath, 'PNG', optimize=True, quality=85)
            logger.info(f"Created thought bubble image: {filepath}")
            
            return filepath
            
        except Exception as e:
            logger.error(f"Error creating thought bubble image: {e}")
            return None
    
    def _get_html_template(self) -> str:
        """HTML template for slides - same as original"""
        return """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{TITLE}}</title>
    <style>
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 0;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: #333;
        }
        
        .slideshow-container {
            max-width: 1000px;
            position: relative;
            margin: auto;
            height: 100vh;
            display: flex;
            align-items: center;
        }
        
        .slide {
            display: none;
            padding: 60px;
            text-align: center;
            background: white;
            border-radius: 10px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.3);
            margin: 20px;
            min-height: 500px;
            display: flex;
            flex-direction: column;
            justify-content: center;
        }
        
        .slide.active {
            display: flex;
        }
        
        .title-slide h1 {
            font-size: 3em;
            color: #667eea;
            margin-bottom: 30px;
        }
        
        .content-slide h2 {
            font-size: 2.5em;
            color: #667eea;
            margin-bottom: 30px;
            border-bottom: 3px solid #667eea;
            padding-bottom: 15px;
        }
        
        .content-slide ul {
            text-align: left;
            font-size: 1.4em;
            line-height: 1.8;
            list-style-type: none;
            padding-left: 0;
        }
        
        .content-slide li {
            margin: 15px 0;
            padding-left: 30px;
            position: relative;
        }
        
        .content-slide li::before {
            content: "→";
            color: #667eea;
            font-weight: bold;
            position: absolute;
            left: 0;
        }
        
        .subtitle {
            font-size: 1.2em;
            color: #666;
            margin: 20px 0;
        }
        
        .date {
            font-size: 1em;
            color: #999;
        }
        
        .navigation {
            position: fixed;
            bottom: 30px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 15px;
        }
        
        .nav-btn {
            background: #667eea;
            color: white;
            border: none;
            padding: 10px 20px;
            border-radius: 25px;
            cursor: pointer;
            font-size: 1em;
            transition: all 0.3s ease;
        }
        
        .nav-btn:hover {
            background: #5a67d8;
            transform: translateY(-2px);
        }
        
        .nav-btn:disabled {
            background: #ccc;
            cursor: not-allowed;
            transform: none;
        }
        
        .slide-counter {
            position: fixed;
            top: 30px;
            right: 30px;
            background: rgba(255,255,255,0.9);
            padding: 10px 15px;
            border-radius: 20px;
            font-weight: bold;
            color: #667eea;
        }
    </style>
</head>
<body>
    <div class="slideshow-container">
        {{SLIDES}}
    </div>
    
    <div class="slide-counter">
        <span id="current-slide">1</span> / <span id="total-slides"></span>
    </div>
    
    <div class="navigation">
        <button class="nav-btn" id="prev-btn" onclick="changeSlide(-1)">← Previous</button>
        <button class="nav-btn" id="next-btn" onclick="changeSlide(1)">Next →</button>
    </div>
    
    <script>
        let currentSlide = 0;
        const slides = document.querySelectorAll('.slide');
        const totalSlides = slides.length;
        
        document.getElementById('total-slides').textContent = totalSlides;
        
        function showSlide(n) {
            slides.forEach(slide => slide.classList.remove('active'));
            
            if (n >= totalSlides) currentSlide = 0;
            if (n < 0) currentSlide = totalSlides - 1;
            
            slides[currentSlide].classList.add('active');
            document.getElementById('current-slide').textContent = currentSlide + 1;
            
            document.getElementById('prev-btn').disabled = currentSlide === 0;
            document.getElementById('next-btn').disabled = currentSlide === totalSlides - 1;
        }
        
        function changeSlide(n) {
            currentSlide += n;
            showSlide(currentSlide);
        }
        
        // Keyboard navigation
        document.addEventListener('keydown', function(e) {
            if (e.key === 'ArrowLeft') changeSlide(-1);
            if (e.key === 'ArrowRight') changeSlide(1);
        });
        
        // Initialize
        showSlide(0);
    </script>
</body>
</html>
        """

    def _generate_simple_sketch(self, prompt_text: str, title: str, slide_number: int) -> str:
        """Generate visual prompt text display (no longer creates images)"""
        logger.info(f"Creating visual prompt text display for slide {slide_number}")
        
        try:
            # Create a text-based visual prompt display (620x504 pixels)
            width, height = 620, 504
            image = Image.new('RGB', (width, height), '#f8f9fa')  # Light gray background
            draw = ImageDraw.Draw(image)
            
            # Load fonts
            try:
                font_paths = [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/System/Library/Fonts/Arial.ttf", 
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/Windows/Fonts/arial.ttf"
                ]
                font_loaded = False
                for font_path in font_paths:
                    try:
                        title_font = ImageFont.truetype(font_path, 18)
                        text_font = ImageFont.truetype(font_path, 14)
                        font_loaded = True
                        break
                    except:
                        continue
                
                if not font_loaded:
                    title_font = ImageFont.load_default()
                    text_font = ImageFont.load_default()
                    
            except:
                title_font = ImageFont.load_default()
                text_font = ImageFont.load_default()
            
            # Draw border
            draw.rectangle([10, 10, width-10, height-10], outline='#d1d5db', width=2)
            
            # Draw title
            title_text = f"Visual Prompt for Slide {slide_number}"
            draw.text((20, 25), title_text, fill='#374151', font=title_font)
            
            # Draw prompt text with word wrapping
            y_position = 60
            margin = 20
            max_width = width - 2 * margin
            
            # Simple word wrapping for the prompt text
            words = prompt_text.split()
            lines = []
            current_line = []
            
            for word in words:
                test_line = ' '.join(current_line + [word])
                bbox = draw.textbbox((0, 0), test_line, font=text_font)
                line_width = bbox[2] - bbox[0]
                
                if line_width <= max_width:
                    current_line.append(word)
                else:
                    if current_line:
                        lines.append(' '.join(current_line))
                        current_line = [word]
                    else:
                        lines.append(word)
            
            if current_line:
                lines.append(' '.join(current_line))
            
            # Draw the wrapped text
            for line in lines[:15]:  # Limit to 15 lines
                draw.text((margin, y_position), line, fill='#374151', font=text_font)
                y_position += 25
                if y_position > height - 50:
                    break
            
            # Add instruction at bottom
            instruction = "💡 Copy this prompt to use with any AI image generator"
            draw.text((margin, height - 35), instruction, fill='#6b7280', font=text_font)
            
            # Save the image to temp directory
            temp_dir = os.path.join(os.getcwd(), 'temp_sketches')
            os.makedirs(temp_dir, exist_ok=True)
            
            # Create filename with slide info
            safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()[:30]
            filename = f"prompt_slide{slide_number:03d}_{safe_title.replace(' ', '_') if safe_title else 'untitled'}.png"
            filepath = os.path.join(temp_dir, filename)
            
            image.save(filepath)
            logger.info(f"Generated visual prompt display: {filename}")
            
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating visual prompt display: {e}")
            return None
    
    def _parse_visual_prompt(self, prompt: str) -> dict:
        """Extract key elements from the visual prompt for sketching"""
        elements = {
            'has_person': False,
            'has_computer': False,
            'has_chart': False,
            'has_arrows': False,
            'has_lightbulb': False,
            'layout': 'single'
        }
        
        prompt_lower = prompt.lower()
        
        # Detect common elements
        if any(word in prompt_lower for word in ['person', 'user', 'developer', 'student', 'learner', 'people']):
            elements['has_person'] = True
            
        if any(word in prompt_lower for word in ['computer', 'laptop', 'screen', 'app', 'interface', 'streamlit']):
            elements['has_computer'] = True
            
        if any(word in prompt_lower for word in ['chart', 'graph', 'data', 'visualization', 'dashboard']):
            elements['has_chart'] = True
            
        if any(word in prompt_lower for word in ['arrow', 'flow', 'process', 'step', 'journey', 'path']):
            elements['has_arrows'] = True
            
        if any(word in prompt_lower for word in ['idea', 'lightbulb', 'concept', 'innovation', 'creative']):
            elements['has_lightbulb'] = True
        
        # Detect layout patterns
        if 'left' in prompt_lower and 'right' in prompt_lower:
            if 'middle' in prompt_lower or 'center' in prompt_lower:
                elements['layout'] = 'left-center-right'
            else:
                elements['layout'] = 'left-right'
            
        return elements
    
    def _draw_sketch_elements(self, draw: ImageDraw.Draw, elements: dict, width: int, height: int):
        """Draw simple sketch elements based on parsed prompt"""
        center_x = width // 2
        center_y = height // 2
        
        # Choose drawing style based on layout
        if elements['layout'] == 'left-center-right':
            # Three section layout with arrows
            positions = [(width//6, center_y), (center_x, center_y), (width*5//6, center_y)]
            for i, (x, y) in enumerate(positions):
                if i == 0 and elements['has_person']:
                    self._draw_person(draw, x, y)
                elif i == 1 and elements['has_computer']:
                    self._draw_computer(draw, x, y)
                elif i == 2 and elements['has_computer']:
                    self._draw_computer(draw, x, y)
                else:
                    self._draw_generic_shape(draw, x, y)
                    
            # Draw arrows
            if elements['has_arrows']:
                self._draw_arrow(draw, width//6 + 40, center_y + 50, center_x - 40, center_y + 50)
                self._draw_arrow(draw, center_x + 40, center_y + 50, width*5//6 - 40, center_y + 50)
                
        elif elements['layout'] == 'left-right':
            # Two section layout
            positions = [(width//3, center_y), (width*2//3, center_y)]
            for i, (x, y) in enumerate(positions):
                if i == 0 and elements['has_person']:
                    self._draw_person(draw, x, y)
                elif elements['has_computer']:
                    self._draw_computer(draw, x, y)
                else:
                    self._draw_generic_shape(draw, x, y)
                    
            if elements['has_arrows']:
                self._draw_arrow(draw, width//3 + 40, center_y + 40, width*2//3 - 40, center_y + 40)
                
        else:
            # Single centered element
            if elements['has_person']:
                self._draw_person(draw, center_x, center_y)
            elif elements['has_computer']:
                self._draw_computer(draw, center_x, center_y)
            elif elements['has_chart']:
                self._draw_chart(draw, center_x, center_y)
            elif elements['has_lightbulb']:
                self._draw_lightbulb(draw, center_x, center_y)
            else:
                self._draw_generic_shape(draw, center_x, center_y)
    
    def _draw_person(self, draw: ImageDraw.Draw, x: int, y: int):
        """Draw a simple stick figure"""
        draw.ellipse([x-10, y-30, x+10, y-10], outline='black', width=2)  # Head
        draw.line([x, y-10, x, y+20], fill='black', width=2)  # Body
        draw.line([x-15, y, x+15, y], fill='black', width=2)  # Arms
        draw.line([x, y+20, x-10, y+40], fill='black', width=2)  # Left leg
        draw.line([x, y+20, x+10, y+40], fill='black', width=2)  # Right leg
    
    def _draw_computer(self, draw: ImageDraw.Draw, x: int, y: int):
        """Draw a simple laptop/computer"""
        draw.rectangle([x-25, y-20, x+25, y+5], outline='black', width=2)  # Screen
        draw.rectangle([x-30, y+5, x+30, y+15], outline='black', width=2)  # Base
        # Screen content
        for i in range(3):
            draw.line([x-20, y-15+i*5, x+20, y-15+i*5], fill='black', width=1)
    
    def _draw_chart(self, draw: ImageDraw.Draw, x: int, y: int):
        """Draw a simple bar chart"""
        bars = [15, 25, 20, 30]
        start_x = x - 20
        for i, height in enumerate(bars):
            bar_x = start_x + i * 12
            draw.rectangle([bar_x, y+20-height, bar_x+8, y+20], fill='black')
        draw.line([x-25, y+25, x+25, y+25], fill='black', width=1)  # X-axis
        draw.line([x-25, y+25, x-25, y-15], fill='black', width=1)  # Y-axis
    
    def _draw_lightbulb(self, draw: ImageDraw.Draw, x: int, y: int):
        """Draw a simple lightbulb"""
        draw.ellipse([x-12, y-25, x+12, y+5], outline='black', width=2)  # Bulb
        draw.rectangle([x-8, y+5, x+8, y+15], outline='black', width=2)  # Base
        # Light rays
        for rx, ry in [(-20, -20), (20, -20), (-15, -30), (15, -30), (0, -35)]:
            draw.line([x, y-10, x+rx, y+ry], fill='black', width=1)
    
    def _draw_generic_shape(self, draw: ImageDraw.Draw, x: int, y: int):
        """Draw a generic shape"""
        draw.ellipse([x-15, y-15, x+15, y+15], outline='black', width=2)
        draw.line([x-10, y, x+10, y], fill='black', width=1)
        draw.line([x, y-10, x, y+10], fill='black', width=1)
    
    def _draw_arrow(self, draw: ImageDraw.Draw, x1: int, y1: int, x2: int, y2: int):
        """Draw a simple arrow"""
        draw.line([x1, y1, x2, y2], fill='black', width=2)
        # Arrow head
        dx, dy = x2 - x1, y2 - y1
        length = (dx*dx + dy*dy) ** 0.5
        if length > 0:
            dx, dy = dx/length, dy/length
            hx1 = x2 - 10 * (dx + 0.5 * dy)
            hy1 = y2 - 10 * (dy - 0.5 * dx)
            hx2 = x2 - 10 * (dx - 0.5 * dy)
            hy2 = y2 - 10 * (dy + 0.5 * dx)
            draw.line([x2, y2, hx1, hy1], fill='black', width=2)
            draw.line([x2, y2, hx2, hy2], fill='black', width=2)


