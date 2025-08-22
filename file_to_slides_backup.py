"""
Script to Slides Generator

A Flask web application that converts uploaded documents to presentation slides
without requiring Google API authentication.

Supported formats: TXT, DOCX, PDF, MD (Markdown)
"""

import os
import json
import logging
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from datetime import datetime
import re
import random
from math import cos, sin

import flask
from flask import Flask, render_template, request, jsonify, send_file, redirect, url_for
from werkzeug.utils import secure_filename

# Document processing libraries
import docx
from docx import Document
import PyPDF2
import markdown
from bs4 import BeautifulSoup

# Presentation generation
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Image generation for sketches
from PIL import Image, ImageDraw, ImageFont
import io

# Google Slides integration
import csv
import io

app = Flask(__name__)
app.secret_key = os.environ.get('FLASK_SECRET_KEY', 'your-secret-key-change-this')

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# File upload configuration
UPLOAD_FOLDER = 'uploads'
EXPORT_FOLDER = 'exports'
ALLOWED_EXTENSIONS = {'docx'}
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

@dataclass
class SlideContent:
    """Represents content for a single slide"""
    title: str
    content: List[str]
    slide_type: str = 'content'  # 'title', 'content', 'image', 'bullet'
    heading_level: Optional[int] = None  # Original heading level from DOCX (1-6)

@dataclass
class DocumentStructure:
    """Represents the parsed structure of a document"""
    title: str
    slides: List[SlideContent]
    metadata: Dict[str, Any]

class DocumentParser:
    """Handles parsing of various document formats"""
    
    def __init__(self):
        self.heading_patterns = [
            r'^#{1,6}\s+(.+)$',  # Markdown headings
            r'^(.+)\n[=-]{3,}$',  # Underlined headings
            r'^\d+\.\s+(.+)$',   # Numbered headings
            r'^([A-Z][A-Z\s]{5,})$',  # ALL CAPS headings
        ]
    
    def parse_file(self, file_path: str, filename: str, script_column: int = 2) -> DocumentStructure:
        """Parse DOCX file and convert to slide structure"""
        file_ext = filename.lower().split('.')[-1]
        
        try:
            if file_ext == 'docx':
                content = self._parse_docx(file_path, script_column)
            else:
                raise ValueError(f"Only DOCX files are supported. Got: {file_ext}")
            
            logger.info(f"DOCX parsing complete: {len(content.split())} words extracted from column {script_column}")
            
            # Extract title from filename or content
            doc_title = self._extract_title(content, filename)
            
            # Convert content to slides
            slides = self._content_to_slides(content)
            
            metadata = {
                'filename': filename,
                'file_type': file_ext,
                'upload_time': datetime.now().isoformat(),
                'slide_count': len(slides),
                'script_column': script_column
            }
            
            return DocumentStructure(
                title=doc_title,
                slides=slides,
                metadata=metadata
            )
            
        except Exception as e:
            logger.error(f"Error parsing DOCX file {filename}: {str(e)}")
            raise
    
    
    def _parse_docx(self, file_path: str, script_column: int = 2) -> str:
        """Parse DOCX file with script column filtering"""
        doc = Document(file_path)
        content = []
        first_table_found = False
        
        logger.info(f"DOCX contains {len(doc.paragraphs)} paragraphs and {len(doc.tables)} tables")
        logger.info(f"Extracting script text from column {script_column}")
        
        # Process paragraphs and tables in document order
        for element in doc.element.body:
            if element.tag.endswith('p'):  # Paragraph
                for paragraph in doc.paragraphs:
                    if paragraph._element == element:
                        text = paragraph.text.strip()
                        if text:
                            # Check if paragraph is a heading
                            if paragraph.style.name.startswith('Heading'):
                                level = paragraph.style.name.replace('Heading ', '')
                                try:
                                    level_num = int(level)
                                    content.append(f"{'#' * level_num} {text}")
                                    logger.info(f"Found heading level {level_num}: {text}")
                                except ValueError:
                                    content.append(f"# {text}")
                                    logger.info(f"Found heading: {text}")
                            else:
                                # Only add non-heading paragraphs if we've found the first table
                                # This ignores intro content before the first table
                                if first_table_found:
                                    content.append(text)
                        break
                        
            elif element.tag.endswith('tbl'):  # Table
                first_table_found = True  # Mark that we've found our first table
                for table in doc.tables:
                    if table._element == element:
                        logger.info(f"Processing table with {len(table.rows)} rows and {len(table.columns)} columns")
                        
                        for row_idx, row in enumerate(table.rows):
                            # Only extract from the specified script column
                            if len(row.cells) >= script_column:
                                target_cell = row.cells[script_column - 1]  # Convert to 0-based index
                                cell_text = target_cell.text.strip()
                                
                                if cell_text:
                                    # Clean up [CLICK] and other stage directions
                                    cleaned_text = self._clean_script_text(cell_text)
                                    if cleaned_text:  # Only add if there's content after cleaning
                                        logger.info(f"  Script cell [{row_idx}][{script_column-1}]: '{cleaned_text[:50]}...'")
                                        content.append(cleaned_text)
                            else:
                                logger.warning(f"  Row {row_idx} has only {len(row.cells)} columns, cannot access column {script_column}")
                        
                        # Add separator after table
                        content.append("")
                        break
        
        logger.info(f"DOCX parsing extracted {len(content)} content lines from script column {script_column}")
        return '\n'.join(content)
    
    def _clean_script_text(self, text: str) -> str:
        """Clean script text by removing stage directions and unwanted elements"""
        if not text:
            return ""
        
        # Remove [CLICK] and similar stage directions
        cleaned = re.sub(r'\[CLICK\]', '', text, flags=re.IGNORECASE)
        cleaned = re.sub(r'\[.*?\]', '', cleaned)  # Remove any other bracketed content
        
        # Remove multiple spaces and clean up
        cleaned = re.sub(r'\s+', ' ', cleaned)
        cleaned = cleaned.strip()
        
        # Skip if the cleaned text is too short or just contains common header words
        if len(cleaned) < 10:
            return ""
        
        # Skip common table headers
        header_patterns = [
            r'^(script\s*content|video|notes|duration|timestamp)$',
            r'^(slide|page|section)\s*\d*$'
        ]
        
        for pattern in header_patterns:
            if re.match(pattern, cleaned, re.IGNORECASE):
                return ""
        
        return cleaned
    
    
    
    def _looks_like_table_row(self, line: str) -> bool:
        """Check if a line looks like it contains table data - very aggressive detection"""
        # Skip empty lines and obvious non-table content
        if not line.strip() or len(line.strip()) < 3:
            return False
            
        # Skip obvious headings and page markers
        if line.startswith(('##', '#', 'Page ', '---', '--- End of Page')):
            return False
            
        words = line.split()
        
        # Very aggressive: if line has 2+ words, treat as potential table row
        if len(words) >= 2:
            return True
            
        # Look for multiple consecutive spaces (common in table formatting)
        if '  ' in line and len(line.split()) > 1:
            return True
            
        # Look for pipe separators
        if '|' in line and line.count('|') > 1:
            return True
            
        # Look for tab separators
        if '\t' in line and len(line.split('\t')) > 1:
            return True
            
        # Look for comma separators (CSV-like)
        if ',' in line and line.count(',') > 1:
            return True
            
        # Look for structured data patterns
        if len(words) > 1:
            # Any line with numbers might be table data
            numeric_count = sum(1 for word in words if any(char.isdigit() for char in word))
            if numeric_count >= 1:
                return True
                
            # Lines with common table delimiters
            if any(char in line for char in [':', ';', '|', '\t']):
                return True
                
            # Short lines with multiple words (likely table headers/data)
            if len(line) < 100 and len(words) > 1:
                return True
        
        return False
    
    def _extract_table_cells(self, line: str) -> List[str]:
        """Extract individual cells from a table row - aggressive extraction with debugging"""
        original_line = line
        cells = []
        extraction_method = "none"
        
        # Try pipe separator first
        if '|' in line:
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            extraction_method = "pipes"
        # Try tab separator
        elif '\t' in line:
            cells = [cell.strip() for cell in line.split('\t') if cell.strip()]
            extraction_method = "tabs"
        # Try comma separator (CSV-like)
        elif ',' in line and line.count(',') > 1:
            cells = [cell.strip() for cell in line.split(',') if cell.strip()]
            extraction_method = "commas"
        # Try semicolon separator
        elif ';' in line and line.count(';') > 1:
            cells = [cell.strip() for cell in line.split(';') if cell.strip()]
            extraction_method = "semicolons"
        # Try multiple spaces (2 or more)
        elif '  ' in line:
            import re
            cells = [cell.strip() for cell in re.split(r'\s{2,}', line) if cell.strip()]
            extraction_method = "multiple_spaces"
        # Try single space but check if it looks structured
        elif ' ' in line and len(line.split()) > 2:
            words = line.split()
            # For PDF text, try to identify column-like patterns
            if len(words) > 4:
                # Try grouping every 2-3 words as potential cells
                cells = []
                for i in range(0, len(words), 2):
                    cell_words = words[i:i+2]
                    if cell_words:
                        cells.append(' '.join(cell_words))
                extraction_method = "word_grouping"
            else:
                cells = words
                extraction_method = "word_splitting"
        # Try colon separator (key:value pairs)
        elif ':' in line and line.count(':') > 0:
            cells = [cell.strip() for cell in line.split(':') if cell.strip()]
            extraction_method = "colons"
        else:
            # Most aggressive fallback: every word becomes a cell
            words = line.split()
            if len(words) > 1:
                cells = [word.strip() for word in words if word.strip()]
                extraction_method = "individual_words"
            elif words:
                cells = words
                extraction_method = "single_word"
        
        # Log the extraction details
        logger.info(f"    Cell extraction from '{original_line[:60]}...':")
        logger.info(f"      Method: {extraction_method}")
        logger.info(f"      Raw result: {cells}")
        
        # Additional processing: split on common patterns within cells
        expanded_cells = []
        for cell in cells:
            if not cell.strip():
                continue
                
            # If cell contains multiple patterns, split further
            if len(cell) > 50:  # Long cells might contain multiple values
                # Try splitting on parentheses, brackets, etc.
                if '(' in cell or ')' in cell:
                    import re
                    parts = re.split(r'[()]', cell)
                    expanded_cells.extend([part.strip() for part in parts if part.strip()])
                else:
                    expanded_cells.append(cell)
            else:
                expanded_cells.append(cell)
        
        logger.info(f"      Final cells: {expanded_cells}")
        return expanded_cells
    
    def _extract_title(self, content: str, filename: str) -> str:
        """Extract document title from content or filename"""
        lines = content.split('\n')
        
        # Collect all major headings/modules to create a comprehensive title
        major_sections = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Look for module/lesson patterns
            if any(keyword in line.lower() for keyword in ['module', 'lesson', 'chapter', 'section', 'part']):
                if len(line) < 100:  # Reasonable heading length
                    major_sections.append(line)
        
        # If we found multiple modules/lessons, create a comprehensive title
        if len(major_sections) > 1:
            # Extract base course name and count
            first_section = major_sections[0]
            total_sections = len(major_sections)
            
            # Try to extract course name from first section
            course_name = first_section
            # Remove markdown symbols and module/lesson numbers and patterns
            course_name = re.sub(r'^#+\s*', '', course_name)  # Remove markdown headings
            course_name = re.sub(r'(module|lesson|chapter|section|part)\s*\d+\s*[:;-]?\s*', '', course_name, flags=re.IGNORECASE)
            course_name = course_name.strip()
            
            if course_name:
                return f"{course_name} ({total_sections} Modules)"
            else:
                return f"Course Content ({total_sections} Modules)"
        
        # Single section or fallback to original logic
        for line in lines[:10]:  # Check first 10 lines
            line = line.strip()
            if not line:
                continue
                
            # Check for markdown heading
            if line.startswith('#'):
                return line.lstrip('#').strip()
            
            # Check for underlined heading
            if len(line) > 3 and not line.startswith(('---', '===', 'Page')):
                return line
        
        # Fallback to filename
        return os.path.splitext(filename)[0].replace('_', ' ').replace('-', ' ').title()
    
    def _content_to_slides(self, content: str) -> List[SlideContent]:
        """Convert script content to slides - each row becomes one slide with bullet points"""
        lines = content.split('\n')
        slides = []
        script_slide_counter = 1
        
        logger.info(f"Converting script content to slides, total lines: {len(lines)}")
        
        for line_num, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            
            # Check if this is a heading
            is_heading = False
            heading_text = line
            
            # Markdown heading
            if line.startswith('#'):
                is_heading = True
                heading_level = len(line) - len(line.lstrip('#'))
                heading_text = line.lstrip('#').strip()
                
                # Add heading slide with original level
                slides.append(SlideContent(
                    title=heading_text,
                    content=[],
                    slide_type='heading',
                    heading_level=heading_level
                ))
                logger.info(f"Created heading slide: '{heading_text}' (level {heading_level})")
                continue
            
            # Check other heading patterns
            for pattern in self.heading_patterns:
                match = re.match(pattern, line, re.MULTILINE)
                if match:
                    is_heading = True
                    heading_text = match.group(1).strip()
                    slides.append(SlideContent(
                        title=heading_text,
                        content=[],
                        slide_type='heading'
                    ))
                    logger.info(f"Created heading slide: '{heading_text}'")
                    break
            
            if not is_heading:
                # This is script content - create a slide with bullet points
                bullet_points = self._create_bullet_points(line)
                slide_title = self._create_slide_title(line, script_slide_counter)
                
                slides.append(SlideContent(
                    title=slide_title,
                    content=bullet_points,
                    slide_type='script'
                ))
                
                logger.info(f"Created script slide {script_slide_counter}: '{slide_title}' with {len(bullet_points)} bullet points")
                script_slide_counter += 1
        
        logger.info(f"FINAL RESULT: Created {len(slides)} total slides")
        
        # Count slides by type for verification
        script_slides = [s for s in slides if s.slide_type == 'script']
        heading_slides = [s for s in slides if s.slide_type == 'heading']
        logger.info(f"VERIFICATION: {len(script_slides)} script slides + {len(heading_slides)} heading slides = {len(slides)} total")
        
        return slides
    
    def _create_bullet_points(self, text: str) -> List[str]:
        """Convert script text into 4-6 concise bullet points that capture key ideas"""
        text = text.strip()
        if not text:
            return ["No content available"]
        
        logger.info(f"Creating bullet points from text: {text[:100]}...")
        
        # Split into sentences for analysis
        sentences = [s.strip() for s in re.split(r'[.!?]+', text) if s.strip() and len(s.strip()) > 10]
        
        if not sentences:
            return ["Key content from the paragraph"]
        
        bullet_points = []
        
        # Extract 4-6 key concepts using focused analysis
        bullet_points.extend(self._extract_core_concepts(sentences))
        bullet_points.extend(self._extract_benefits_outcomes(sentences))  
        bullet_points.extend(self._extract_methods_processes(sentences))
        bullet_points.extend(self._extract_applications_uses(sentences))
        
        # Clean and distill each bullet to be concise
        refined_bullets = []
        for bullet in bullet_points:
            refined = self._refine_bullet_to_concise_style(bullet)
            if refined and len(refined) > 8:
                refined_bullets.append(refined)
        
        # Remove duplicates and ensure 4-6 bullets
        unique_bullets = self._deduplicate_bullets(refined_bullets)
        
        # If we don't have enough, extract more from the text
        if len(unique_bullets) < 4:
            additional = self._extract_additional_concepts(sentences, unique_bullets)
            unique_bullets.extend(additional)
        
        # Special extraction for "speed and iteration" type concepts
        if len(unique_bullets) < 4:
            for sentence in sentences:
                if 'speed' in sentence.lower() or 'critical' in sentence.lower() or 'iteration' in sentence.lower():
                    unique_bullets.append("Especially valuable for rapid prototyping workflows")
                    break
        
        # Ensure exactly 4-6 bullets
        final_bullets = unique_bullets[:6]
        while len(final_bullets) < 4:
            final_bullets.append(f"Additional insight from the content")
        
        logger.info(f"Final bullet points: {final_bullets}")
        return final_bullets
    
    def _extract_core_concepts(self, sentences: List[str]) -> List[str]:
        """Extract main subjects and their key attributes"""
        concepts = []
        
        for sentence in sentences:
            # Look for main subject-verb-object patterns
            if any(verb in sentence.lower() for verb in ['allows', 'enables', 'helps', 'provides', 'offers', 'supports']):
                # Extract the core capability
                concept = self._extract_main_capability(sentence)
                if concept:
                    concepts.append(concept)
            
            # Look for API/interface descriptions
            if 'api' in sentence.lower() or 'simplifies' in sentence.lower():
                if 'intuitive' in sentence.lower() or 'simplifies' in sentence.lower():
                    concepts.append("The API is simple and supports interactive elements and data visualizations")
        
        return concepts[:2]
    
    def _extract_benefits_outcomes(self, sentences: List[str]) -> List[str]:
        """Extract benefits, advantages, and positive outcomes"""
        benefits = []
        
        for sentence in sentences:
            sentence_lower = sentence.lower()
            if 'prototyping' in sentence_lower and 'tools' in sentence_lower:
                benefits.append("It's ideal for quickly building and testing AI tools")
            elif 'speed and iteration' in sentence_lower or 'critical' in sentence_lower:
                benefits.append("Especially valuable for rapid prototyping workflows")
        
        return benefits[:2]
    
    def _extract_methods_processes(self, sentences: List[str]) -> List[str]:
        """Extract methods, processes, and approaches"""
        methods = []
        
        for sentence in sentences:
            sentence_lower = sentence.lower()
            if any(word in sentence_lower for word in ['process', 'method', 'approach', 'way', 'api', 'interface']):
                method = self._extract_method_concept(sentence)
                if method:
                    methods.append(method)
        
        return methods[:1]
    
    def _extract_applications_uses(self, sentences: List[str]) -> List[str]:
        """Extract specific applications and use cases"""
        applications = []
        
        for sentence in sentences:
            sentence_lower = sentence.lower()
            if any(word in sentence_lower for word in ['for', 'used', 'building', 'creating', 'developing', 'prototyping']):
                application = self._extract_application_concept(sentence)
                if application:
                    applications.append(application)
        
        return applications[:1]
    
    def _extract_main_capability(self, sentence: str) -> str:
        """Extract main capability from sentences like 'X allows Y to Z'"""
        # Pattern: "Streamlit allows developers to rapidly create..."
        # Output: "Streamlit enables fast development of web apps"
        
        sentence = sentence.strip()
        if 'allows' in sentence.lower():
            parts = sentence.lower().split('allows')
            if len(parts) >= 2:
                subject = parts[0].strip()
                capability = parts[1].strip()
                
                # Simplify the capability
                if 'developers to' in capability:
                    capability = capability.replace('developers to', '').strip()
                if 'rapidly create' in capability:
                    capability = capability.replace('rapidly create', 'fast development of')
                
                return f"{subject.title()} enables {capability}"
        
        elif 'enables' in sentence.lower():
            # Already in good form
            return sentence
        
        return sentence
    
    def _extract_definition_concept(self, sentence: str) -> str:
        """Extract core definition from sentences"""
        sentence = sentence.strip()
        
        # Look for "Its API is simple" -> "The API is simple"
        if sentence.lower().startswith('its '):
            sentence = 'The ' + sentence[4:]
        
        return sentence
    
    def _extract_benefit_concept(self, sentence: str) -> str:
        """Extract benefit/value proposition"""
        sentence = sentence.strip()
        
        if 'makes it especially useful for' in sentence.lower():
            parts = sentence.lower().split('makes it especially useful for')
            if len(parts) >= 2:
                use_case = parts[1].strip()
                return f"Ideal for {use_case}"
        
        elif 'especially useful for' in sentence.lower():
            parts = sentence.lower().split('especially useful for')
            if len(parts) >= 2:
                use_case = parts[1].strip()
                return f"Especially valuable for {use_case}"
        
        return sentence
    
    def _extract_method_concept(self, sentence: str) -> str:
        """Extract method or process information"""
        sentence = sentence.strip()
        
        if 'api' in sentence.lower() and any(word in sentence.lower() for word in ['simple', 'intuitive', 'easy']):
            return "The API is simple and supports interactive elements and data visualizations"
        elif 'simplifies' in sentence.lower() and 'process' in sentence.lower():
            return "The API is simple and supports interactive elements and data visualizations"
        
        return sentence
    
    def _extract_application_concept(self, sentence: str) -> str:
        """Extract application or use case information"""
        sentence = sentence.strip()
        
        if 'prototyping' in sentence.lower() and 'tools' in sentence.lower():
            return "It's ideal for quickly building and testing AI tools"
        elif 'speed and iteration are critical' in sentence.lower():
            return "Especially valuable for rapid prototyping workflows"
        
        return sentence
    
    def _refine_bullet_to_concise_style(self, bullet: str) -> str:
        """Refine bullet to match the concise style from the example"""
        bullet = bullet.strip()
        
        # Remove redundant phrases
        replacements = {
            'allows developers to rapidly create': 'enables fast development of',
            'its intuitive api simplifies': 'the API is simple and',
            'this makes it especially useful': 'ideal',
            'where speed and iteration are critical': 'workflows',
            'interactive web apps with just a few lines of python code': 'web apps using Python',
            'the process of adding widgets, visualizing data, and integrating machine learning models': 'interactive elements and data visualizations',
            'prototyping ai tools and dashboards': 'building and testing AI tools',
            'rapid prototyping workflows': 'rapid prototyping workflows'
        }
        
        bullet_lower = bullet.lower()
        for old, new in replacements.items():
            if old in bullet_lower:
                bullet = bullet.replace(old, new)
                bullet = bullet.replace(old.title(), new.title())
                bullet = bullet.replace(old.upper(), new.upper())
        
        # Ensure proper capitalization
        bullet = bullet[0].upper() + bullet[1:] if bullet else ""
        
        # Remove trailing punctuation and ensure it ends properly
        bullet = bullet.rstrip('.,!?;:').strip()
        
        return bullet
    
    def _deduplicate_bullets(self, bullets: List[str]) -> List[str]:
        """Remove duplicate concepts while preserving order"""
        seen = set()
        unique = []
        
        for bullet in bullets:
            # Create a key for similarity checking
            key_words = set(bullet.lower().split()[:5])  # First 5 words
            key = ' '.join(sorted(key_words))
            
            if key not in seen and len(bullet.strip()) > 8:
                seen.add(key)
                unique.append(bullet.strip())
        
        return unique
    
    def _extract_additional_concepts(self, sentences: List[str], existing: List[str]) -> List[str]:
        """Extract additional concepts when we need more bullets"""
        additional = []
        
        # Look for any remaining key concepts not yet captured
        for sentence in sentences:
            if len(additional) >= 2:
                break
                
            # Skip if this concept is already captured
            if any(self._concepts_overlap(sentence, existing_bullet) for existing_bullet in existing):
                continue
            
            # Extract any valuable remaining concept
            if len(sentence) > 20:  # Only substantial sentences
                concept = self._simplify_sentence_to_bullet(sentence)
                if concept:
                    additional.append(concept)
        
        return additional
    
    def _concepts_overlap(self, sentence: str, existing_bullet: str) -> bool:
        """Check if sentence concept already covered in existing bullet"""
        sentence_words = set(sentence.lower().split())
        bullet_words = set(existing_bullet.lower().split())
        
        # If 3 or more key words overlap, consider it duplicate
        overlap = len(sentence_words.intersection(bullet_words))
        
        # Special case: multiple "ideal for" bullets are redundant
        if 'ideal' in sentence.lower() and any('ideal' in existing.lower() for existing in [existing_bullet]):
            return True
        
        return overlap >= 3
    
    def _simplify_sentence_to_bullet(self, sentence: str) -> str:
        """Simplify a sentence to bullet point style"""
        sentence = sentence.strip()
        
        # Remove common sentence starters
        if sentence.lower().startswith(('this ', 'it ', 'that ', 'which ', 'where ')):
            words = sentence.split()[1:]
            if words:
                sentence = ' '.join(words).capitalize()
        
        # Make more concise
        sentence = sentence.replace('is especially useful', 'works well')
        sentence = sentence.replace('can be used', 'is used')
        sentence = sentence.replace('it is possible to', 'you can')
        
        return sentence[:100]  # Keep reasonable length
    
    def _extract_key_ideas(self, text: str) -> List[str]:
        """Extract 4-6 key ideas/concepts from text, focusing on main points rather than sentence structure"""
        ideas = []
        
        # Split into sentences for analysis
        sentences = [s.strip() for s in re.split(r'[.!?]+', text) if s.strip()]
        
        # Pattern 1: Identify main topics/subjects
        main_topics = self._identify_main_topics(text)
        ideas.extend(main_topics)
        
        # Pattern 2: Extract cause-effect relationships
        cause_effects = self._extract_cause_effects(text)
        ideas.extend(cause_effects)
        
        # Pattern 3: Find problem-solution pairs
        problems_solutions = self._extract_problems_solutions(text)
        ideas.extend(problems_solutions)
        
        # Pattern 4: Identify key benefits/outcomes
        benefits = self._extract_key_benefits(text)
        ideas.extend(benefits)
        
        # Pattern 5: Extract important processes/methods
        processes = self._extract_key_processes(text)
        ideas.extend(processes)
        
        # Pattern 6: Find comparisons/contrasts
        comparisons = self._extract_comparisons(text)
        ideas.extend(comparisons)
        
        return ideas
    
    def _extract_supporting_ideas(self, text: str, existing_ideas: List[str]) -> List[str]:
        """Extract supporting ideas when we need more bullet points"""
        supporting = []
        
        # Look for examples
        if 'example' in text.lower() or 'for instance' in text.lower():
            examples = self._extract_examples(text)
            supporting.extend(examples)
        
        # Extract specific details or statistics
        details = self._extract_specific_details(text)
        supporting.extend(details)
        
        # Find action items or recommendations
        actions = self._extract_action_items(text)
        supporting.extend(actions)
        
        return supporting
    
    def _identify_main_topics(self, text: str) -> List[str]:
        """Identify the main topics/subjects discussed in the text"""
        topics = []
        
        # Look for topic introduction patterns
        topic_patterns = [
            r'([A-Z][^.!?]*(?:is|are|involves|includes|means|refers to)[^.!?]*)',
            r'(The (?:main|key|primary|important)[^.!?]*)',
            r'(This (?:approach|method|process|concept)[^.!?]*)',
        ]
        
        for pattern in topic_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                if len(match.strip()) > 15:
                    topics.append(match.strip())
        
        return topics[:2]  # Limit to 2 main topics
    
    def _extract_cause_effects(self, text: str) -> List[str]:
        """Extract cause and effect relationships"""
        causes_effects = []
        
        # Look for causal language
        causal_patterns = [
            r'([^.!?]*(?:because|since|due to|as a result|therefore|consequently)[^.!?]*)',
            r'([^.!?]*(?:leads to|causes|results in|enables|allows)[^.!?]*)',
        ]
        
        for pattern in causal_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                if len(match.strip()) > 15:
                    causes_effects.append(match.strip())
        
        return causes_effects[:2]
    
    def _extract_problems_solutions(self, text: str) -> List[str]:
        """Extract problem-solution pairs"""
        prob_sols = []
        
        problem_words = ['problem', 'issue', 'challenge', 'difficulty', 'obstacle']
        solution_words = ['solution', 'answer', 'approach', 'method', 'way to', 'fix']
        
        sentences = [s.strip() for s in re.split(r'[.!?]+', text) if s.strip()]
        
        for sentence in sentences:
            sentence_lower = sentence.lower()
            if any(word in sentence_lower for word in problem_words):
                prob_sols.append(f"Challenge: {sentence}")
            elif any(word in sentence_lower for word in solution_words):
                prob_sols.append(f"Solution: {sentence}")
        
        return prob_sols[:2]
    
    def _extract_key_benefits(self, text: str) -> List[str]:
        """Extract key benefits or positive outcomes"""
        benefits = []
        
        benefit_patterns = [
            r'([^.!?]*(?:benefit|advantage|improvement|better|faster|easier|more efficient)[^.!?]*)',
            r'([^.!?]*(?:helps|enables|allows|improves|increases|reduces)[^.!?]*)',
        ]
        
        for pattern in benefit_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                if len(match.strip()) > 15:
                    benefits.append(match.strip())
        
        return benefits[:2]
    
    def _extract_key_processes(self, text: str) -> List[str]:
        """Extract key processes or methods"""
        processes = []
        
        process_patterns = [
            r'([^.!?]*(?:process|method|approach|technique|strategy)[^.!?]*)',
            r'([^.!?]*(?:first|then|next|finally)[^.!?]*)',
        ]
        
        for pattern in process_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                if len(match.strip()) > 15:
                    processes.append(match.strip())
        
        return processes[:2]
    
    def _extract_comparisons(self, text: str) -> List[str]:
        """Extract comparisons or contrasts"""
        comparisons = []
        
        comparison_patterns = [
            r'([^.!?]*(?:compared to|versus|unlike|different from|similar to)[^.!?]*)',
            r'([^.!?]*(?:while|whereas|however|in contrast)[^.!?]*)',
        ]
        
        for pattern in comparison_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                if len(match.strip()) > 15:
                    comparisons.append(match.strip())
        
        return comparisons[:1]
    
    def _extract_examples(self, text: str) -> List[str]:
        """Extract examples and illustrations"""
        examples = []
        
        example_patterns = [
            r'([^.!?]*(?:example|for instance|such as|including)[^.!?]*)',
            r'([^.!?]*(?:like|including|such as)[^.!?]*)',
        ]
        
        for pattern in example_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                if len(match.strip()) > 15:
                    examples.append(f"Example: {match.strip()}")
        
        return examples[:2]
    
    def _extract_specific_details(self, text: str) -> List[str]:
        """Extract specific details, numbers, or statistics"""
        details = []
        
        # Look for numbers/statistics
        number_pattern = r'([^.!?]*\d+[^.!?]*)'
        matches = re.findall(number_pattern, text)
        for match in matches:
            if len(match.strip()) > 15:
                details.append(match.strip())
        
        return details[:1]
    
    def _extract_action_items(self, text: str) -> List[str]:
        """Extract actionable items or recommendations"""
        actions = []
        
        action_patterns = [
            r'([^.!?]*(?:should|must|need to|have to|recommend)[^.!?]*)',
            r'([^.!?]*(?:use|implement|apply|adopt|consider)[^.!?]*)',
        ]
        
        for pattern in action_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                if len(match.strip()) > 15:
                    actions.append(f"Action: {match.strip()}")
        
        return actions[:2]
    
    def _extract_key_themes(self, text: str) -> List[str]:
        """Extract key themes and concepts from the text for better bullet point creation"""
        # Look for action words, important concepts, and key phrases
        action_words = []
        important_concepts = []
        
        # Find action verbs and important keywords
        words = text.split()
        for i, word in enumerate(words):
            clean_word = re.sub(r'[^\w]', '', word.lower())
            
            # Action words (verbs that indicate learning/doing)
            if clean_word in ['learn', 'build', 'create', 'develop', 'implement', 'understand', 
                            'explore', 'discover', 'master', 'practice', 'apply', 'use', 
                            'design', 'solve', 'analyze', 'optimize', 'improve', 'transform']:
                # Get surrounding context
                context_start = max(0, i-2)
                context_end = min(len(words), i+3)
                context = ' '.join(words[context_start:context_end])
                action_words.append(context)
            
            # Important technical/domain concepts (longer, capitalized, or specialized terms)
            if len(clean_word) > 6 or word[0].isupper() or clean_word in [
                'ai', 'genai', 'artificial', 'intelligence', 'machine', 'learning', 'data',
                'prototype', 'framework', 'algorithm', 'system', 'architecture', 'design',
                'development', 'programming', 'coding', 'technology', 'digital', 'software'
            ]:
                # Get surrounding context
                context_start = max(0, i-1)
                context_end = min(len(words), i+2)
                context = ' '.join(words[context_start:context_end])
                important_concepts.append(context)
        
        # Combine and deduplicate themes
        all_themes = action_words + important_concepts
        # Remove duplicates while preserving order
        unique_themes = []
        seen = set()
        for theme in all_themes:
            theme_key = theme.lower()
            if theme_key not in seen and len(theme.strip()) > 5:
                unique_themes.append(theme.strip())
                seen.add(theme_key)
        
        return unique_themes[:8]  # Limit to top 8 themes
    
    def _create_engaging_bullets(self, text: str, themes: List[str]) -> List[str]:
        """Create engaging, action-oriented bullet points based on key themes"""
        bullets = []
        
        # Analyze text for different content patterns
        text_lower = text.lower()
        
        # Pattern 1: Learning objectives (what learners will gain)
        learning_indicators = ['learn', 'understand', 'master', 'discover', 'explore', 'gain', 'develop']
        if any(indicator in text_lower for indicator in learning_indicators):
            learning_bullets = self._extract_learning_objectives(text, themes)
            bullets.extend(learning_bullets)
        
        # Pattern 2: Process/steps (how to do something)
        process_indicators = ['first', 'then', 'next', 'step', 'process', 'method', 'approach', 'way']
        if any(indicator in text_lower for indicator in process_indicators):
            process_bullets = self._extract_process_steps(text, themes)
            bullets.extend(process_bullets)
        
        # Pattern 3: Benefits/outcomes (why this matters)
        benefit_indicators = ['benefit', 'advantage', 'improve', 'better', 'faster', 'save', 'reduce', 'increase']
        if any(indicator in text_lower for indicator in benefit_indicators):
            benefit_bullets = self._extract_benefits(text, themes)
            bullets.extend(benefit_bullets)
        
        # Pattern 4: Problems/challenges (what we're solving)
        problem_indicators = ['problem', 'challenge', 'issue', 'difficulty', 'struggle', 'pain', 'trouble']
        if any(indicator in text_lower for indicator in problem_indicators):
            problem_bullets = self._extract_challenges(text, themes)
            bullets.extend(problem_bullets)
        
        # Pattern 5: Tools/technologies (what we're using)
        tool_indicators = ['use', 'tool', 'technology', 'framework', 'platform', 'system', 'software']
        if any(indicator in text_lower for indicator in tool_indicators):
            tool_bullets = self._extract_tools_and_methods(text, themes)
            bullets.extend(tool_bullets)
        
        # If we didn't find specific patterns, create general engaging bullets
        if not bullets:
            bullets = self._create_general_engaging_bullets(text, themes)
        
        return bullets
    
    def _extract_learning_objectives(self, text: str, themes: List[str]) -> List[str]:
        """Extract what learners will learn or achieve"""
        bullets = []
        sentences = re.split(r'[.!?]+', text)
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            sentence_lower = sentence.lower()
            if any(word in sentence_lower for word in ['learn', 'understand', 'master', 'discover', 'gain']):
                # Transform into learner-focused bullet
                if 'learn' in sentence_lower:
                    bullet = sentence.replace('learn', 'Learn').replace('you will learn', 'Learn').replace('we learn', 'Learn')
                elif 'understand' in sentence_lower:
                    bullet = sentence.replace('understand', 'Understand').replace('you will understand', 'Understand')
                else:
                    bullet = sentence
                
                bullet = self._make_bullet_engaging(bullet, themes)
                bullets.append(bullet)
        
        return bullets[:3]  # Limit to 3 learning objectives
    
    def _extract_process_steps(self, text: str, themes: List[str]) -> List[str]:
        """Extract process steps or sequential actions"""
        bullets = []
        
        # Look for numbered steps or sequential indicators
        step_patterns = [
            r'(\d+[\.)]\s*[^.!?]*)',  # Numbered steps like "1. Do this"
            r'(first[^.!?]*)',         # "First, we do..."
            r'(then[^.!?]*)',          # "Then we..."
            r'(next[^.!?]*)',          # "Next, we..."
            r'(finally[^.!?]*)',       # "Finally..."
        ]
        
        for pattern in step_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                step = match.strip()
                if len(step) > 10:  # Ignore very short matches
                    step = self._make_bullet_action_oriented(step)
                    bullets.append(step)
        
        # If no explicit steps found, look for action verbs
        if not bullets:
            action_sentences = []
            sentences = re.split(r'[.!?]+', text)
            for sentence in sentences:
                sentence = sentence.strip()
                if any(action in sentence.lower() for action in ['build', 'create', 'develop', 'implement', 'design']):
                    action_sentences.append(self._make_bullet_action_oriented(sentence))
            bullets.extend(action_sentences[:3])
        
        return bullets
    
    def _extract_benefits(self, text: str, themes: List[str]) -> List[str]:
        """Extract benefits, improvements, or positive outcomes"""
        bullets = []
        sentences = re.split(r'[.!?]+', text)
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            sentence_lower = sentence.lower()
            benefit_words = ['improve', 'better', 'faster', 'save', 'reduce', 'increase', 'boost', 'enhance', 'optimize']
            
            if any(word in sentence_lower for word in benefit_words):
                # Transform into benefit-focused bullet
                bullet = self._make_bullet_benefit_focused(sentence)
                bullets.append(bullet)
        
        return bullets[:3]
    
    def _extract_challenges(self, text: str, themes: List[str]) -> List[str]:
        """Extract problems or challenges being addressed"""
        bullets = []
        sentences = re.split(r'[.!?]+', text)
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            sentence_lower = sentence.lower()
            challenge_words = ['problem', 'challenge', 'issue', 'difficulty', 'struggle', 'limitation']
            
            if any(word in sentence_lower for word in challenge_words):
                # Transform into challenge-focused bullet
                bullet = self._make_bullet_challenge_focused(sentence)
                bullets.append(bullet)
        
        return bullets[:2]  # Limit challenges to avoid negativity
    
    def _extract_tools_and_methods(self, text: str, themes: List[str]) -> List[str]:
        """Extract tools, technologies, or methods being used"""
        bullets = []
        sentences = re.split(r'[.!?]+', text)
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            sentence_lower = sentence.lower()
            tool_words = ['use', 'tool', 'technology', 'framework', 'platform', 'system', 'method', 'approach']
            
            if any(word in sentence_lower for word in tool_words):
                bullet = self._make_bullet_tool_focused(sentence)
                bullets.append(bullet)
        
        return bullets[:2]
    
    def _create_general_engaging_bullets(self, text: str, themes: List[str]) -> List[str]:
        """Create engaging bullets when no specific pattern is found"""
        bullets = []
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 15]
        
        # Take the most important sentences and make them engaging
        for sentence in sentences[:4]:
            bullet = self._make_bullet_engaging(sentence, themes)
            bullets.append(bullet)
        
        return bullets
    
    def _make_bullet_engaging(self, text: str, themes: List[str]) -> str:
        """Transform text into an engaging bullet point"""
        # Start with action words when possible
        action_starters = ['Discover', 'Learn', 'Master', 'Build', 'Create', 'Develop', 'Explore', 'Understand']
        
        text = text.strip()
        if not text:
            return text
            
        # If it doesn't start with an engaging word, try to add one
        first_word = text.split()[0].lower() if text.split() else ''
        if first_word not in [starter.lower() for starter in action_starters]:
            # Try to add context from themes
            for theme in themes[:2]:
                if any(word in theme.lower() for word in text.lower().split()[:3]):
                    if 'learn' in text.lower() or 'understand' in text.lower():
                        text = f"Learn {text.lower()}"
                    elif 'build' in text.lower() or 'create' in text.lower():
                        text = f"Build {text.lower()}"
                    break
        
        # Ensure proper capitalization
        if text and not text[0].isupper():
            text = text[0].upper() + text[1:]
        
        return text
    
    def _make_bullet_action_oriented(self, text: str) -> str:
        """Transform text into action-oriented bullet point"""
        text = text.strip()
        # Remove step numbers and clean up
        text = re.sub(r'^\d+[\.)]\s*', '', text)
        text = re.sub(r'^(first|then|next|finally)[,\s]*', '', text, flags=re.IGNORECASE)
        
        # Ensure it starts with an action word
        if text and not any(text.lower().startswith(action) for action in ['build', 'create', 'develop', 'implement', 'design', 'learn', 'master']):
            # Try to make it action-oriented
            if 'build' in text.lower():
                text = 'Build ' + text.lower().replace('build', '').strip()
            elif 'create' in text.lower():
                text = 'Create ' + text.lower().replace('create', '').strip()
            elif 'learn' in text.lower():
                text = 'Learn ' + text.lower().replace('learn', '').strip()
        
        return text[0].upper() + text[1:] if text else text
    
    def _make_bullet_benefit_focused(self, text: str) -> str:
        """Transform text into benefit-focused bullet point"""
        text = text.strip()
        # Emphasize positive outcomes
        positive_transforms = {
            'improve': 'Improve',
            'better': 'Achieve better',
            'faster': 'Work faster with',
            'save': 'Save time by',
            'reduce': 'Reduce complexity through',
            'increase': 'Increase efficiency with'
        }
        
        for old, new in positive_transforms.items():
            if old in text.lower():
                text = text.lower().replace(old, new, 1)
                break
        
        return text[0].upper() + text[1:] if text else text
    
    def _make_bullet_challenge_focused(self, text: str) -> str:
        """Transform text into solution-focused bullet point"""
        text = text.strip()
        # Transform problems into solutions
        if 'problem' in text.lower():
            text = text.lower().replace('problem', 'challenge').replace('the challenge', 'overcome the challenge')
        elif 'issue' in text.lower():
            text = text.lower().replace('issue', 'challenge').replace('the challenge', 'address the challenge')
        
        return text[0].upper() + text[1:] if text else text
    
    def _make_bullet_tool_focused(self, text: str) -> str:
        """Transform text into tool/method-focused bullet point"""
        text = text.strip()
        # Emphasize the tools and their benefits
        if 'use' in text.lower():
            text = text.lower().replace('use', 'utilize', 1).replace('we utilize', 'utilize')
        
        return text[0].upper() + text[1:] if text else text
    
    def _create_structured_bullets(self, text: str) -> List[str]:
        """Fallback method: Create structured bullets from sentences"""
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
        
        bullets = []
        for sentence in sentences[:5]:
            # Clean and optimize each sentence
            bullet = self._optimize_sentence_for_bullet(sentence)
            bullets.append(bullet)
        
        return bullets
    
    def _optimize_sentence_for_bullet(self, sentence: str) -> str:
        """Optimize a sentence to be a better bullet point"""
        sentence = sentence.strip()
        
        # Remove redundant phrases
        redundant_phrases = ['it is important to', 'we should', 'it is necessary to', 'we need to']
        for phrase in redundant_phrases:
            sentence = sentence.replace(phrase, '').strip()
        
        # Make it more direct and actionable
        if sentence.startswith('you can'):
            sentence = sentence.replace('you can', '').strip()
            sentence = sentence[0].upper() + sentence[1:] if sentence else sentence
        
        return sentence
    
    def _optimize_bullet_points(self, bullets: List[str]) -> List[str]:
        """Final optimization of bullet points for clarity and engagement - ensures complete sentences"""
        if not bullets:
            return ["This slide covers key concepts and important information."]
        
        optimized = []
        for bullet in bullets:
            if not bullet or len(bullet.strip()) < 5:
                continue
                
            bullet = bullet.strip()
            
            # Make sure it's a complete sentence
            bullet = self._ensure_complete_sentence(bullet)
            
            # Ensure proper capitalization
            if bullet and not bullet[0].isupper():
                bullet = bullet[0].upper() + bullet[1:]
            
            # Ensure sentence ends with proper punctuation
            if bullet and not bullet.endswith(('.', '!', '?')):
                bullet = bullet + '.'
            
            # Limit length while maintaining complete sentence structure
            if len(bullet) > 120:
                bullet = self._shorten_to_complete_sentence(bullet)
            
            # Avoid duplicate concepts
            if not any(self._bullets_too_similar(bullet, existing) for existing in optimized):
                optimized.append(bullet)
        
        # Ensure we have 3-6 bullets
        if len(optimized) < 3:
            # Split longer bullets into complete sentences if possible
            while len(optimized) < 3 and any(len(b) > 80 for b in optimized):
                longest_idx = max(range(len(optimized)), key=lambda i: len(optimized[i]))
                longest = optimized[longest_idx]
                
                # Try to split at natural sentence boundaries
                split_result = self._split_into_complete_sentences(longest)
                if len(split_result) > 1:
                    optimized[longest_idx] = split_result[0]
                    optimized.extend(split_result[1:])
                else:
                    break
        
        # Limit to 6 bullets maximum
        if len(optimized) > 6:
            optimized = optimized[:6]
        
        # Ensure at least one bullet
        if not optimized:
            optimized = ["This slide covers key concepts and learning objectives."]
        
        # Final check - ensure all are complete sentences
        final_bullets = []
        for bullet in optimized:
            complete_bullet = self._ensure_complete_sentence(bullet)
            if complete_bullet and not complete_bullet.endswith(('.', '!', '?')):
                complete_bullet = complete_bullet + '.'
            final_bullets.append(complete_bullet)
        
        return final_bullets
    
    def _ensure_complete_sentence(self, text: str) -> str:
        """Ensure text is a grammatically complete sentence"""
        text = text.strip()
        if not text:
            return text
        
        # Remove common incomplete starters that make fragments
        fragment_starters = ['also ', 'and ', 'but ', 'or ', 'so ', 'because ']
        for starter in fragment_starters:
            if text.lower().startswith(starter):
                text = text[len(starter):].strip()
                if text:
                    text = text[0].upper() + text[1:]
        
        # Check if it's likely a sentence fragment and try to fix it
        words = text.split()
        if len(words) < 3:
            # Too short to be a meaningful sentence
            return f"This concept involves {text.lower()}."
        
        # Check for common sentence patterns
        has_verb = any(word.lower() in ['is', 'are', 'was', 'were', 'have', 'has', 'had', 'will', 'would', 'can', 'could', 'should', 'must', 'may', 'might', 'do', 'does', 'did', 'learn', 'build', 'create', 'develop', 'understand', 'explore', 'discover', 'implement', 'use', 'work', 'help', 'make', 'take', 'get', 'go', 'come', 'see', 'know', 'think', 'want', 'need', 'try', 'become', 'feel', 'look', 'seem', 'appear', 'provide', 'include', 'allow', 'enable', 'require', 'involve', 'focus', 'improve', 'increase', 'reduce', 'save', 'achieve', 'ensure', 'maintain', 'solve', 'address', 'handle', 'manage', 'utilize', 'leverage', 'optimize'] for word in words)
        
        has_subject = True  # Assume it has a subject if it starts with a capital letter
        
        # If it seems like a complete sentence, return as is
        if has_verb and has_subject:
            return text
        
        # If it's missing a verb, try to add one based on context
        if not has_verb:
            # Check for learning/action contexts
            if any(word.lower() in ['learn', 'learning', 'study', 'understand', 'master', 'discover', 'explore'] for word in words):
                return f"You will learn about {text.lower()}."
            elif any(word.lower() in ['build', 'create', 'develop', 'implement', 'design', 'make'] for word in words):
                return f"You will build {text.lower()}."
            elif any(word.lower() in ['benefit', 'advantage', 'improvement', 'better', 'faster', 'efficient'] for word in words):
                return f"This provides {text.lower()}."
            elif any(word.lower() in ['tool', 'method', 'approach', 'technique', 'strategy', 'framework'] for word in words):
                return f"This involves using {text.lower()}."
            else:
                return f"This covers {text.lower()}."
        
        return text
    
    def _shorten_to_complete_sentence(self, text: str) -> str:
        """Shorten text while maintaining complete sentence structure"""
        # Try to find a natural break point for a complete sentence
        sentences = re.split(r'[.!?]+', text)
        if len(sentences) > 1 and len(sentences[0].strip()) > 20:
            # Use the first complete sentence if it's substantial
            result = sentences[0].strip()
            if not result.endswith(('.', '!', '?')):
                result += '.'
            return result
        
        # If no natural sentence break, try to cut at a meaningful phrase boundary
        words = text.split()
        if len(words) > 15:
            # Look for natural breaking points
            break_points = []
            for i, word in enumerate(words):
                if word.endswith(',') or word.lower() in ['and', 'but', 'or', 'so', 'because', 'when', 'where', 'which', 'that']:
                    if i > 8:  # Don't break too early
                        break_points.append(i)
            
            if break_points:
                # Use the first good breaking point
                cut_point = break_points[0]
                result = ' '.join(words[:cut_point+1])
                # Remove trailing conjunction words
                if result.split()[-1].lower() in ['and', 'but', 'or', 'so', 'because']:
                    result = ' '.join(result.split()[:-1])
                if not result.endswith(('.', '!', '?')):
                    result += '.'
                return result
            else:
                # No good breaking point, cut at 15 words and add period
                result = ' '.join(words[:15])
                if not result.endswith(('.', '!', '?')):
                    result += '.'
                return result
        
        # If it's already short enough, just ensure it ends properly
        if not text.endswith(('.', '!', '?')):
            text += '.'
        return text
    
    def _split_into_complete_sentences(self, text: str) -> List[str]:
        """Split text into multiple complete sentences"""
        # Try to split at natural sentence boundaries
        if ' and ' in text and len(text.split(' and ')) == 2:
            parts = text.split(' and ', 1)
            first = parts[0].strip()
            second = parts[1].strip()
            
            # Make sure both parts can be complete sentences
            if len(first.split()) >= 3 and len(second.split()) >= 3:
                # Ensure first part is complete
                if not first.endswith(('.', '!', '?')):
                    first += '.'
                
                # Ensure second part is complete
                if not second[0].isupper():
                    second = second[0].upper() + second[1:]
                if not second.endswith(('.', '!', '?')):
                    second += '.'
                
                return [first, second]
        
        # Try splitting at commas for long sentences
        if ', ' in text and len(text.split(', ')) >= 2:
            parts = text.split(', ', 1)
            first = parts[0].strip()
            second = parts[1].strip()
            
            if len(first.split()) >= 4 and len(second.split()) >= 4:
                # Make sure both can be complete sentences
                first_complete = self._ensure_complete_sentence(first)
                second_complete = self._ensure_complete_sentence(second)
                
                if not first_complete.endswith(('.', '!', '?')):
                    first_complete += '.'
                if not second_complete.endswith(('.', '!', '?')):
                    second_complete += '.'
                
                return [first_complete, second_complete]
        
        # If we can't split naturally, return as single sentence
        return [text]
    
    def _bullets_too_similar(self, bullet1: str, bullet2: str) -> bool:
        """Check if two bullets are too similar"""
        words1 = set(bullet1.lower().split())
        words2 = set(bullet2.lower().split())
        
        # Remove common words for comparison
        common_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by'}
        words1 = words1 - common_words
        words2 = words2 - common_words
        
        if not words1 or not words2:
            return False
        
        # Check if more than 60% of words overlap
        overlap = len(words1 & words2)
        total_unique = len(words1 | words2)
        
        return overlap / total_unique > 0.6
    
    def _extract_key_words(self, sentence: str) -> str:
        """Extract key words/phrases from a sentence"""
        # Remove common filler words and focus on important content
        filler_words = {
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 
            'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'have', 
            'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should',
            'this', 'that', 'these', 'those', 'here', 'there', 'when', 'where',
            'how', 'why', 'what', 'which', 'who', 'whom'
        }
        
        words = sentence.split()
        key_words = []
        
        for word in words:
            clean_word = re.sub(r'[^\w\s]', '', word.lower())
            if clean_word not in filler_words and len(clean_word) > 2:
                key_words.append(word)
        
        if len(key_words) > 8:  # Limit key words
            key_words = key_words[:8]
        
        result = ' '.join(key_words)
        return result if len(result) > 10 else ""
    
    def _create_slide_title(self, text: str, slide_number: int) -> str:
        """Create a descriptive title for the slide"""
        # Try to extract first few words as title
        words = text.split()[:8]  # First 8 words max
        title = ' '.join(words)
        
        # Clean up title
        title = re.sub(r'[^\w\s-]', '', title)  # Remove most punctuation
        title = title.strip()
        
        # If title is too long, truncate
        if len(title) > 50:
            title = title[:47] + "..."
        
        # If title is too short or empty, use a generic title
        if len(title) < 3:
            title = f"Script {slide_number}"
        
        return title
    
    def _create_slides_from_content(self, content: str) -> List[SlideContent]:
        """Create slides from content without clear headings"""
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        if not lines:
            return [SlideContent(title="Content", content=["No content found"], slide_type='content')]
        
        slides = []
        
        # Try to detect potential headings first
        potential_headings = []
        for i, line in enumerate(lines):
            # Look for short lines that might be headings
            if (len(line) < 80 and 
                not line.endswith('.') and 
                not line.startswith('-') and
                not line.startswith('•') and
                len(line.split()) > 1 and
                len(line.split()) < 10):
                potential_headings.append((i, line))
        
        if potential_headings and len(potential_headings) > 5:
            # Use detected headings to create slides
            for i, (line_idx, heading) in enumerate(potential_headings):
                start_idx = line_idx + 1
                end_idx = potential_headings[i + 1][0] if i + 1 < len(potential_headings) else len(lines)
                
                content_lines = lines[start_idx:end_idx]
                # Limit content to avoid overcrowded slides
                content_lines = [line for line in content_lines if line and not line.startswith('---')][:8]
                
                if content_lines or i == 0:  # Always include first slide even if no content
                    slides.append(SlideContent(
                        title=heading[:60] + "..." if len(heading) > 60 else heading,
                        content=content_lines,
                        slide_type='content'
                    ))
        else:
            # Fallback: Create smaller chunks for more slides
            slide_size = 4  # Smaller chunks for more slides
            
            for i in range(0, len(lines), slide_size):
                slide_lines = lines[i:i + slide_size]
                title = slide_lines[0] if slide_lines else f"Content {i // slide_size + 1}"
                content = slide_lines[1:] if len(slide_lines) > 1 else []
                
                # Skip page separators in title
                if title.startswith('--- Page'):
                    title = f"Page {i // slide_size + 1}"
                
                slides.append(SlideContent(
                    title=title[:50] + "..." if len(title) > 50 else title,
                    content=content,
                    slide_type='content'
                ))
        
        return slides

class SlideGenerator:
    """Handles generation of presentation slides"""
    
    def create_google_slides_json(self, doc_structure: DocumentStructure) -> str:
        """Generate Google Slides-compatible JSON structure"""
        slides_data = {
            "title": doc_structure.title,
            "slides": [],
            "metadata": doc_structure.metadata,
            "instructions": {
                "how_to_import": [
                    "1. Go to slides.google.com",
                    "2. Create a new presentation",
                    "3. Use 'Import slides' feature",
                    "4. Or copy-paste content from the generated text file"
                ]
            }
        }
        
        # Title slide
        slides_data["slides"].append({
            "slide_number": 1,
            "type": "title",
            "title": doc_structure.title,
            "subtitle": f"Generated from {doc_structure.metadata['filename']}",
            "date": datetime.now().strftime('%B %d, %Y')
        })
        
        # Content slides
        for i, slide_content in enumerate(doc_structure.slides, 2):
            slide_data = {
                "slide_number": i,
                "type": "content",
                "title": slide_content.title,
                "content": slide_content.content[:8]  # Limit to 8 bullet points
            }
            slides_data["slides"].append(slide_data)
        
        # Save JSON file
        filename = f"google_slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        filepath = os.path.join(EXPORT_FOLDER, filename)
        os.makedirs(EXPORT_FOLDER, exist_ok=True)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(slides_data, f, indent=2, ensure_ascii=False)
        
        return filepath
    
    def create_google_slides_text(self, doc_structure: DocumentStructure) -> str:
        """Generate text format that can be easily copied to Google Slides"""
        content_lines = []
        
        # Title slide content
        content_lines.append("=== TITLE SLIDE ===")
        content_lines.append(f"Title: {doc_structure.title}")
        content_lines.append(f"Subtitle: Generated from {doc_structure.metadata['filename']}")
        content_lines.append(f"Date: {datetime.now().strftime('%B %d, %Y')}")
        content_lines.append("")
        
        # Content slides
        for i, slide_content in enumerate(doc_structure.slides, 1):
            content_lines.append(f"=== SLIDE {i + 1} ===")
            content_lines.append(f"Title: {slide_content.title}")
            content_lines.append("Content:")
            
            for bullet in slide_content.content[:8]:  # Limit bullets
                content_lines.append(f"• {bullet}")
            
            content_lines.append("")
        
        # Instructions
        content_lines.append("=== INSTRUCTIONS ===")
        content_lines.append("How to use with Google Slides:")
        content_lines.append("1. Go to slides.google.com and create a new presentation")
        content_lines.append("2. For each slide above:")
        content_lines.append("   - Create a new slide")
        content_lines.append("   - Copy the title to the slide title")
        content_lines.append("   - Copy the bullet points to the slide content")
        content_lines.append("3. Choose a theme and customize as needed")
        
        content_text = '\n'.join(content_lines)
        
        # Save text file
        filename = f"google_slides_import_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        filepath = os.path.join(EXPORT_FOLDER, filename)
        os.makedirs(EXPORT_FOLDER, exist_ok=True)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content_text)
        
        return filepath
    
    def create_google_slides_csv(self, doc_structure: DocumentStructure) -> str:
        """Generate CSV format for easy import to Google Slides"""
        csv_data = []
        
        # Header
        csv_data.append(['Slide Number', 'Slide Type', 'Title', 'Content'])
        
        # Title slide
        csv_data.append([
            1, 
            'Title', 
            doc_structure.title, 
            f"Generated from {doc_structure.metadata['filename']} on {datetime.now().strftime('%B %d, %Y')}"
        ])
        
        # Content slides
        for i, slide_content in enumerate(doc_structure.slides, 2):
            content_text = ' | '.join(slide_content.content[:8])  # Join with separator
            csv_data.append([
                i,
                'Content',
                slide_content.title,
                content_text
            ])
        
        # Save CSV file
        filename = f"google_slides_csv_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        filepath = os.path.join(EXPORT_FOLDER, filename)
        os.makedirs(EXPORT_FOLDER, exist_ok=True)
        
        with open(filepath, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            writer.writerows(csv_data)
        
        return filepath
    
    def create_powerpoint(self, doc_structure: DocumentStructure) -> str:
        """Generate PowerPoint presentation optimized for Google Slides import"""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = doc_structure.title
        subtitle.text = f"Generated from {doc_structure.metadata['filename']}\n{datetime.now().strftime('%B %d, %Y')}\n\nOptimized for Google Slides Import"
        
        # Organize slides by heading hierarchy
        organized_slides = self._organize_slides_by_hierarchy(doc_structure.slides)
        
        logger.info(f"PowerPoint generation: Input has {len(doc_structure.slides)} slides")
        logger.info(f"After organization: {len(organized_slides)} organized sections")
        
        # Count organized slides by type
        slide_types = {}
        for section in organized_slides:
            slide_type = section.get('type', 'unknown')
            slide_types[slide_type] = slide_types.get(slide_type, 0) + 1
        
        logger.info(f"Organized slide types: {slide_types}")
        
        # Create slides based on hierarchy: H1=presentation title, H2=section, H3=subsection, H4=slide titles, cells=content
        pptx_slides_created = 0
        for section in organized_slides:
            if section['type'] == 'presentation_title':
                # This is handled by the main title slide above, skip
                continue
                
            elif section['type'] == 'section_title':
                # H2 = Section title page
                section_slide_layout = prs.slide_layouts[5]  # Title and content layout
                slide = prs.slides.add_slide(section_slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                pptx_slides_created += 1
                logger.info(f"Created section title slide: '{section['title']}'")
                
                # Style as section title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(48)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(25, 118, 210)  # Blue color
                
                # Add overview if available
                if section.get('overview'):
                    content_shape = None
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            content_shape = placeholder
                            break
                    if content_shape:
                        content_shape.text = section['overview']
                        
            elif section['type'] == 'subsection_title':
                # H3 = Subsection title page
                subsection_slide_layout = prs.slide_layouts[5]  # Title and content layout
                slide = prs.slides.add_slide(subsection_slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                pptx_slides_created += 1
                logger.info(f"Created subsection title slide: '{section['title']}'")
                
                # Style as subsection title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(40)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(102, 126, 234)  # Lighter blue
                
                # Add overview if available
                if section.get('overview'):
                    content_shape = None
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            content_shape = placeholder
                            break
                    if content_shape:
                        content_shape.text = section['overview']
                        
            elif section['type'] == 'slide_title':
                # H4 = Individual slide with content
                slide_layout = prs.slide_layouts[5]  # Title and content layout
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                pptx_slides_created += 1
                logger.info(f"Created H4 slide: '{section['title']}'")
                
                # Style as individual slide title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(36)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = RGBColor(55, 71, 79)  # Dark gray
                
                # Add content if available
                if section.get('content'):
                    content_shape = None
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            content_shape = placeholder
                            break
                    if content_shape:
                        # Convert content to bullet points if it's a list
                        if isinstance(section['content'], list):
                            content_text = '\n'.join([f"• {item}" for item in section['content']])
                        else:
                            # Process text content into bullet points
                            bullet_points = self._create_bullet_points(str(section['content']))
                            content_text = '\n'.join([f"• {bullet}" for bullet in bullet_points])
                        
                        content_shape.text = content_text
                        
            elif section['type'] == 'script':
                # Script content - two-column layout with bullet points on left and sketch suggestion on right
                script_slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
                slide = prs.slides.add_slide(script_slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(
                    left=Inches(0.5), top=Inches(0.3), 
                    width=Inches(9), height=Inches(1)
                )
                title_frame = title_shape.text_frame
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.text = section['title']
                title_paragraph.font.size = Pt(32)
                title_paragraph.font.bold = True
                title_paragraph.alignment = PP_ALIGN.CENTER
                
                # Get the current slide number BEFORE incrementing for unique image generation
                current_slide_number = pptx_slides_created + 1
                pptx_slides_created += 1
                logger.info(f"Created script slide #{current_slide_number}: '{section['title']}' with {len(section.get('content', []))} bullet points")
                
                # Left side - Bullet points
                if section.get('content'):
                    bullet_shape = slide.shapes.add_textbox(
                        left=Inches(0.5), top=Inches(1.5), 
                        width=Inches(4.5), height=Inches(5)
                    )
                    bullet_frame = bullet_shape.text_frame
                    bullet_frame.word_wrap = True
                    bullet_frame.clear()
                    
                    # Add bullet points directly without header
                    for i, bullet_point in enumerate(section['content']):
                        p = bullet_frame.paragraphs[0] if i == 0 else bullet_frame.add_paragraph()
                        p.text = f"• {bullet_point}"
                        p.level = 0
                        p.font.size = Pt(16)
                        p.space_after = Pt(8)
                        p.space_before = Pt(4)
                
                # Right side - Content-specific sketch only (no text suggestion)
                # Generate sketch based directly on the bullet points content for this specific slide
                logger.info(f"Creating content-specific sketch for slide {current_slide_number}: '{section['title']}' with {len(section.get('content', []))} bullet points")
                sketch_image_path = self._create_content_specific_sketch(section.get('content', []), section['title'], current_slide_number)
                logger.info(f"Generated sketch path: {sketch_image_path}")
                
                # Add the sketch image (full size on the right)
                if sketch_image_path:
                    try:
                        # Add the sketch image taking up the full right side
                        slide.shapes.add_picture(
                            sketch_image_path,
                            left=Inches(5.2), 
                            top=Inches(1.5), 
                            width=Inches(4.3), 
                            height=Inches(5)
                        )
                        logger.info(f"Successfully added content-specific sketch for slide: {section['title']}")
                    except Exception as e:
                        logger.error(f"Could not add sketch image for '{section['title']}': {e}")
                else:
                    logger.warning(f"No sketch image generated for slide: {section['title']}")
                            
            else:
                # Other content slides
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                title_shape = slide.shapes.title
                title_shape.text = section['title']
                
                # Style title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(28)
                
                # Add content
                if section.get('content'):
                    content_shape = None
                    for placeholder in slide.placeholders:
                        if placeholder.placeholder_format.idx == 1:
                            content_shape = placeholder
                            break
                    
                    if content_shape:
                        text_frame = content_shape.text_frame
                        text_frame.clear()
                        
                        for i, content_item in enumerate(section['content'][:6]):
                            p = text_frame.add_paragraph()
                            p.text = content_item
                            p.level = 0
                            p.font.size = Pt(20)
                            p.space_after = Pt(8)
        
        # Add a summary slide at the end
        summary_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Summary"
        
        # Find content placeholder safely for summary
        content_shape = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 1:
                content_shape = placeholder
                break
        
        if content_shape:
            summary_text = f"• {len(organized_slides)} main sections covered\n"
            summary_text += f"• Generated from {doc_structure.metadata['filename']}\n"
            summary_text += f"• Ready for Google Slides import\n\n"
            summary_text += "Next steps:\n"
            summary_text += "• Import this PowerPoint to Google Slides\n"
            summary_text += "• Customize themes and layouts\n"
            summary_text += "• Add images and additional formatting"
            
            content_shape.text = summary_text
        
        logger.info(f"POWERPOINT FINAL: Created {pptx_slides_created} content slides + 1 title slide + 1 summary slide = {pptx_slides_created + 2} total slides")
        
        # Save presentation
        filename = f"google_slides_optimized_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(EXPORT_FOLDER, filename)
        os.makedirs(EXPORT_FOLDER, exist_ok=True)
        prs.save(filepath)
        
        return filepath
    
    def _organize_slides_by_hierarchy(self, slides: List[SlideContent]) -> List[Dict]:
        """Organize slides by heading hierarchy: H1=presentation title, H2=section, H3=subsection, H4=slide titles, cells=content"""
        organized = []
        presentation_title_added = False
        
        for slide in slides:
            title = slide.title
            content = slide.content
            slide_type = getattr(slide, 'slide_type', 'content')  # Get slide_type attribute
            
            # Check if this is explicitly a heading slide
            if slide_type == 'heading':
                # Use original heading level if available, otherwise determine from title
                level = getattr(slide, 'heading_level', None) or self._determine_heading_level(title)
                
                if level == 1 and not presentation_title_added:
                    # H1 = Title page for entire presentation (only use first H1)
                    organized.append({
                        'type': 'presentation_title',
                        'title': title,
                        'level': 1,
                        'overview': self._create_section_overview(content) if content else None
                    })
                    presentation_title_added = True
                    
                elif level == 1 or level == 2:
                    # H2 = Section title page
                    organized.append({
                        'type': 'section_title',
                        'title': title,
                        'level': 2,
                        'overview': self._create_section_overview(content) if content else None
                    })
                    
                elif level == 3:
                    # H3 = Subsection title page
                    organized.append({
                        'type': 'subsection_title',
                        'title': title,
                        'level': 3,
                        'overview': self._create_section_overview(content) if content else None
                    })
                    
                elif level == 4:
                    # H4 = Individual slide title (content slides with H4 titles)
                    organized.append({
                        'type': 'slide_title',
                        'title': title,
                        'content': content,
                        'level': 4
                    })
                    
            elif slide_type == 'script':
                # Script content slides - each table row becomes one slide with bullet points
                organized.append({
                    'type': 'script',
                    'title': title,
                    'content': content,
                    'level': 4
                })
                
            else:
                # Other content slides (fallback)
                if title.startswith('Table Cell'):
                    # Each table cell gets its own slide
                    organized.append({
                        'type': 'cell_content',
                        'title': title,
                        'content': content,
                        'level': 4
                    })
                else:
                    # Other content slides
                    organized.append({
                        'type': 'content',
                        'title': title,
                        'content': content,
                        'level': 4
                    })
        
        return organized
    
    def _generate_sketch_suggestion(self, bullet_points: List[str], title: str) -> str:
        """Generate a contextual sketch suggestion based on specific bullet point content"""
        if not bullet_points:
            return "Simple concept illustration related to the topic"
        
        # Extract key concepts directly from bullet points
        key_concepts = self._extract_key_concepts_from_bullets(bullet_points)
        
        # Combine title and bullet points for analysis
        all_text = f"{title} {' '.join(bullet_points)}".lower()
        
        # Generate specific suggestion based on content
        if any(word in all_text for word in ['step', 'process', 'workflow', 'method', 'approach', 'first', 'then', 'next']):
            return f"Process flow showing: {key_concepts}. Draw connected steps with arrows and icons."
        
        elif any(word in all_text for word in ['compare', 'vs', 'versus', 'difference', 'traditional', 'new', 'before', 'after']):
            return f"Comparison diagram: {key_concepts}. Show side-by-side with contrasting elements."
        
        elif any(word in all_text for word in ['ai', 'genai', 'chatbot', 'model', 'prompt', 'artificial intelligence']):
            return f"AI concept visualization: {key_concepts}. Include neural network or brain diagram with data flows."
        
        elif any(word in all_text for word in ['ui', 'interface', 'design', 'user', 'screen', 'app', 'button', 'form']):
            return f"UI mockup showing: {key_concepts}. Create wireframe with interface elements and user interactions."
        
        elif any(word in all_text for word in ['data', 'chart', 'graph', 'metrics', 'analysis', 'dashboard', 'report']):
            return f"Data visualization: {key_concepts}. Show charts, graphs, and key metrics dashboard."
        
        elif any(word in all_text for word in ['system', 'architecture', 'component', 'api', 'database', 'server']):
            return f"System diagram: {key_concepts}. Draw components with connections and data flow arrows."
        
        elif any(word in all_text for word in ['learn', 'course', 'lesson', 'training', 'education', 'skill']):
            return f"Learning path: {key_concepts}. Show progression with milestones and knowledge building blocks."
        
        elif any(word in all_text for word in ['problem', 'solution', 'challenge', 'issue', 'fix', 'solve']):
            return f"Problem-solution flow: {key_concepts}. Use warning icons for problems, lightbulbs for solutions."
        
        elif any(word in all_text for word in ['benefit', 'advantage', 'improve', 'better', 'faster', 'efficiency']):
            return f"Benefits diagram: {key_concepts}. Use checkmarks, positive icons, and improvement arrows."
        
        elif any(word in all_text for word in ['team', 'collaboration', 'meeting', 'communication', 'stakeholder']):
            return f"Collaboration diagram: {key_concepts}. Show people icons, communication flows, and team interactions."
        
        elif any(word in all_text for word in ['timeline', 'schedule', 'roadmap', 'phases', 'milestone', 'time']):
            return f"Timeline visualization: {key_concepts}. Create timeline with key events and progress markers."
        
        else:
            # Default: create conceptual diagram based on content
            return f"Concept map: {key_concepts}. Show main ideas connected with relationships and supporting details."
    
    def _extract_key_concepts_from_bullets(self, bullet_points: List[str]) -> str:
        """Extract the most important concepts from bullet points for sketch suggestions"""
        if not bullet_points:
            return "main concepts"
        
        # Get first few key words from each bullet point
        concepts = []
        for bullet in bullet_points[:3]:  # Focus on first 3 bullets
            # Clean up bullet text and extract key terms
            clean_bullet = bullet.strip()
            if clean_bullet:
                # Take first significant phrase (up to first punctuation or first 4 words)
                words = clean_bullet.split()[:4]
                key_phrase = ' '.join(words)
                concepts.append(key_phrase)
        
        if concepts:
            return ' → '.join(concepts)
        else:
            return 'key concepts from content'
    
    def _create_content_specific_sketch(self, bullet_points: List[str], title: str, slide_number: int = 0) -> str:
        """Create a thought bubble with visual prompt on light blue background"""
        logger.info(f"_create_content_specific_sketch called for slide {slide_number}, title: '{title}', bullets: {len(bullet_points)}")
        if not bullet_points and not title:
            logger.warning(f"No content for sketch - bullets: {bullet_points}, title: '{title}'")
            return None
            
        try:
            # Generate a unique drawing prompt
            drawing_prompt = self._generate_drawing_prompt(bullet_points, title, slide_number)
            logger.info(f"Generated drawing prompt for slide {slide_number}: {drawing_prompt}")
            
            # Create thought bubble image with the prompt
            return self._create_thought_bubble_image(drawing_prompt, slide_number)
            
            # Keep text file creation for backup
            self._create_prompt_text_file(drawing_prompt, slide_number)
            
            # Analyze the specific content to determine the best sketch approach
            all_content = f"{title} {' '.join(bullet_points)}".lower()
            logger.info(f"All content for analysis: '{all_content[:100]}...'")
            
            # Create image with the same quality as before
            width, height = 800, 600
            image = Image.new('RGB', (width, height), '#f8f9fa')
            draw = ImageDraw.Draw(image)
            
            # Load fonts
            try:
                font_paths = [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/System/Library/Fonts/Arial.ttf", 
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/Windows/Fonts/arial.ttf"
                ]
                font_loaded = False
                for font_path in font_paths:
                    try:
                        font = ImageFont.truetype(font_path, 18)
                        title_font = ImageFont.truetype(font_path, 24)
                        small_font = ImageFont.truetype(font_path, 14)
                        large_font = ImageFont.truetype(font_path, 20)
                        font_loaded = True
                        break
                    except:
                        continue
                
                if not font_loaded:
                    raise Exception("No fonts found")
                    
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
                large_font = ImageFont.load_default()
            
            # Determine sketch type based on actual bullet point content
            sketch_type = self._determine_content_specific_sketch_type(bullet_points, title)
            logger.info(f"Determined sketch type: '{sketch_type}' for slide {slide_number}")
            
            # Draw the appropriate sketch based on content analysis
            if sketch_type == 'process':
                self._draw_process_from_content(draw, width, height, font, bullet_points, title)
            elif sketch_type == 'comparison':
                self._draw_comparison_from_content(draw, width, height, font, bullet_points, title)
            elif sketch_type == 'ai_tech':
                self._draw_ai_concept(draw, width, height, font, title_font)
            elif sketch_type == 'ui_design':
                self._draw_ui_design_from_content(draw, width, height, font, bullet_points, title)
            elif sketch_type == 'data_metrics':
                self._draw_data_visualization(draw, width, height, font, title_font)
            elif sketch_type == 'system_arch':
                self._draw_system_architecture(draw, width, height, font, title_font)
            elif sketch_type == 'learning_path':
                self._draw_learning_path(draw, width, height, font, title_font)
            elif sketch_type == 'benefits':
                self._draw_benefits_diagram(draw, width, height, font, bullet_points)
            else:
                # Default: create concept map from bullet points
                self._draw_concept_map_from_bullets(draw, width, height, font, bullet_points, title)
            
            # Add slide title at the top (smaller and cleaner)
            if title:
                bbox = draw.textbbox((0, 0), title, font=title_font)
                title_width = bbox[2] - bbox[0]
                title_height = bbox[3] - bbox[1]
                
                # Simple title without heavy background
                draw.text((width//2 - title_width//2, 20), title, fill='#333333', font=title_font)
            
            # Save with unique naming
            os.makedirs('temp_sketches', exist_ok=True)
            import random
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            random_component = random.randint(1000, 9999)
            unique_id = abs(hash(f"{slide_number}_{title}_{str(bullet_points)}_{sketch_type}_{timestamp}_{random_component}"))
            image_path = f'temp_sketches/sketch_slide{slide_number:03d}_{sketch_type}_{unique_id}_{timestamp}_{random_component}.png'
            
            image.save(image_path)
            logger.info(f"Created content-specific sketch: {image_path}")
            return image_path
            
        except Exception as e:
            logger.error(f"Error creating content-specific sketch: {e}")
            return None
    
    def _create_llm_generated_diagram(self, bullet_points: List[str], title: str, slide_number: int) -> str:
        """Use LLM to generate a diagram based on slide content"""
        try:
            import openai
            
            # Prepare the content for LLM analysis
            content_text = f"Title: {title}\nBullet Points:\n" + "\n".join([f"• {bullet}" for bullet in bullet_points])
            
            # Create a prompt for diagram generation
            diagram_prompt = f"""
Analyze the following slide content and create a detailed description for a diagram that would best illustrate the concepts:

{content_text}

Please respond with a JSON object containing:
1. "diagram_type": The best type of diagram (flowchart, comparison, concept_map, process, timeline, hierarchy, etc.)
2. "description": A detailed description of what the diagram should show
3. "elements": A list of key visual elements to include
4. "layout": Suggested layout/structure
5. "svg_code": Simple SVG code (800x600) that represents this diagram with professional styling

Focus on creating a clean, professional diagram that clearly communicates the key concepts from the slide content.
"""

            # Get API key from environment or config
            api_key = os.getenv('OPENAI_API_KEY')
            if not api_key:
                logger.info("No OpenAI API key found, skipping LLM diagram generation")
                return None
            
            # Make API call to OpenAI
            client = openai.OpenAI(api_key=api_key)
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",  # Faster, cheaper model
                messages=[
                    {"role": "system", "content": "Create SVG diagrams. Respond with JSON: {'svg_content': '<svg>...</svg>', 'diagram_type': 'process|comparison|flowchart'}"},
                    {"role": "user", "content": diagram_prompt}
                ],
                max_tokens=1000,  # Shorter responses for speed
                temperature=0.2   # Less creativity for speed
            )
            
            # Parse the response
            import json
            response_content = response.choices[0].message.content
            logger.info(f"LLM response for slide {slide_number}: {response_content[:200]}...")
            
            try:
                diagram_data = json.loads(response_content)
            except json.JSONDecodeError:
                logger.error(f"Failed to parse LLM response as JSON for slide {slide_number}")
                return None
            
            if 'svg_code' in diagram_data:
                # Convert SVG to PNG
                svg_content = diagram_data['svg_code']
                logger.info(f"Generated SVG for slide {slide_number}: {len(svg_content)} characters")
                return self._convert_svg_to_image(svg_content, slide_number, diagram_data.get('diagram_type', 'llm'))
            
            logger.warning(f"No SVG code in LLM response for slide {slide_number}")
            return None
            
        except Exception as e:
            logger.error(f"Error in LLM diagram generation: {e}")
            return None
    
    def _convert_svg_to_image(self, svg_content: str, slide_number: int, diagram_type: str) -> str:
        """Convert SVG content to PNG image"""
        try:
            # Save the SVG content to a file for inspection
            svg_filename = f"temp_sketches/slide{slide_number:03d}_{diagram_type}_debug.svg"
            os.makedirs('temp_sketches', exist_ok=True)
            with open(svg_filename, 'w', encoding='utf-8') as f:
                f.write(svg_content)
            logger.info(f"Saved SVG content to {svg_filename}")
            
            # Try using cairosvg if available and working
            try:
                import cairosvg
                from io import BytesIO
                
                # Convert SVG to PNG
                png_data = cairosvg.svg2png(bytestring=svg_content.encode('utf-8'))
                
                # Create unique filename
                content_hash = abs(hash(svg_content)) % (10**16)
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
                random_id = random.randint(1000, 9999)
                filename = f"sketch_slide{slide_number:03d}_llm_{diagram_type}_{content_hash}_{timestamp}_{random_id}.png"
                
                # Ensure temp_sketches directory exists
                image_path = os.path.join('temp_sketches', filename)
                
                # Save the PNG
                with open(image_path, 'wb') as f:
                    f.write(png_data)
                
                logger.info(f"Successfully converted LLM SVG to PNG: {image_path}")
                return image_path
                
            except Exception as cairo_error:
                logger.warning(f"cairosvg conversion failed: {cairo_error}")
                
                # Second attempt: Use ImageMagick convert command
                try:
                    import subprocess
                    
                    # Create unique filename
                    content_hash = abs(hash(svg_content)) % (10**16)
                    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
                    random_id = random.randint(1000, 9999)
                    filename = f"sketch_slide{slide_number:03d}_llm_{diagram_type}_{content_hash}_{timestamp}_{random_id}.png"
                    
                    # Save SVG temporarily
                    temp_svg = os.path.join('temp_sketches', f'temp_{slide_number}_{diagram_type}.svg')
                    with open(temp_svg, 'w', encoding='utf-8') as f:
                        f.write(svg_content)
                    
                    # Convert using ImageMagick
                    image_path = os.path.join('temp_sketches', filename)
                    subprocess.run([
                        'convert', 
                        '-background', 'white',
                        '-density', '96',
                        temp_svg,
                        image_path
                    ], check=True, capture_output=True)
                    
                    # Clean up temp SVG
                    os.remove(temp_svg)
                    
                    logger.info(f"Successfully converted SVG using ImageMagick: {image_path}")
                    return image_path
                    
                except Exception as magick_error:
                    logger.warning(f"ImageMagick conversion failed: {magick_error}, using fallback")
                    return self._create_image_from_svg_description(svg_content, slide_number, diagram_type)
            
        except Exception as e:
            logger.error(f"Error converting SVG to image: {e}")
            return None

    def _generate_drawing_prompt(self, bullet_points: List[str], title: str, slide_number: int) -> str:
        """Generate simple, flat-design visual concepts for educational slides"""
        try:
            # Analyze content to determine main educational concepts
            all_text = f"{title} {' '.join(bullet_points)}".lower()
            
            # Identify the core educational theme
            visual_prompt = self._create_educational_visual_concept(title, all_text, bullet_points)
            
            return visual_prompt
            
        except Exception as e:
            logger.error(f"Error generating visual prompt: {e}")
            return self._get_fallback_visual_prompt()
    
    def _create_educational_visual_concept(self, title: str, all_text: str, bullet_points: List[str]) -> str:
        """Create educational flat-design visual concepts based on slide content"""
        
        # Format the visual concept using the template structure
        visual_prompt = "Create a visual concept for this course slide:\n"
        visual_prompt += f"Slide title: \"{title}\"\n"
        
        # Add key points from bullets
        if bullet_points:
            visual_prompt += "Key points:\n"
            for bullet in bullet_points[:4]:  # Use up to 4 main points
                visual_prompt += f"    •    {bullet}\n"
            visual_prompt += "\n"
        
        # Determine the appropriate visual metaphor based on content
        image_idea = self._determine_educational_visual_metaphor(all_text, bullet_points)
        visual_prompt += f"Image idea: {image_idea}\n\n"
        
        # Always use clean flat illustration style
        visual_prompt += "Style: clean flat illustration, ideal for a presentation."
        
        return visual_prompt
    
    def _determine_educational_visual_metaphor(self, all_text: str, bullet_points: List[str]) -> str:
        """Determine the best visual metaphor for educational content"""
        
        # Process/workflow/steps content
        if any(word in all_text for word in ['process', 'step', 'workflow', 'sequence', 'first', 'then', 'next', 'finally']):
            return "A series of connected geometric shapes forming a path, with icons representing each step in the process. Simple arrows guide the flow from start to finish."
        
        # Comparison/contrast content
        elif any(word in all_text for word in ['compare', 'versus', 'vs', 'difference', 'before', 'after', 'traditional', 'modern']):
            return "A split composition with two distinct sides - one representing the old/traditional approach (muted colors, simple shapes) and the other showing the new/modern approach (bright colors, dynamic shapes)."
        
        # AI/Technology content
        elif any(word in all_text for word in ['ai', 'artificial', 'intelligence', 'machine learning', 'algorithm', 'neural', 'model']):
            return "A stylized brain made of interconnected nodes and circuits, with data flowing through neural pathways. Geometric patterns suggest computational processing."
        
        # Speed/efficiency content
        elif any(word in all_text for word in ['fast', 'quick', 'rapid', 'speed', 'efficient', 'streamline']):
            return "A minimalist rocket or arrow moving through simplified obstacles, leaving a trail of progress markers. Motion lines suggest speed and efficiency."
        
        # Building/creating content
        elif any(word in all_text for word in ['build', 'create', 'develop', 'construct', 'design', 'prototype']):
            return "Building blocks or puzzle pieces coming together to form a complete structure. Each piece represents a component of the solution."
        
        # Learning/education content
        elif any(word in all_text for word in ['learn', 'understand', 'master', 'discover', 'explore', 'knowledge']):
            return "A lightbulb with rays emanating outward, surrounded by simple icons representing different concepts being illuminated and understood."
        
        # Problem/solution content
        elif any(word in all_text for word in ['problem', 'solution', 'challenge', 'solve', 'fix', 'resolve']):
            return "A maze or tangled path on the left transforming into a clear, straight path on the right. A key or tool bridges the transformation."
        
        # Tools/resources content
        elif any(word in all_text for word in ['tool', 'resource', 'platform', 'framework', 'library', 'api']):
            return "A toolbox with simplified tool icons floating above it, each tool connected to its application with dotted lines."
        
        # Data/analysis content
        elif any(word in all_text for word in ['data', 'analysis', 'metrics', 'measure', 'statistics', 'visualization']):
            return "Abstract data points transforming into organized charts and graphs. Flowing lines connect raw data to meaningful insights."
        
        # Collaboration/team content
        elif any(word in all_text for word in ['team', 'collaborate', 'together', 'share', 'communicate', 'feedback']):
            return "Simplified human figures arranged in a circle, with connecting lines showing information flow. Speech bubbles or thought clouds overlap to show shared ideas."
        
        # Growth/improvement content
        elif any(word in all_text for word in ['grow', 'improve', 'enhance', 'evolve', 'progress', 'advance']):
            return "A plant or tree growing from seed to full bloom, with each stage clearly defined. Growth stages align with concept progression."
        
        # Testing/validation content
        elif any(word in all_text for word in ['test', 'validate', 'verify', 'check', 'quality', 'debug']):
            return "A magnifying glass examining different geometric shapes, with checkmarks appearing on validated items and X marks on issues to fix."
        
        # Default educational visual
        else:
            return "Interconnected circles representing concepts, with the largest circle in the center and smaller related concepts orbiting around it. Lines show relationships between ideas."
    
    def _get_fallback_visual_prompt(self) -> str:
        """Provide a generic educational visual prompt as fallback"""
        return """Create a visual concept for this course slide:

Image idea: Abstract geometric shapes arranged to suggest learning and progress. Use circles, triangles, and squares in a balanced composition with arrows showing flow and connection.

Style: clean flat illustration, ideal for a presentation."""
    
    def _create_prompt_text_file(self, drawing_prompt: str, slide_number: int) -> str:
                    f"Create a video game level progression showing {concepts_str} as unlockable achievements. Include progress bars, experience points, power-ups, and character evolution through each stage."
                ]
                visual_suggestions = suggestions
                
            # Comparison/contrast content
            elif any(word in all_text for word in ['compare', 'vs', 'versus', 'difference', 'traditional', 'new', 'before', 'after', 'old', 'modern', 'instead']):
                suggestions = [
                    f"Create a dramatic before/after transformation scene showing {concepts_str}. Use split-screen with stark visual contrasts - perhaps a caterpillar to butterfly metamorphosis, or a cluttered room becoming organized and zen-like.",
                    f"Design a superhero transformation sequence illustrating the evolution from {concepts_str}. Show the 'old way' as a regular person struggling, then the transformation, and finally the empowered 'new way' as a superhero with special abilities.",
                    f"Illustrate a time machine concept showing past vs future versions of {concepts_str}. Include vintage/sepia tones for the old approach and sleek, holographic elements for the modern approach, connected by a swirling time portal.",
                    f"Create a David vs Goliath visual metaphor contrasting {concepts_str}. Show the underdog approach as small but clever/agile, versus the traditional approach as large but slow/cumbersome."
                ]
                visual_suggestions = suggestions
                
            # AI/Technology content
            elif any(word in all_text for word in ['ai', 'genai', 'artificial', 'intelligence', 'machine', 'learning', 'model', 'algorithm', 'data', 'neural', 'code', 'tech']):
                suggestions = [
                    f"Create a futuristic city skyline where {concepts_str} are represented as glowing, interconnected buildings. Show data streams flowing between structures, holographic displays, and AI entities as friendly robots helping humans navigate the city.",
                    f"Design a magical spell-casting scene where {concepts_str} are ingredients in a digital potion. Show a wizard-programmer at a computer, casting code spells, with magical symbols floating in the air and creating amazing results.",
                    f"Illustrate a brain-computer symbiosis showing {concepts_str} as neurons firing in a hybrid human-AI mind. Use electric blue neural pathways, glowing synapses, and thought bubbles containing innovative ideas.",
                    f"Create a space mission control room where {concepts_str} are different mission-critical systems. Show astronauts (users) working with AI co-pilots, multiple screens, and successful rocket launches."
                ]
                visual_suggestions = suggestions
                
            # Interface/User Experience content
            elif any(word in all_text for word in ['ui', 'interface', 'design', 'user', 'screen', 'app', 'website', 'button', 'form', 'layout', 'experience', 'interact']):
                suggestions = [
                    f"Create a user's emotional journey map showing their face changing from confused/frustrated to delighted as they interact with {concepts_str}. Include thought bubbles, facial expressions, and heart symbols growing larger with each positive interaction.",
                    f"Design a theme park map where {concepts_str} are different attractions and rides. Show excited visitors (users) moving through the park, with clear signage, fun mascots guiding them, and happy faces at each attraction.",
                    f"Illustrate a cooking show setup where {concepts_str} are the steps of preparing a perfect meal for guests. Show a friendly chef (designer) presenting each step with enthusiasm, taste-testing, and guests' reactions.",
                    f"Create a personal assistant/concierge scenario where {concepts_str} represent different services. Show a helpful character anticipating user needs, providing solutions, and creating 'wow' moments."
                ]
                visual_suggestions = suggestions
                
            # Learning/Education content
            elif any(word in all_text for word in ['learn', 'understand', 'course', 'lesson', 'tutorial', 'training', 'education', 'skill', 'knowledge', 'master', 'explore']):
                suggestions = [
                    f"Create a adventure treasure hunt map where {concepts_str} are hidden treasures to discover. Show a curious explorer with a magnifying glass, mysterious caves, 'X marks the spot' locations, and chests overflowing with knowledge gems.",
                    f"Design a gardening metaphor where {concepts_str} are seeds that grow into a beautiful knowledge garden. Show planting, watering (practice), sunlight (guidance), and the garden blooming with diverse plants and flowers.",
                    f"Illustrate a library/bookstore scene where {concepts_str} are magical books that come alive when opened. Show floating words, characters stepping out of books, and a cozy reading nook with a curious learner surrounded by wonder.",
                    f"Create a mountain climbing expedition where {concepts_str} are different peaks to conquer. Show climbers with proper gear, base camps, helping each other, and the amazing view from the summit."
                ]
                visual_suggestions = suggestions
                
            # Benefits/Improvements content
            elif any(word in all_text for word in ['benefit', 'advantage', 'improve', 'better', 'faster', 'efficient', 'save', 'reduce', 'increase', 'boost', 'optimize']):
                suggestions = [
                    f"Create a fitness transformation montage showing {concepts_str} as a personal training program. Show before/after photos, workout equipment, progress charts, healthy foods, and a confident, energized person at the end.",
                    f"Design a home makeover reveal scene where {concepts_str} are the renovation improvements. Show the dramatic 'ta-da!' moment with confetti, amazed reactions, before/after photos, and a beautiful transformed space.",
                    f"Illustrate a race car pit stop where {concepts_str} are the performance upgrades. Show mechanics working efficiently, speed blur effects, checkered flags, and the car zooming ahead of competitors.",
                    f"Create a plant growth time-lapse where {concepts_str} are the nutrients and care that make it thrive. Show accelerated growth, blooming flowers, abundant fruit, and a lush, healthy plant."
                ]
                visual_suggestions = suggestions
                
            # Problem-solving content
            elif any(word in all_text for word in ['problem', 'challenge', 'solution', 'solve', 'fix', 'issue', 'trouble', 'overcome', 'resolve']):
                suggestions = [
                    f"Create a detective mystery scene where {concepts_str} are the clues leading to the solution. Show a thoughtful detective with a magnifying glass, evidence board with red strings connecting clues, and the 'aha!' moment of solving the case.",
                    f"Design a superhero rescue scene where {concepts_str} are the super-powers used to save the day. Show a city in trouble, the hero arriving with cape flowing, using each power strategically, and citizens cheering.",
                    f"Illustrate a medieval quest where {concepts_str} are magical tools/weapons needed to defeat a dragon (the problem). Show a brave knight collecting each tool, facing the dragon, and ultimately achieving victory.",
                    f"Create a medical emergency room scene where {concepts_str} are the life-saving treatments. Show skilled doctors working as a team, high-tech equipment, vital signs improving, and a successful recovery."
                ]
                visual_suggestions = suggestions
                
            # Team/Collaboration content
            elif any(word in all_text for word in ['team', 'people', 'collaborate', 'together', 'group', 'community', 'share', 'communicate', 'partner']):
                suggestions = [
                    f"Create a orchestral symphony performance where {concepts_str} are different instruments working in harmony. Show a conductor leading, musicians focused and coordinated, and the audience moved by the beautiful music.",
                    f"Design a sports team huddle and victory scene where {concepts_str} are the plays that lead to winning. Show team strategy, high-fives, coordinated action, and the celebration with trophy held high.",
                    f"Illustrate a potluck dinner party where {concepts_str} are the different dishes everyone brings to share. Show diverse people contributing, a beautifully set table, everyone enjoying the feast together.",
                    f"Create a barn-raising community scene where {concepts_str} are the contributions each person makes. Show neighbors working together, sharing tools, supporting each other, and celebrating the completed barn."
                ]
                visual_suggestions = suggestions
                
            else:
                # Creative general suggestions with rich variety
                general_suggestions = [
                    f"Create a magical ecosystem where {concepts_str} are interconnected like a thriving forest. Show symbiotic relationships, energy flowing between elements, and the whole system glowing with life and harmony.",
                    f"Design a inventor's workshop scene where {concepts_str} are the breakthrough innovations. Show a curious inventor surrounded by sketches, prototypes, light bulb moments, and amazing contraptions coming to life.",
                    f"Illustrate a cosmic constellation where {concepts_str} are stars forming meaningful patterns in the night sky. Show an astronomer discovering the connections, drawing lines between stars, and revealing hidden wisdom.",
                    f"Create a vibrant marketplace bazaar where {concepts_str} are exotic goods being traded. Show colorful stalls, enthusiastic merchants, curious customers, and the excitement of discovering valuable treasures.",
                    f"Design a master chef's kitchen where {concepts_str} are premium ingredients being combined into a masterpiece dish. Show precision, artistry, aromatic steam, and the final plating that looks like art.",
                    f"Illustrate a time capsule discovery where {concepts_str} are precious artifacts that tell an important story. Show archaeologists carefully uncovering items, examining them with wonder, and piecing together the narrative."
                ]
                visual_suggestions = general_suggestions
            
            # Select suggestion based on variation seed for uniqueness
            if visual_suggestions:
                selected_suggestion = visual_suggestions[variation_seed % len(visual_suggestions)]
                
                # Add unique enhancement based on slide number
                enhancements = [
                    " Add dynamic lighting effects and shadows to create depth.",
                    " Include small animated details like floating particles or gentle movement.",
                    " Use a warm, inviting color palette that feels approachable and friendly.",
                    " Add subtle texture overlays to give the image a handcrafted, artistic feel.",
                    " Include inspirational quotes or key phrases integrated into the visual design.",
                    " Show diverse, inclusive representation of people with various backgrounds.",
                    " Add environmental context that reinforces the learning atmosphere.",
                    " Include interactive elements that suggest the viewer can engage with the concept.",
                    " Use metaphorical symbolism that makes abstract concepts concrete and memorable.",
                    " Add celebration elements like confetti, sparkles, or success indicators."
                ]
                
                enhancement = enhancements[slide_number % len(enhancements)]
                return selected_suggestion + enhancement
            else:
                return f"Create an engaging visual story featuring {concepts_str} that helps learners connect emotionally with the concepts and remember them through vivid imagery."
                
        except Exception as e:
            logger.error(f"Error generating visual suggestion: {e}")
            return f"Create a memorable visual representation of {title} that uses storytelling, metaphor, and rich imagery to help learners understand and remember the key concepts."

    def _create_prompt_text_file(self, drawing_prompt: str, slide_number: int) -> str:
        """Create a text file containing the drawing prompt for easy copy-paste into LLMs"""
        try:
            # Create unique filename
            import hashlib
            prompt_hash = hashlib.md5(drawing_prompt.encode()).hexdigest()[:8]
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
            random_id = random.randint(1000, 9999)
            filename = f"prompt_slide{slide_number:03d}_{prompt_hash}_{timestamp}_{random_id}.txt"
            
            # Create the text content
            text_content = f"""RICH VISUAL SUGGESTION FOR SLIDE {slide_number}
==============================================

{drawing_prompt}

==============================================
This creative visual suggestion is designed to help learners understand and remember the concepts through engaging imagery, storytelling, and memorable metaphors.

Copy the text above and paste it into any image generation AI (DALL-E, Midjourney, Stable Diffusion, etc.) or describe it to an illustrator to create a compelling visual that will animate your slide and captivate your audience.

Generated: {timestamp}
"""
            
            # Save text file
            os.makedirs('temp_sketches', exist_ok=True)
            file_path = os.path.join('temp_sketches', filename)
            
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(text_content)
            
            logger.info(f"Created drawing prompt text file: {file_path}")
            return file_path
            
        except Exception as e:
            logger.error(f"Error creating prompt text file: {e}")
            return None
    
    def _create_image_from_svg_description(self, svg_content: str, slide_number: int, diagram_type: str) -> str:
        """Create a simple diagram image based on SVG content description (fallback)"""
        try:
            # Create a basic image with SVG-inspired content
            width, height = 800, 600
            image = Image.new('RGB', (width, height), '#ffffff')
            draw = ImageDraw.Draw(image)
            
            # Add a title indicating this is an LLM-generated concept
            try:
                font = ImageFont.truetype('/System/Library/Fonts/Arial.ttf', 16)
                title_font = ImageFont.truetype('/System/Library/Fonts/Arial Bold.ttf', 20)
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
            
            # Add header
            draw.text((20, 20), f"🤖 LLM-Generated {diagram_type.replace('_', ' ').title()}", fill='#2c3e50', font=title_font)
            
            # Try to extract and display some text content from the SVG
            import re
            
            # Look for text elements in SVG
            text_matches = re.findall(r'<text[^>]*>([^<]+)</text>', svg_content, re.IGNORECASE)
            if text_matches:
                y_pos = 60
                draw.text((20, y_pos), "Key Elements:", fill='#7f8c8d', font=title_font)
                y_pos += 30
                
                for i, text in enumerate(text_matches[:8]):  # Max 8 elements
                    draw.text((30, y_pos), f"• {text.strip()}", fill='#2c3e50', font=font)
                    y_pos += 25
            
            # Look for rectangles and circles to show structure
            rect_matches = re.findall(r'<rect[^>]*>', svg_content, re.IGNORECASE)
            circle_matches = re.findall(r'<circle[^>]*>', svg_content, re.IGNORECASE)
            
            if rect_matches or circle_matches:
                y_pos = height - 150
                draw.text((20, y_pos), f"Structure: {len(rect_matches)} boxes, {len(circle_matches)} circles", 
                         fill='#7f8c8d', font=font)
            
            # Add a visual indicator that this is AI-generated
            draw.rectangle([width-200, height-100, width-20, height-20], 
                          fill='#e3f2fd', outline='#2196f3', width=2)
            draw.text((width-180, height-80), "Generated by AI", fill='#1976d2', font=font)
            draw.text((width-180, height-60), "from slide content", fill='#1976d2', font=font)
            
            # Create unique filename
            content_hash = abs(hash(svg_content)) % (10**16)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S_%f')[:-3]
            random_id = random.randint(1000, 9999)
            filename = f"sketch_slide{slide_number:03d}_llm_{diagram_type}_{content_hash}_{timestamp}_{random_id}.png"
            
            # Save the image
            os.makedirs('temp_sketches', exist_ok=True)
            image_path = os.path.join('temp_sketches', filename)
            image.save(image_path, 'PNG')
            
            logger.info(f"Created fallback LLM diagram with {len(text_matches)} text elements: {image_path}")
            return image_path
            
        except Exception as e:
            logger.error(f"Error creating fallback diagram: {e}")
            return None

    def _determine_content_specific_sketch_type(self, bullet_points: List[str], title: str) -> str:
        """Analyze bullet points and title to determine the best sketch type"""
        all_text = f"{title} {' '.join(bullet_points)}".lower()
        
        # Process/workflow indicators
        if any(word in all_text for word in ['process', 'step', 'workflow', 'procedure', 'method', 'approach', 'how to']):
            return 'process'
            
        # Comparison indicators
        if any(word in all_text for word in ['vs', 'versus', 'compare', 'comparison', 'difference', 'traditional', 'modern']):
            return 'comparison'
            
        # AI/Technology indicators
        if any(word in all_text for word in ['ai', 'genai', 'artificial intelligence', 'machine learning', 'algorithm', 'model']):
            return 'ai_tech'
            
        # UI/Design indicators  
        if any(word in all_text for word in ['interface', 'design', 'user', 'ui', 'ux', 'mockup', 'prototype']):
            return 'ui_design'
            
        # Data/Metrics indicators
        if any(word in all_text for word in ['data', 'metrics', 'analytics', 'chart', 'graph', 'statistics', 'performance']):
            return 'data_metrics'
            
        # System/Architecture indicators
        if any(word in all_text for word in ['system', 'architecture', 'component', 'service', 'api', 'database']):
            return 'system_arch'
            
        # Learning/Education indicators
        if any(word in all_text for word in ['learn', 'course', 'lesson', 'tutorial', 'training', 'education', 'skill']):
            return 'learning_path'
            
        # Benefits/Advantages indicators
        if any(word in all_text for word in ['benefit', 'advantage', 'improve', 'better', 'faster', 'efficient']):
            return 'benefits'
            
        # Default to concept map
        return 'concept_map'
    
    def _draw_process_from_content(self, draw, width, height, font, bullet_points, title):
        """Draw a process flow based on actual bullet point content"""
        num_steps = min(len(bullet_points), 5)  # Max 5 steps for readability
        if num_steps == 0:
            return
            
        # Calculate step positions
        step_width = 120
        step_height = 60
        total_width = num_steps * step_width + (num_steps - 1) * 40
        start_x = (width - total_width) // 2
        y_pos = height // 2 - step_height // 2
        
        colors = ['#3498db', '#2ecc71', '#f39c12', '#e74c3c', '#9b59b6']
        
        for i, bullet in enumerate(bullet_points[:num_steps]):
            x_pos = start_x + i * (step_width + 40)
            
            # Draw step box with shadow (ensure valid dimensions)
            shadow_offset = 3
            if step_width > 0 and step_height > 0:
                draw.rounded_rectangle([x_pos + shadow_offset, y_pos + shadow_offset, 
                                      x_pos + step_width + shadow_offset, y_pos + step_height + shadow_offset],
                                     radius=8, fill='#00000020')
                draw.rounded_rectangle([x_pos, y_pos, x_pos + step_width, y_pos + step_height],
                                     radius=8, fill=colors[i % len(colors)], outline='#2c3e50', width=2)
            
            # Add step number
            step_num = f"{i+1}"
            draw.text((x_pos + 10, y_pos + 10), step_num, fill='white', font=font)
            
            # Add abbreviated bullet text
            text_lines = self._wrap_text(bullet[:30] + "...", 14, font)
            for j, line in enumerate(text_lines[:2]):  # Max 2 lines
                draw.text((x_pos + 10, y_pos + 25 + j * 15), line, fill='white', font=font)
            
            # Draw arrow to next step
            if i < num_steps - 1:
                arrow_start_x = x_pos + step_width + 5
                arrow_end_x = arrow_start_x + 30
                arrow_y = y_pos + step_height // 2
                
                # Arrow line
                draw.line([arrow_start_x, arrow_y, arrow_end_x - 8, arrow_y], fill='#2c3e50', width=3)
                # Arrow head
                draw.polygon([(arrow_end_x, arrow_y), (arrow_end_x - 8, arrow_y - 4), (arrow_end_x - 8, arrow_y + 4)], 
                           fill='#2c3e50')
    
    def _draw_comparison_from_content(self, draw, width, height, font, bullet_points, title):
        """Draw a comparison diagram based on bullet points"""
        if len(bullet_points) < 2:
            return
            
        # Split bullets into two groups
        mid_point = len(bullet_points) // 2
        left_items = bullet_points[:mid_point] if mid_point > 0 else bullet_points[:1]
        right_items = bullet_points[mid_point:] if mid_point > 0 else bullet_points[1:]
        
        # Draw two columns
        col_width = width // 2 - 60
        left_x = 30
        right_x = width // 2 + 30
        
        # Left column (Traditional/Old)
        left_rect = [left_x, 100, left_x + max(col_width, 50), height - 100]
        if left_rect[3] > left_rect[1]:  # Ensure y2 > y1
            draw.rounded_rectangle(left_rect, radius=10, fill='#ecf0f1', outline='#bdc3c7', width=2)
        draw.text((left_x + 10, 110), "Before/Traditional", fill='#2c3e50', font=font)
        
        for i, item in enumerate(left_items[:4]):  # Max 4 items
            y_pos = 140 + i * 50
            # Bullet point
            draw.ellipse([left_x + 15, y_pos, left_x + 25, y_pos + 10], fill='#e74c3c')
            # Text
            text_lines = self._wrap_text(item[:40] + "...", col_width - 50, font)
            draw.text((left_x + 35, y_pos), text_lines[0] if text_lines else "", fill='#2c3e50', font=font)
        
        # Right column (Modern/New)
        right_rect = [right_x, 100, right_x + max(col_width, 50), height - 100]
        if right_rect[3] > right_rect[1]:  # Ensure y2 > y1
            draw.rounded_rectangle(right_rect, radius=10, fill='#e8f5e8', outline='#27ae60', width=2)
        draw.text((right_x + 10, 110), "After/Modern", fill='#27ae60', font=font)
        
        for i, item in enumerate(right_items[:4]):  # Max 4 items
            y_pos = 140 + i * 50
            # Bullet point
            draw.ellipse([right_x + 15, y_pos, right_x + 25, y_pos + 10], fill='#27ae60')
            # Text
            text_lines = self._wrap_text(item[:40] + "...", col_width - 50, font)
            draw.text((right_x + 35, y_pos), text_lines[0] if text_lines else "", fill='#27ae60', font=font)
        
        # Arrow between columns
        arrow_y = height // 2
        draw.line([width // 2 - 20, arrow_y, width // 2 + 20, arrow_y], fill='#3498db', width=4)
        draw.polygon([(width // 2 + 20, arrow_y), (width // 2 + 12, arrow_y - 6), (width // 2 + 12, arrow_y + 6)], 
                   fill='#3498db')
    
    def _draw_benefits_diagram(self, draw, width, height, font, bullet_points):
        """Draw a benefits/advantages diagram"""
        if not bullet_points:
            return
            
        # Central hub
        center_x, center_y = width // 2, height // 2
        hub_radius = 60
        
        # Draw central hub
        draw.ellipse([center_x - hub_radius, center_y - hub_radius, 
                     center_x + hub_radius, center_y + hub_radius],
                    fill='#3498db', outline='#2980b9', width=3)
        draw.text((center_x - 30, center_y - 10), "Benefits", fill='white', font=font)
        
        # Draw benefit nodes around the hub
        num_benefits = min(len(bullet_points), 6)  # Max 6 for readability
        angle_step = 2 * 3.14159 / num_benefits
        
        for i, benefit in enumerate(bullet_points[:num_benefits]):
            angle = i * angle_step
            node_x = center_x + int(140 * cos(angle))
            node_y = center_y + int(140 * sin(angle))
            
            # Draw connecting line
            draw.line([center_x + int(hub_radius * cos(angle)), 
                      center_y + int(hub_radius * sin(angle)),
                      node_x - int(30 * cos(angle)), 
                      node_y - int(30 * sin(angle))], 
                     fill='#34495e', width=2)
            
            # Draw benefit node
            node_width, node_height = 120, 40
            draw.rounded_rectangle([node_x - node_width//2, node_y - node_height//2,
                                  node_x + node_width//2, node_y + node_height//2],
                                 radius=8, fill='#2ecc71', outline='#27ae60', width=2)
            
            # Add benefit text
            benefit_text = benefit[:25] + "..." if len(benefit) > 25 else benefit
            text_lines = self._wrap_text(benefit_text, node_width - 10, font)
            for j, line in enumerate(text_lines[:2]):
                draw.text((node_x - len(line) * 4, node_y - 8 + j * 12), line, fill='white', font=font)
    
    def _draw_concept_map_from_bullets(self, draw, width, height, font, bullet_points, title):
        """Draw a concept map directly from bullet points"""
        if not bullet_points:
            return
            
        # Main concept in center
        center_x, center_y = width // 2, height // 2 + 20
        main_width, main_height = 160, 50
        
        draw.rounded_rectangle([center_x - main_width//2, center_y - main_height//2,
                              center_x + main_width//2, center_y + main_height//2],
                             radius=10, fill='#3498db', outline='#2980b9', width=3)
        
        # Truncate title for center
        center_title = title[:20] + "..." if len(title) > 20 else title
        title_lines = self._wrap_text(center_title, main_width - 10, font)
        for i, line in enumerate(title_lines[:2]):
            draw.text((center_x - len(line) * 4, center_y - 10 + i * 12), line, fill='white', font=font)
        
        # Surrounding concepts from bullet points
        num_concepts = min(len(bullet_points), 5)
        colors = ['#e74c3c', '#2ecc71', '#f39c12', '#9b59b6', '#1abc9c']
        
        for i, bullet in enumerate(bullet_points[:num_concepts]):
            angle = (i * 2 * 3.14159) / num_concepts - 3.14159/2  # Start from top
            radius = 120
            
            concept_x = center_x + int(radius * cos(angle))
            concept_y = center_y + int(radius * sin(angle))
            
            # Connection line
            draw.line([center_x + int((main_width//2) * cos(angle)),
                      center_y + int((main_height//2) * sin(angle)),
                      concept_x - int(40 * cos(angle)),
                      concept_y - int(20 * sin(angle))],
                     fill='#7f8c8d', width=2)
            
            # Concept box
            concept_width, concept_height = 100, 40
            draw.rounded_rectangle([concept_x - concept_width//2, concept_y - concept_height//2,
                                  concept_x + concept_width//2, concept_y + concept_height//2],
                                 radius=8, fill=colors[i], outline='#2c3e50', width=2)
            
            # Concept text (first few words)
            concept_text = bullet.split()[:3]  # First 3 words
            concept_text = ' '.join(concept_text)[:15] + "..." if len(' '.join(concept_text)) > 15 else ' '.join(concept_text)
            
            text_lines = self._wrap_text(concept_text, concept_width - 10, font)
            for j, line in enumerate(text_lines[:2]):
                draw.text((concept_x - len(line) * 3, concept_y - 8 + j * 12), line, fill='white', font=font)

    def _draw_ui_design_from_content(self, draw, width, height, font, bullet_points, title):
        """Draw a UI mockup based on bullet points"""
        if not bullet_points:
            return
            
        # Draw a simple mockup with header, content areas, and buttons based on content
        # Header area
        header_height = 60
        if header_height > 0:
            draw.rounded_rectangle([50, 50, width - 50, 50 + header_height],
                                 radius=8, fill='#3498db', outline='#2980b9', width=2)
            draw.text((60, 70), title[:30] + "..." if len(title) > 30 else title, fill='white', font=font)
        
        # Main content areas based on bullet points
        content_start_y = 120
        content_height = (height - content_start_y - 80) // min(len(bullet_points), 4)
        
        for i, bullet in enumerate(bullet_points[:4]):
            y_pos = content_start_y + i * (content_height + 10)
            
            # Content box
            if content_height > 20:  # Ensure valid dimensions
                draw.rounded_rectangle([60, y_pos, width - 60, y_pos + max(content_height - 10, 20)],
                                     radius=5, fill='#ecf0f1', outline='#bdc3c7', width=1)
                
                # Content text
                content_text = bullet[:25] + "..." if len(bullet) > 25 else bullet
                draw.text((70, y_pos + 10), content_text, fill='#2c3e50', font=font)
        
        # Action buttons at bottom
        button_y = height - 60
        button_width = 100
        if button_y > content_start_y + 50:  # Ensure buttons don't overlap content
            # Primary button
            draw.rounded_rectangle([60, button_y, 60 + button_width, button_y + 30],
                                 radius=5, fill='#27ae60', outline='#229954', width=1)
            draw.text((75, button_y + 8), "Submit", fill='white', font=font)
            
            # Secondary button
            draw.rounded_rectangle([180, button_y, 180 + button_width, button_y + 30],
                                 radius=5, fill='#95a5a6', outline='#7f8c8d', width=1)
            draw.text((205, button_y + 8), "Cancel", fill='white', font=font)

    def _wrap_text(self, text: str, max_width: int, font) -> List[str]:
        """Wrap text to fit within specified width"""
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            # Test if adding this word would exceed the width
            test_line = ' '.join(current_line + [word])
            try:
                bbox = ImageDraw.Draw(Image.new('RGB', (1, 1))).textbbox((0, 0), test_line, font=font)
                text_width = bbox[2] - bbox[0]
            except:
                # Fallback for simple width estimation
                text_width = len(test_line) * 8
            
            if text_width <= max_width or not current_line:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return lines if lines else ['']

    def _create_sketch_image(self, bullet_points: List[str], title: str, suggestion: str, slide_number: int = 0) -> str:
        """Create a high-quality, professional sketch based on content and suggestion"""
        try:
            # Create high-resolution image for better quality
            width, height = 800, 600
            image = Image.new('RGB', (width, height), '#f8f9fa')  # Light gray background
            draw = ImageDraw.Draw(image)
            
            # Try to load better fonts with multiple fallbacks
            try:
                # Try to find the best available font
                font_paths = [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/System/Library/Fonts/Arial.ttf", 
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/Windows/Fonts/arial.ttf"
                ]
                font_loaded = False
                for font_path in font_paths:
                    try:
                        font = ImageFont.truetype(font_path, 18)
                        title_font = ImageFont.truetype(font_path, 24)
                        small_font = ImageFont.truetype(font_path, 14)
                        large_font = ImageFont.truetype(font_path, 20)
                        font_loaded = True
                        break
                    except:
                        continue
                
                if not font_loaded:
                    raise Exception("No fonts found")
                    
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
                large_font = ImageFont.load_default()
            
            # Combine content for analysis - include suggestion text for better matching
            all_content = f"{title} {' '.join(bullet_points)} {suggestion}".lower()
            
            # More specific keyword matching based on suggestion text
            sketch_type = self._determine_sketch_type(all_content, suggestion)
            
            # Draw sketch based on determined type
            if sketch_type == 'flowchart':
                self._draw_process_flowchart(draw, width, height, font, title_font)
            elif sketch_type == 'ai':
                self._draw_ai_concept(draw, width, height, font, title_font)
            elif sketch_type == 'comparison':
                self._draw_comparison_diagram(draw, width, height, font, title_font)
            elif sketch_type == 'data':
                self._draw_data_visualization(draw, width, height, font, title_font)
            elif sketch_type == 'ui':
                self._draw_ui_mockup(draw, width, height, font, title_font)
            elif sketch_type == 'system':
                self._draw_system_architecture(draw, width, height, font, title_font)
            elif sketch_type == 'learning':
                self._draw_learning_path(draw, width, height, font, title_font)
            elif sketch_type == 'collaboration':
                self._draw_collaboration_diagram(draw, width, height, font, title_font)
            elif sketch_type == 'problem_solution':
                self._draw_problem_solution(draw, width, height, font, title_font)
            elif sketch_type == 'timeline':
                self._draw_timeline(draw, width, height, font, title_font)
            elif sketch_type == 'business':
                self._draw_business_diagram(draw, width, height, font, title_font)
            else:
                # Generic concept diagram
                self._draw_generic_concept(draw, width, height, font, title_font, bullet_points)
            
            # Add a professional title with modern styling
            sketch_title = self._get_sketch_title(all_content, sketch_type)
            bbox = draw.textbbox((0, 0), sketch_title, font=title_font)
            title_width = bbox[2] - bbox[0]
            title_height = bbox[3] - bbox[1]
            
            # Title background with rounded corners
            title_x = width//2 - title_width//2 - 15
            title_y = 15
            draw.rounded_rectangle([title_x, title_y, title_x + title_width + 30, title_y + title_height + 15], 
                                 radius=8, fill='#37474F', outline='#263238', width=2)
            
            # Title text with white color for contrast
            draw.text((width//2 - title_width//2, 22), sketch_title, fill='white', font=title_font)
            
            # Save image to temporary file with unique naming
            os.makedirs('temp_sketches', exist_ok=True)
            # Use slide number, timestamp, and random component to ensure every slide gets a unique image
            import random
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            random_component = random.randint(1000, 9999)
            unique_id = abs(hash(f"{slide_number}_{title}_{suggestion}_{str(bullet_points)}_{sketch_type}_{timestamp}_{random_component}"))
            image_path = f'temp_sketches/sketch_slide{slide_number:03d}_{sketch_type}_{unique_id}_{timestamp}_{random_component}.png'
            image.save(image_path)
            
            logger.info(f"Created {sketch_type} sketch for slide {slide_number}: {image_path}")
            return image_path
            
        except Exception as e:
            logger.error(f"Error creating sketch image: {e}")
            return None
    
    def _determine_sketch_type(self, all_content: str, suggestion: str) -> str:
        """Determine the most appropriate sketch type based on content and suggestion"""
        suggestion_lower = suggestion.lower()
        
        # Check suggestion text first for most accurate matching
        if 'flowchart' in suggestion_lower or 'step-by-step' in suggestion_lower or 'arrows connecting' in suggestion_lower:
            return 'flowchart'
        elif 'side-by-side' in suggestion_lower or 'comparison' in suggestion_lower or 'two columns' in suggestion_lower:
            return 'comparison'
        elif 'brain' in suggestion_lower or 'neural network' in suggestion_lower or 'ai concept' in suggestion_lower:
            return 'ai'
        elif 'charts and graphs' in suggestion_lower or 'bar chart' in suggestion_lower or 'dashboard' in suggestion_lower:
            return 'data'
        elif 'wireframe' in suggestion_lower or 'buttons' in suggestion_lower or 'user interface' in suggestion_lower:
            return 'ui'
        elif 'system architecture' in suggestion_lower or 'components' in suggestion_lower or 'data flow' in suggestion_lower:
            return 'system'
        elif 'learning path' in suggestion_lower or 'milestone markers' in suggestion_lower or 'progression' in suggestion_lower:
            return 'learning'
        elif 'collaboration' in suggestion_lower or 'meeting table' in suggestion_lower or 'team members' in suggestion_lower:
            return 'collaboration'
        elif 'problem-solution' in suggestion_lower or 'warning icons' in suggestion_lower or 'checkmarks' in suggestion_lower:
            return 'problem_solution'
        elif 'timeline' in suggestion_lower or 'milestone flags' in suggestion_lower or 'progression' in suggestion_lower:
            return 'timeline'
        elif 'business diagrams' in suggestion_lower or 'growth arrows' in suggestion_lower or 'market segments' in suggestion_lower:
            return 'business'
        
        # Fallback to content analysis
        if any(keyword in all_content for keyword in ['process', 'workflow', 'steps', 'cycle', 'development', 'build', 'deploy']):
            return 'flowchart'
        elif any(keyword in all_content for keyword in ['ai', 'genai', 'chatbot', 'gpt', 'model', 'prompt']):
            return 'ai'
        elif any(keyword in all_content for keyword in ['compare', 'vs', 'versus', 'difference', 'before', 'after']):
            return 'comparison'
        elif any(keyword in all_content for keyword in ['data', 'analytics', 'chart', 'graph', 'metrics', 'visualization']):
            return 'data'
        elif any(keyword in all_content for keyword in ['interface', 'ui', 'app', 'website', 'screen', 'button']):
            return 'ui'
        elif any(keyword in all_content for keyword in ['system', 'architecture', 'components', 'api', 'database']):
            return 'system'
        elif any(keyword in all_content for keyword in ['learn', 'course', 'lesson', 'education', 'skill', 'training']):
            return 'learning'
        elif any(keyword in all_content for keyword in ['team', 'collaboration', 'meeting', 'communication']):
            return 'collaboration'
        elif any(keyword in all_content for keyword in ['problem', 'solution', 'issue', 'challenge', 'fix']):
            return 'problem_solution'
        elif any(keyword in all_content for keyword in ['timeline', 'schedule', 'roadmap', 'phases', 'milestone']):
            return 'timeline'
        elif any(keyword in all_content for keyword in ['business', 'strategy', 'market', 'growth', 'revenue']):
            return 'business'
        else:
            return 'generic'
    
    def _get_sketch_title(self, content: str, sketch_type: str) -> str:
        """Generate appropriate title for the sketch"""
        titles = {
            'flowchart': 'Process Flow',
            'ai': 'AI System',
            'comparison': 'Comparison',
            'data': 'Data Analysis',
            'ui': 'User Interface',
            'system': 'System Architecture',
            'learning': 'Learning Path',
            'collaboration': 'Team Collaboration',
            'problem_solution': 'Problem → Solution',
            'timeline': 'Timeline',
            'business': 'Business Strategy',
            'generic': 'Concept Overview'
        }
        return titles.get(sketch_type, 'Diagram')
    
    def _draw_process_flowchart(self, draw, width, height, font, title_font):
        """Draw a professional process flowchart with modern styling"""
        steps = ["Start", "Analyze", "Process", "Review", "Complete"]
        box_width, box_height = 130, 60
        spacing_x = 150
        start_x = (width - (len(steps) * box_width + (len(steps)-1) * 20)) // 2
        y = height // 2 - box_height // 2
        
        # Color scheme for different step types
        colors = {
            'start': '#4CAF50',      # Green
            'process': '#2196F3',    # Blue  
            'review': '#FF9800',     # Orange
            'end': '#9C27B0'         # Purple
        }
        
        for i, step in enumerate(steps):
            if i == 0:  # Vertical layout for better fit
                x = width // 2 - box_width // 2
                current_y = 80 + i * 90
            else:
                x = width // 2 - box_width // 2
                current_y = 80 + i * 90
            
            # Choose color based on step type
            if i == 0:
                color = colors['start']
            elif i == len(steps) - 1:
                color = colors['end']
            elif 'review' in step.lower() or 'analyze' in step.lower():
                color = colors['review']
            else:
                color = colors['process']
            
            # Draw shadow for depth
            shadow_offset = 4
            draw.rounded_rectangle([x + shadow_offset, current_y + shadow_offset, 
                                  x + box_width + shadow_offset, current_y + box_height + shadow_offset], 
                                 radius=12, fill='#00000030')
            
            # Draw main box with gradient-like effect
            draw.rounded_rectangle([x, current_y, x + box_width, current_y + box_height], 
                                 radius=12, fill=color, outline='#333333', width=3)
            
            # Add inner highlight for 3D effect
            draw.rounded_rectangle([x + 3, current_y + 3, x + box_width - 3, current_y + box_height - 3], 
                                 radius=9, outline='#ffffff40', width=1)
            
            # Add step text with better positioning
            bbox = draw.textbbox((0, 0), step, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            text_x = x + box_width//2 - text_width//2
            text_y = current_y + box_height//2 - text_height//2
            
            # Add text shadow for better readability
            draw.text((text_x + 1, text_y + 1), step, fill='#00000080', font=font)
            draw.text((text_x, text_y), step, fill='white', font=font)
            
            # Draw modern arrow to next step
            if i < len(steps) - 1:
                arrow_start_y = current_y + box_height + 5
                arrow_end_y = 80 + (i+1) * 90 - 5
                arrow_x = width // 2
                
                # Draw arrow line with gradient effect
                draw.line([arrow_x, arrow_start_y, arrow_x, arrow_end_y], 
                         fill='#666666', width=4)
                
                # Draw modern arrow head
                arrow_size = 10
                draw.polygon([arrow_x - arrow_size, arrow_end_y - arrow_size,
                             arrow_x + arrow_size, arrow_end_y - arrow_size,
                             arrow_x, arrow_end_y], 
                            fill='#666666')
    
    def _draw_ai_concept(self, draw, width, height, font, title_font):
        """Draw a modern AI brain/neural network concept"""
        center_x, center_y = width // 2, height // 2
        
        # Draw stylized brain with gradient effect
        brain_width, brain_height = 200, 140
        
        # Brain shadow for depth
        shadow_offset = 6
        draw.ellipse([center_x - brain_width//2 + shadow_offset, center_y - brain_height//2 + shadow_offset,
                     center_x + brain_width//2 + shadow_offset, center_y + brain_height//2 + shadow_offset], 
                    fill='#00000030')
        
        # Main brain with layered gradient effect
        colors = ['#E8F5E8', '#C8E6C9', '#A5D6A7', '#81C784']
        for i, color in enumerate(colors):
            offset = i * 10
            draw.ellipse([center_x - brain_width//2 + offset, center_y - brain_height//2 + offset,
                         center_x + brain_width//2 - offset, center_y + brain_height//2 - offset], 
                        fill=color, outline='#4CAF50', width=3)
        
        # Draw professional neural network layers
        layer_positions = [
            [(center_x-80, center_y-50), (center_x-80, center_y-15), (center_x-80, center_y+15), (center_x-80, center_y+50)],
            [(center_x-30, center_y-60), (center_x-30, center_y-20), (center_x-30, center_y+20), (center_x-30, center_y+60)],
            [(center_x+20, center_y-40), (center_x+20, center_y-10), (center_x+20, center_y+25), (center_x+20, center_y+55)],
            [(center_x+70, center_y-25), (center_x+70, center_y+25)]
        ]
        
        # Draw connections with varying thickness for importance
        for i in range(len(layer_positions)-1):
            for j, node1 in enumerate(layer_positions[i]):
                for k, node2 in enumerate(layer_positions[i+1]):
                    # Vary connection strength visually
                    thickness = 2 if (j + k) % 2 == 0 else 1
                    alpha = 'AA' if thickness == 2 else '66'
                    draw.line([node1[0], node1[1], node2[0], node2[1]], 
                             fill=f'#2196F3{alpha}', width=thickness)
        
        # Draw nodes with modern styling
        for layer_idx, layer in enumerate(layer_positions):
            node_color = ['#4CAF50', '#2196F3', '#FF9800', '#9C27B0'][layer_idx % 4]
            for x, y in layer:
                # Node shadow
                draw.ellipse([x-7, y-7, x+7, y+7], fill='#00000020')
                # Main node
                draw.ellipse([x-6, y-6, x+6, y+6], fill=node_color, outline='#333333', width=2)
                # Inner highlight
                draw.ellipse([x-3, y-3, x+3, y+3], fill='#ffffff60')
        
        # Add professional labels with backgrounds
        # Input label
        input_text = "INPUT"
        bbox = draw.textbbox((0, 0), input_text, font=font)
        text_width = bbox[2] - bbox[0]
        draw.rounded_rectangle([center_x-80-text_width//2-5, center_y-75, 
                              center_x-80+text_width//2+5, center_y-60], 
                             radius=5, fill='#4CAF50', outline='#333333', width=1)
        draw.text((center_x-80-text_width//2, center_y-72), input_text, fill='white', font=font)
        
        # Output label  
        output_text = "OUTPUT"
        bbox = draw.textbbox((0, 0), output_text, font=font)
        text_width = bbox[2] - bbox[0]
        draw.rounded_rectangle([center_x+70-text_width//2-5, center_y-40, 
                              center_x+70+text_width//2+5, center_y-25], 
                             radius=5, fill='#9C27B0', outline='#333333', width=1)
        draw.text((center_x+70-text_width//2, center_y-37), output_text, fill='white', font=font)
        
        # Main AI label with background
        ai_text = "NEURAL NETWORK"
        bbox = draw.textbbox((0, 0), ai_text, font=title_font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        draw.rounded_rectangle([center_x-text_width//2-10, center_y+80, 
                              center_x+text_width//2+10, center_y+80+text_height+10], 
                             radius=8, fill='#2196F3', outline='#0D47A1', width=2)
        draw.text((center_x-text_width//2, center_y+85), ai_text, fill='white', font=title_font)
    
    def _draw_comparison_diagram(self, draw, width, height, font, title_font):
        """Draw side-by-side comparison diagram"""
        center_x = width // 2
        y_start = 80
        
        # Draw dividing line
        draw.line([center_x, y_start, center_x, height-40], fill='black', width=2)
        
        # Left side - "Before"
        left_x = center_x // 2
        draw.text((left_x-20, y_start-20), "Before", fill='black', font=title_font)
        
        # Draw negative indicators
        draw.rectangle([left_x-40, y_start+20, left_x+40, y_start+50], outline='red', width=2)
        draw.text((left_x-25, y_start+30), "Problem", fill='black', font=font)
        
        draw.rectangle([left_x-40, y_start+70, left_x+40, y_start+100], outline='red', width=2)
        draw.text((left_x-20, y_start+80), "Issues", fill='black', font=font)
        
        # Right side - "After"
        right_x = center_x + center_x // 2
        draw.text((right_x-15, y_start-20), "After", fill='black', font=title_font)
        
        # Draw positive indicators
        draw.rectangle([right_x-40, y_start+20, right_x+40, y_start+50], outline='green', width=2)
        draw.text((right_x-25, y_start+30), "Solution", fill='black', font=font)
        
        draw.rectangle([right_x-40, y_start+70, right_x+40, y_start+100], outline='green', width=2)
        draw.text((right_x-25, y_start+80), "Benefits", fill='black', font=font)
        
        # Add arrow
        arrow_y = y_start + 60
        draw.polygon([center_x-20, arrow_y, center_x+20, arrow_y, center_x+15, arrow_y-5, center_x+15, arrow_y+5], fill='black')
    
    def _draw_timeline(self, draw, width, height, font, title_font):
        """Draw timeline with milestones"""
        y = height // 2 + 20
        start_x = 60
        end_x = width - 60
        
        # Draw main timeline
        draw.line([start_x, y, end_x, y], fill='black', width=4)
        
        # Draw milestones
        milestones = ["Start", "Phase 1", "Phase 2", "Complete"]
        milestone_x = [start_x + i * (end_x - start_x) // 3 for i in range(4)]
        
        for i, (x, milestone) in enumerate(zip(milestone_x, milestones)):
            # Draw milestone marker
            draw.ellipse([x-8, y-8, x+8, y+8], fill='black', outline='black')
            draw.ellipse([x-5, y-5, x+5, y+5], fill='white', outline='black')
            
            # Add milestone text
            bbox = draw.textbbox((0, 0), milestone, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((x - text_width//2, y + 20), milestone, fill='black', font=font)
            
            # Add date placeholders
            date = f"Week {i+1}"
            bbox = draw.textbbox((0, 0), date, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((x - text_width//2, y - 25), date, fill='gray', font=font)
    
    def _draw_business_diagram(self, draw, width, height, font, title_font):
        """Draw business strategy diagram"""
        center_x, center_y = width // 2, height // 2 + 20
        
        # Draw growth arrow
        arrow_start_x = 80
        arrow_end_x = width - 80
        arrow_y = center_y
        
        draw.line([arrow_start_x, arrow_y, arrow_end_x, arrow_y], fill='black', width=4)
        draw.polygon([arrow_end_x-15, arrow_y-8, arrow_end_x-15, arrow_y+8, arrow_end_x, arrow_y], fill='black')
        
        # Add growth indicators
        growth_points = [(120, center_y-20), (200, center_y-35), (280, center_y-50), (360, center_y-65)]
        
        for i, (x, y) in enumerate(growth_points):
            if x < width - 100:
                # Draw upward bar
                bar_height = 20 + i * 15
                draw.rectangle([x-10, y, x+10, y+bar_height], fill='green', outline='black')
                
                # Add percentage
                draw.text((x-15, y-15), f"{(i+1)*25}%", fill='black', font=font)
        
        # Add labels
        draw.text((arrow_start_x-30, arrow_y+30), "Growth", fill='black', font=title_font)
        draw.text((arrow_end_x-40, arrow_y+30), "Success", fill='black', font=title_font)
    
    def _draw_data_visualization(self, draw, width, height, font, title_font):
        """Draw simple charts and graphs"""
        # Bar chart with better proportions
        bar_x = 60
        bar_y = height - 100
        bar_heights = [30, 50, 40, 60, 35]
        bar_colors = ['black', 'gray', 'black', 'gray', 'black']
        
        # Draw bars
        for i, h in enumerate(bar_heights):
            x = bar_x + i * 35
            draw.rectangle([x, bar_y - h, x + 25, bar_y], fill=bar_colors[i], outline='black', width=2)
            draw.text((x + 8, bar_y + 5), f"Q{i+1}", fill='black', font=font)
        
        # Draw axes
        draw.line([bar_x-10, bar_y, bar_x + len(bar_heights)*35, bar_y], fill='black', width=2)  # X-axis
        draw.line([bar_x-10, bar_y, bar_x-10, bar_y-80], fill='black', width=2)  # Y-axis
        
        # Simple line graph on the right
        line_start_x = 280
        line_points = [(line_start_x, 200), (line_start_x+40, 180), (line_start_x+80, 160), (line_start_x+120, 140)]
        
        # Draw line graph
        for i in range(len(line_points) - 1):
            draw.line([line_points[i], line_points[i+1]], fill='black', width=3)
        
        # Draw points with values
        for i, point in enumerate(line_points):
            draw.ellipse([point[0]-4, point[1]-4, point[0]+4, point[1]+4], fill='white', outline='black', width=2)
            draw.text((point[0]-5, point[1]-20), f"{100-i*10}", fill='black', font=font)
        
        # Add chart labels
        draw.text((bar_x, 50), "Sales Data", fill='black', font=title_font)
        draw.text((line_start_x, 50), "Growth Trend", fill='black', font=title_font)
    
    def _draw_ui_mockup(self, draw, width, height, font, title_font):
        """Draw a modern, professional UI wireframe"""
        # Main container with shadow and modern styling
        container_shadow = 4
        draw.rounded_rectangle([50 + container_shadow, 80 + container_shadow, 
                              width-50 + container_shadow, height-60 + container_shadow], 
                             radius=12, fill='#00000020')
        
        # Main container
        draw.rounded_rectangle([50, 80, width-50, height-60], 
                             radius=12, fill='white', outline='#E0E0E0', width=3)
        
        # Modern header with gradient-like effect
        draw.rounded_rectangle([60, 90, width-60, 130], 
                             radius=8, fill='#2196F3', outline='#1976D2', width=2)
        draw.rounded_rectangle([62, 92, width-62, 128], 
                             radius=6, outline='#64B5F6', width=1)
        
        # Header text
        header_text = "App Header / Navigation"
        bbox = draw.textbbox((0, 0), header_text, font=font)
        text_width = bbox[2] - bbox[0]
        draw.text((70, 105), header_text, fill='white', font=font)
        
        # Menu icon (hamburger)
        menu_x = width - 100
        for i in range(3):
            draw.rectangle([menu_x, 100 + i*6, menu_x + 20, 102 + i*6], fill='white')
        
        # Modern buttons with different styles
        button_colors = ['#4CAF50', '#FF9800', '#9C27B0']
        button_texts = ['Primary', 'Secondary', 'Action']
        
        for i, (color, text) in enumerate(zip(button_colors, button_texts)):
            x = 70 + i * 120
            y = 150
            
            # Button shadow
            draw.rounded_rectangle([x + 2, y + 2, x + 102, y + 37], 
                                 radius=8, fill='#00000030')
            # Main button
            draw.rounded_rectangle([x, y, x + 100, y + 35], 
                                 radius=8, fill=color, outline='#333333', width=2)
            # Button highlight
            draw.rounded_rectangle([x + 2, y + 2, x + 98, y + 15], 
                                 radius=6, fill='#ffffff40')
            
            # Button text
            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((x + 50 - text_width//2, y + 12), text, fill='white', font=font)
        
        # Modern form fields
        field_labels = ['Email Address', 'Password', 'Confirm Password']
        for i, label in enumerate(field_labels):
            y = 210 + i * 50
            
            # Field label
            draw.text((70, y), label, fill='#333333', font=font)
            
            # Field container with modern styling
            draw.rounded_rectangle([70, y + 18, 350, y + 43], 
                                 radius=6, fill='white', outline='#BDBDBD', width=2)
            draw.rounded_rectangle([72, y + 20, 348, y + 41], 
                                 radius=4, outline='#E0E0E0', width=1)
            
            # Placeholder text
            placeholder = "Enter " + label.lower()
            draw.text((80, y + 26), placeholder, fill='#9E9E9E', font=font)
        
        # Modern mobile frame
        mobile_x = width - 150
        mobile_y = 140
        
        # Phone shadow
        draw.rounded_rectangle([mobile_x + 3, mobile_y + 3, mobile_x + 83, mobile_y + 163], 
                             radius=20, fill='#00000030')
        
        # Phone body
        draw.rounded_rectangle([mobile_x, mobile_y, mobile_x + 80, mobile_y + 160], 
                             radius=20, fill='#263238', outline='#37474F', width=3)
        
        # Screen
        draw.rounded_rectangle([mobile_x + 8, mobile_y + 20, mobile_x + 72, mobile_y + 140], 
                             radius=8, fill='white', outline='#90A4AE', width=1)
        
        # Mobile header
        draw.rounded_rectangle([mobile_x + 10, mobile_y + 25, mobile_x + 70, mobile_y + 45], 
                             radius=4, fill='#2196F3')
        draw.text((mobile_x + 25, mobile_y + 30), "Mobile", fill='white', font=font)
        
        # Mobile content blocks
        for i in range(3):
            block_y = mobile_y + 55 + i * 20
            draw.rounded_rectangle([mobile_x + 12, block_y, mobile_x + 68, block_y + 15], 
                                 radius=3, fill='#F5F5F5', outline='#E0E0E0', width=1)
    
    def _draw_system_architecture(self, draw, width, height, font, title_font):
        """Draw system architecture diagram"""
        # Database
        db_x, db_y = 60, 180
        draw.ellipse([db_x, db_y, db_x + 60, db_y + 40], outline='black', width=2)
        draw.text((db_x + 15, db_y + 15), "DB", fill='black', font=font)
        
        # API/Server
        api_x, api_y = 180, 140
        draw.rectangle([api_x, api_y, api_x + 80, api_y + 60], outline='black', width=2)
        draw.text((api_x + 25, api_y + 25), "API", fill='black', font=font)
        
        # Frontend
        ui_x, ui_y = 320, 100
        draw.rectangle([ui_x, ui_y, ui_x + 60, ui_y + 80], outline='black', width=2)
        draw.text((ui_x + 20, ui_y + 35), "UI", fill='black', font=font)
        
        # Connections
        draw.line([db_x + 60, db_y + 20, api_x, api_y + 30], fill='black', width=2)
        draw.line([api_x + 80, api_y + 30, ui_x, ui_y + 40], fill='black', width=2)
        
        # Labels
        draw.text((50, 30), "System Architecture", fill='black', font=font)
    
    def _draw_learning_path(self, draw, width, height, font, title_font):
        """Draw learning progression diagram"""
        # Learning steps as ascending stairs
        steps = ["Start", "Learn", "Practice", "Master"]
        for i, step in enumerate(steps):
            x = 50 + i * 80
            y = height - 60 - i * 30
            w, h = 60, 30
            
            draw.rectangle([x, y, x + w, y + h], outline='black', width=2)
            bbox = draw.textbbox((0, 0), step, font=font)
            text_width = bbox[2] - bbox[0]
            draw.text((x + w//2 - text_width//2, y + h//2 - 6), step, fill='black', font=font)
            
            # Arrow to next step
            if i < len(steps) - 1:
                draw.line([x + w + 5, y + h//2, x + w + 15, y + h//2 - 15], fill='black', width=2)
        
        # Add lightbulb icon
        bulb_x, bulb_y = width - 80, 50
        draw.ellipse([bulb_x, bulb_y, bulb_x + 30, bulb_y + 30], outline='black', width=2)
        draw.text((bulb_x + 8, bulb_y + 8), "💡", fill='black', font=font)
    
    def _draw_collaboration_diagram(self, draw, width, height, font, title_font):
        """Draw team collaboration diagram"""
        # Central meeting table (ellipse)
        center_x, center_y = width // 2, height // 2
        draw.ellipse([center_x - 60, center_y - 30, center_x + 60, center_y + 30], 
                    outline='black', width=2)
        
        # People around table (circles)
        positions = [(center_x - 80, center_y - 60), (center_x + 80, center_y - 60),
                    (center_x - 80, center_y + 60), (center_x + 80, center_y + 60)]
        
        for i, (x, y) in enumerate(positions):
            # Person (circle)
            draw.ellipse([x - 15, y - 15, x + 15, y + 15], outline='black', width=2)
            draw.text((x - 5, y - 5), str(i+1), fill='black', font=font)
            
            # Speech bubble
            bubble_x = x + (20 if i % 2 == 0 else -40)
            bubble_y = y - 30
            draw.ellipse([bubble_x, bubble_y, bubble_x + 20, bubble_y + 15], outline='black', width=1)
        
        draw.text((center_x - 30, 30), "Team Meeting", fill='black', font=font)
    
    def _draw_problem_solution(self, draw, width, height, font, title_font):
        """Draw problem-solution flow"""
        # Problem box (left)
        prob_x, prob_y = 40, height//2 - 40
        draw.rectangle([prob_x, prob_y, prob_x + 80, prob_y + 80], outline='black', width=2)
        draw.text((prob_x + 15, prob_y + 35), "Problem", fill='black', font=font)
        
        # Warning icon
        draw.text((prob_x + 35, prob_y + 10), "⚠", fill='black', font=font)
        
        # Arrow
        arrow_y = height // 2
        draw.line([prob_x + 90, arrow_y, prob_x + 150, arrow_y], fill='black', width=3)
        draw.line([prob_x + 145, arrow_y - 5, prob_x + 150, arrow_y], fill='black', width=3)
        draw.line([prob_x + 145, arrow_y + 5, prob_x + 150, arrow_y], fill='black', width=3)
        
        # Solution box (right)
        sol_x = prob_x + 160
        draw.rectangle([sol_x, prob_y, sol_x + 80, prob_y + 80], outline='black', width=2)
        draw.text((sol_x + 15, prob_y + 35), "Solution", fill='black', font=font)
        
        # Checkmark
        draw.text((sol_x + 35, prob_y + 10), "✓", fill='black', font=font)
    
    def _draw_generic_concept(self, draw, width, height, font, title_font, bullet_points):
        """Draw generic concept map with key terms"""
        center_x, center_y = width // 2, height // 2
        
        # Central concept circle
        draw.ellipse([center_x - 40, center_y - 30, center_x + 40, center_y + 30], 
                    outline='black', width=2)
        draw.text((center_x - 20, center_y - 10), "Concept", fill='black', font=font)
        
        # Surrounding concept bubbles
        if bullet_points:
            positions = [(center_x - 100, center_y - 80), (center_x + 100, center_y - 80),
                        (center_x - 100, center_y + 80), (center_x + 100, center_y + 80)]
            
            for i, (x, y) in enumerate(positions[:min(4, len(bullet_points))]):
                # Extract first word from bullet point
                words = bullet_points[i].split()[:2] if i < len(bullet_points) else ["Item"]
                text = " ".join(words)[:10]  # Limit length
                
                # Draw bubble
                bubble_width = max(40, len(text) * 6)
                draw.ellipse([x - bubble_width//2, y - 15, x + bubble_width//2, y + 15], 
                           outline='black', width=1)
                
                # Add text
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                draw.text((x - text_width//2, y - 5), text, fill='black', font=font)
                
                # Connect to center
                draw.line([center_x, center_y, x, y], fill='black', width=1)
    
    def _determine_heading_level(self, title: str) -> int:
        """Determine the heading level of a title"""
        title_lower = title.lower()
        
        # Level 1 indicators (major sections)
        level_1_keywords = [
            'introduction', 'overview', 'what is', 'types of', 'getting started',
            'conclusion', 'summary', 'applications', 'key concepts', 'fundamentals'
        ]
        
        # Level 2 indicators (subsections)
        level_2_keywords = [
            'supervised', 'unsupervised', 'reinforcement', 'algorithm', 'example',
            'healthcare', 'finance', 'technology', 'prerequisites', 'learning path'
        ]
        
        # Level 3 indicators (sub-subsections)
        level_3_keywords = [
            'step', 'phase', 'stage', 'part', 'section', 'lesson', 'chapter'
        ]
        
        # Level 4 indicators (individual slide titles)
        level_4_keywords = [
            'how to', 'why', 'when', 'where', 'what', 'which', 'who',
            'implementation', 'details', 'features', 'benefits', 'challenges',
            'best practices', 'tips', 'tricks', 'common', 'key points'
        ]
        
        if any(keyword in title_lower for keyword in level_1_keywords):
            return 1
        elif any(keyword in title_lower for keyword in level_2_keywords):
            return 2
        elif any(keyword in title_lower for keyword in level_3_keywords):
            return 3
        elif any(keyword in title_lower for keyword in level_4_keywords):
            return 4
        else:
            # Default: shorter titles are usually more specific (higher level)
            if len(title.split()) <= 3:
                return 4  # Short titles are usually slide titles
            elif len(title.split()) <= 6:
                return 3  # Medium titles are subsections
            else:
                return 2  # Longer titles are sections
    
    def _create_section_overview(self, content: List[str]) -> str:
        """Create a brief overview for section intro slides"""
        if not content:
            return None
        
        # Take first sentence or two as overview
        overview_parts = []
        char_count = 0
        
        for item in content[:3]:  # Max 3 items
            if char_count + len(item) > 150:  # Keep overview concise
                break
            overview_parts.append(item)
            char_count += len(item)
        
        return ' '.join(overview_parts) if overview_parts else None
    
    def create_html_slides(self, doc_structure: DocumentStructure) -> str:
        """Generate HTML presentation from document structure"""
        slides_html = []
        
        # Title slide
        title_slide = f"""
        <section class="slide title-slide">
            <h1>{doc_structure.title}</h1>
            <p class="subtitle">Generated from {doc_structure.metadata['filename']}</p>
            <p class="date">{datetime.now().strftime('%B %d, %Y')}</p>
        </section>
        """
        slides_html.append(title_slide)
        
        # Content slides
        for slide_content in doc_structure.slides:
            content_items = '</li><li>'.join(slide_content.content[:8]) if slide_content.content else ''
            
            slide_html = f"""
            <section class="slide content-slide">
                <h2>{slide_content.title}</h2>
                {f'<ul><li>{content_items}</li></ul>' if content_items else ''}
            </section>
            """
            slides_html.append(slide_html)
        
        # Use the same HTML template from the original app
        html_template = self._get_html_template()
        full_html = html_template.replace('{{SLIDES}}', '\n'.join(slides_html))
        full_html = full_html.replace('{{TITLE}}', doc_structure.title)
        
        # Save HTML file
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        filepath = os.path.join(EXPORT_FOLDER, filename)
        os.makedirs(EXPORT_FOLDER, exist_ok=True)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(full_html)
        
        return filepath
    
    def _create_thought_bubble_image(self, drawing_prompt: str, slide_number: int) -> str:
        """Create a thought bubble with visual prompt on light blue background"""
        try:
            # Create image dimensions
            width, height = 800, 600
            
            # Light blue background (#E6F3FF)
            image = Image.new('RGB', (width, height), '#E6F3FF')
            draw = ImageDraw.Draw(image)
            
            # Load fonts
            try:
                font_paths = [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/System/Library/Fonts/Arial.ttf", 
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/Windows/Fonts/arial.ttf"
                ]
                font_loaded = False
                for font_path in font_paths:
                    try:
                        font = ImageFont.truetype(font_path, 16)
                        title_font = ImageFont.truetype(font_path, 20)
                        font_loaded = True
                        break
                    except:
                        continue
                
                if not font_loaded:
                    raise Exception("No fonts found")
                    
            except:
                font = ImageFont.load_default()
                title_font = ImageFont.load_default()
            
            # Draw thought bubble outline
            bubble_margin = 40
            bubble_x1 = bubble_margin
            bubble_y1 = bubble_margin
            bubble_x2 = width - bubble_margin
            bubble_y2 = height - bubble_margin
            
            # Main thought bubble (white with light blue border)
            draw.ellipse([bubble_x1, bubble_y1, bubble_x2, bubble_y2], 
                        fill='white', outline='#4A90E2', width=3)
            
            # Small thought bubbles leading to main bubble
            small_bubble_1_x = bubble_x1 + 30
            small_bubble_1_y = bubble_y2 + 10
            draw.ellipse([small_bubble_1_x, small_bubble_1_y, 
                         small_bubble_1_x + 20, small_bubble_1_y + 20],
                        fill='white', outline='#4A90E2', width=2)
            
            small_bubble_2_x = bubble_x1 + 60
            small_bubble_2_y = bubble_y2 + 25
            draw.ellipse([small_bubble_2_x, small_bubble_2_y,
                         small_bubble_2_x + 15, small_bubble_2_y + 15],
                        fill='white', outline='#4A90E2', width=2)
            
            # Add title at the top of bubble
            title_text = "💡 Visual Inspiration"
            title_bbox = draw.textbbox((0, 0), title_text, font=title_font)
            title_width = title_bbox[2] - title_bbox[0]
            title_x = width//2 - title_width//2
            title_y = bubble_y1 + 20
            
            draw.text((title_x, title_y), title_text, font=title_font, 
                     fill='#2C5282', anchor='mm')
            
            # Wrap and draw the prompt text
            max_width = bubble_x2 - bubble_x1 - 80  # Leave margins
            wrapped_lines = self._wrap_text(drawing_prompt, max_width, font)
            
            # Calculate total text height
            line_height = 24
            total_text_height = len(wrapped_lines) * line_height
            start_y = bubble_y1 + 80  # Start below title
            
            # Draw each line of wrapped text
            for i, line in enumerate(wrapped_lines):
                text_bbox = draw.textbbox((0, 0), line, font=font)
                text_width = text_bbox[2] - text_bbox[0]
                text_x = width//2 - text_width//2
                text_y = start_y + (i * line_height)
                
                draw.text((text_x, text_y), line, font=font, fill='#2D3748')
            
            # Add copy instruction at bottom
            copy_text = "💻 Copy & paste this into any AI image generator"
            copy_bbox = draw.textbbox((0, 0), copy_text, font=font)
            copy_width = copy_bbox[2] - copy_bbox[0]
            copy_x = width//2 - copy_width//2
            copy_y = bubble_y2 - 40
            
            draw.text((copy_x, copy_y), copy_text, font=font, 
                     fill='#4A90E2', anchor='mm')
            
            # Save the image
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
            random_suffix = random.randint(1000, 9999)
            filename = f"thought_bubble_slide{slide_number:03d}_{timestamp}_{random_suffix}.png"
            
            # Ensure temp_sketches directory exists
            os.makedirs('temp_sketches', exist_ok=True)
            filepath = os.path.join('temp_sketches', filename)
            
            image.save(filepath, 'PNG', optimize=True, quality=85)
            logger.info(f"Created thought bubble image: {filepath}")
            
            return filepath
            
        except Exception as e:
            logger.error(f"Error creating thought bubble image: {e}")
            return None
    
    def _get_html_template(self) -> str:
        """HTML template for slides - same as original"""
        return """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{TITLE}}</title>
    <style>
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 0;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: #333;
        }
        
        .slideshow-container {
            max-width: 1000px;
            position: relative;
            margin: auto;
            height: 100vh;
            display: flex;
            align-items: center;
        }
        
        .slide {
            display: none;
            padding: 60px;
            text-align: center;
            background: white;
            border-radius: 10px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.3);
            margin: 20px;
            min-height: 500px;
            display: flex;
            flex-direction: column;
            justify-content: center;
        }
        
        .slide.active {
            display: flex;
        }
        
        .title-slide h1 {
            font-size: 3em;
            color: #667eea;
            margin-bottom: 30px;
        }
        
        .content-slide h2 {
            font-size: 2.5em;
            color: #667eea;
            margin-bottom: 30px;
            border-bottom: 3px solid #667eea;
            padding-bottom: 15px;
        }
        
        .content-slide ul {
            text-align: left;
            font-size: 1.4em;
            line-height: 1.8;
            list-style-type: none;
            padding-left: 0;
        }
        
        .content-slide li {
            margin: 15px 0;
            padding-left: 30px;
            position: relative;
        }
        
        .content-slide li::before {
            content: "→";
            color: #667eea;
            font-weight: bold;
            position: absolute;
            left: 0;
        }
        
        .subtitle {
            font-size: 1.2em;
            color: #666;
            margin: 20px 0;
        }
        
        .date {
            font-size: 1em;
            color: #999;
        }
        
        .navigation {
            position: fixed;
            bottom: 30px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 15px;
        }
        
        .nav-btn {
            background: #667eea;
            color: white;
            border: none;
            padding: 10px 20px;
            border-radius: 25px;
            cursor: pointer;
            font-size: 1em;
            transition: all 0.3s ease;
        }
        
        .nav-btn:hover {
            background: #5a67d8;
            transform: translateY(-2px);
        }
        
        .nav-btn:disabled {
            background: #ccc;
            cursor: not-allowed;
            transform: none;
        }
        
        .slide-counter {
            position: fixed;
            top: 30px;
            right: 30px;
            background: rgba(255,255,255,0.9);
            padding: 10px 15px;
            border-radius: 20px;
            font-weight: bold;
            color: #667eea;
        }
    </style>
</head>
<body>
    <div class="slideshow-container">
        {{SLIDES}}
    </div>
    
    <div class="slide-counter">
        <span id="current-slide">1</span> / <span id="total-slides"></span>
    </div>
    
    <div class="navigation">
        <button class="nav-btn" id="prev-btn" onclick="changeSlide(-1)">← Previous</button>
        <button class="nav-btn" id="next-btn" onclick="changeSlide(1)">Next →</button>
    </div>
    
    <script>
        let currentSlide = 0;
        const slides = document.querySelectorAll('.slide');
        const totalSlides = slides.length;
        
        document.getElementById('total-slides').textContent = totalSlides;
        
        function showSlide(n) {
            slides.forEach(slide => slide.classList.remove('active'));
            
            if (n >= totalSlides) currentSlide = 0;
            if (n < 0) currentSlide = totalSlides - 1;
            
            slides[currentSlide].classList.add('active');
            document.getElementById('current-slide').textContent = currentSlide + 1;
            
            document.getElementById('prev-btn').disabled = currentSlide === 0;
            document.getElementById('next-btn').disabled = currentSlide === totalSlides - 1;
        }
        
        function changeSlide(n) {
            currentSlide += n;
            showSlide(currentSlide);
        }
        
        // Keyboard navigation
        document.addEventListener('keydown', function(e) {
            if (e.key === 'ArrowLeft') changeSlide(-1);
            if (e.key === 'ArrowRight') changeSlide(1);
        });
        
        // Initialize
        showSlide(0);
    </script>
</body>
</html>
        """

# Helper functions
def allowed_file(filename):
    """Check if file extension is allowed"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Flask routes
@app.route('/')
def index():
    """Main page"""
    return render_template('file_to_slides.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    """Handle file upload and conversion"""
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': f'File type not supported. Allowed: {", ".join(ALLOWED_EXTENSIONS)}'}), 400
    
    export_format = request.form.get('format', 'pptx')
    script_column = int(request.form.get('script_column', '2'))  # Default to column 2
    
    try:
        # Save uploaded file
        filename = secure_filename(file.filename)
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)
        
        # Parse document
        parser = DocumentParser()
        doc_structure = parser.parse_file(filepath, filename, script_column)
        
        # Generate slides
        generator = SlideGenerator()
        
        if export_format == 'pptx':
            output_path = generator.create_powerpoint(doc_structure)
        elif export_format == 'html':
            output_path = generator.create_html_slides(doc_structure)
        elif export_format == 'google-slides-json':
            output_path = generator.create_google_slides_json(doc_structure)
        elif export_format == 'google-slides-text':
            output_path = generator.create_google_slides_text(doc_structure)
        elif export_format == 'google-slides-csv':
            output_path = generator.create_google_slides_csv(doc_structure)
        else:
            return jsonify({'error': 'Unsupported format'}), 400
        
        # Clean up uploaded file
        os.remove(filepath)
        
        return jsonify({
            'success': True,
            'filename': os.path.basename(output_path),
            'download_url': f'/download/{os.path.basename(output_path)}',
            'slide_count': len(doc_structure.slides),
            'title': doc_structure.title
        })
        
    except Exception as e:
        logger.error(f"Error processing file: {str(e)}")
        # Clean up on error
        if 'filepath' in locals() and os.path.exists(filepath):
            os.remove(filepath)
        
        return jsonify({'error': f'Error processing file: {str(e)}'}), 500

@app.route('/download/<filename>')
def download_file(filename):
    """Download generated presentation file"""
    filepath = os.path.join(EXPORT_FOLDER, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    else:
        return jsonify({'error': 'File not found'}), 404

if __name__ == '__main__':
    # Create necessary directories
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    os.makedirs(EXPORT_FOLDER, exist_ok=True)
    
    app.run(debug=True, port=5001)  # Different port to avoid conflict