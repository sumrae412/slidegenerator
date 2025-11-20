#!/usr/bin/env python3
"""
Demo: Data Visualization Feature

Demonstrates AI-powered chart detection and generation from slide content.

Features demonstrated:
1. Automatic chart type detection (pie, bar, line, scatter)
2. Data extraction from text
3. PowerPoint chart generation
4. Integration with slide generation pipeline
"""

import os
import sys
from typing import List

# Add parent directory to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_generator_pkg.data_intelligence import DataIntelligence
from pptx import Presentation


# Demo content examples

DEMO_EXAMPLES = [
    {
        'title': '📊 Revenue Distribution by Customer Segment',
        'text': """
        Our fiscal year 2024 revenue breakdown demonstrates strong enterprise growth:
        Enterprise customers contribute 45% of total revenue, showing our strength in large accounts.
        SMB customers represent 30% of revenue, providing a stable middle market.
        Startup customers account for 15% of revenue, giving us early-stage exposure.
        Non-profit customers make up the remaining 10%, supporting our mission-driven goals.
        """,
        'bullets': [
            "Enterprise customers: 45%",
            "SMB customers: 30%",
            "Startup customers: 15%",
            "Non-profit customers: 10%"
        ],
        'expected_chart': 'pie'
    },
    {
        'title': '⭐ Customer Satisfaction Scores by Department',
        'text': """
        Q4 2024 customer satisfaction survey results across departments:
        Support team achieved 85 out of 100, showing room for improvement.
        Product team scored highest at 92 out of 100, reflecting quality engineering.
        Documentation scored 78 out of 100, highlighting need for content updates.
        Onboarding team achieved 88 out of 100, demonstrating strong first impressions.
        """,
        'bullets': [
            "Support: 85 out of 100",
            "Product: 92 out of 100",
            "Documentation: 78 out of 100",
            "Onboarding: 88 out of 100"
        ],
        'expected_chart': 'bar/column'
    },
    {
        'title': '📈 Monthly Active Users Growth Trajectory',
        'text': """
        Five-month user growth demonstrates strong product-market fit:
        January started with 1,000 active users as our baseline.
        February showed 20% growth to 1,200 active users.
        March accelerated to 1,450 active users with viral features.
        April reached 1,800 active users with paid campaigns.
        May achieved 2,200 active users, our highest month yet.
        """,
        'bullets': [
            "January: 1,000 active users",
            "February: 1,200 active users",
            "March: 1,450 active users",
            "April: 1,800 active users",
            "May: 2,200 active users"
        ],
        'expected_chart': 'line'
    },
    {
        'title': '🔄 Team Size vs Delivery Velocity',
        'text': """
        Analysis of team size impact on feature delivery shows diminishing returns:
        3-person teams deliver 5 features per month, establishing our baseline.
        5-person teams deliver 12 features per month, showing strong efficiency.
        8-person teams deliver 18 features per month, approaching optimal size.
        12-person teams deliver only 22 features per month, indicating coordination overhead.
        Data suggests 8-person teams may be the sweet spot for productivity.
        """,
        'bullets': [
            "3 person team: 5 features/month",
            "5 person team: 12 features/month",
            "8 person team: 18 features/month",
            "12 person team: 22 features/month"
        ],
        'expected_chart': 'scatter/line'
    },
    {
        'title': '💡 Design Principles (No Chart)',
        'text': """
        Our product design philosophy focuses on user-centric principles:
        Intuitive navigation guides users naturally through workflows.
        Clean interface reduces cognitive load and visual clutter.
        Responsive design ensures consistency across all devices.
        Accessible features make our product usable for everyone.
        These principles guide every design decision we make.
        """,
        'bullets': [
            "Intuitive navigation that guides users naturally",
            "Clean interface that reduces cognitive load",
            "Responsive design that works on all devices",
            "Accessible features for all users"
        ],
        'expected_chart': 'none (qualitative)'
    },
    {
        'title': '💰 Cloud Cost Savings by Company Size',
        'text': """
        Our cloud optimization service delivers measurable savings across all company sizes:
        Small companies (under 50 employees) save 40% on cloud costs.
        Medium companies (50-200 employees) save 50% on cloud costs.
        Large companies (200-1000 employees) save 55% on cloud costs.
        Enterprise companies (over 1000 employees) save 60% on cloud costs.
        Larger companies benefit more due to economies of scale.
        """,
        'bullets': [
            "Small companies: 40% cost savings",
            "Medium companies: 50% cost savings",
            "Large companies: 55% cost savings",
            "Enterprise: 60% cost savings"
        ],
        'expected_chart': 'column'
    }
]


def print_header(text: str):
    """Print a formatted header"""
    print("\n" + "="*80)
    print(text)
    print("="*80 + "\n")


def print_section(text: str):
    """Print a formatted section header"""
    print(f"\n{text}")
    print("-"*80)


def demo_chart_detection():
    """Demonstrate chart detection across different content types"""
    print_header("DATA VISUALIZATION FEATURE DEMO")

    # Check for API key
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        print("❌ ERROR: ANTHROPIC_API_KEY environment variable not set")
        print("   Please set your API key to run this demo")
        print("   Example: export ANTHROPIC_API_KEY='your-key-here'")
        return

    # Initialize DataIntelligence
    print("🔧 Initializing DataIntelligence with Claude API...")
    from anthropic import Anthropic
    client = Anthropic(api_key=api_key)
    data_intelligence = DataIntelligence(client=client)
    print("✅ DataIntelligence initialized successfully\n")

    # Analyze each example
    total_cost = 0.0
    charts_detected = 0

    for i, example in enumerate(DEMO_EXAMPLES, 1):
        print_section(f"Example {i}: {example['title']}")

        print(f"\n📝 Content Preview:")
        print(f"   {example['text'].strip()[:150]}...")

        print(f"\n📌 Bullets:")
        for bullet in example['bullets'][:3]:
            print(f"   • {bullet}")
        if len(example['bullets']) > 3:
            print(f"   ... and {len(example['bullets']) - 3} more")

        # Analyze for visualization
        print(f"\n🔍 Analyzing for data visualization...")
        viz_config = data_intelligence.suggest_visualization(
            text=example['text'],
            bullets=example['bullets'],
            slide_title=example['title']
        )

        total_cost += viz_config.cost

        # Display results
        if viz_config.should_visualize:
            charts_detected += 1
            print(f"\n✅ CHART RECOMMENDED")
            print(f"   Chart Type: {viz_config.chart_type.upper()}")
            print(f"   Chart Title: {viz_config.chart_title}")
            print(f"   Confidence: {viz_config.confidence:.1%}")
            print(f"   X-axis: {viz_config.x_label}")
            print(f"   Y-axis: {viz_config.y_label}")

            print(f"\n📊 Extracted Data:")
            labels = viz_config.data['labels']
            values = viz_config.data['series'][0]['values']
            for label, value in zip(labels, values):
                print(f"   {label}: {value}")

            if len(viz_config.data['series']) > 1:
                print(f"\n   Multiple series detected:")
                for series in viz_config.data['series']:
                    print(f"   - {series['name']}: {series['values']}")

            print(f"\n💭 Reasoning: {viz_config.reasoning}")
            print(f"💵 Cost: ${viz_config.cost:.4f}")

        else:
            print(f"\n⚪ NO CHART RECOMMENDED")
            print(f"   Reasoning: {viz_config.reasoning}")
            print(f"💵 Cost: ${viz_config.cost:.4f}")

        print(f"\n{'─'*80}")

    # Summary
    print_section("SUMMARY")
    print(f"📊 Charts detected: {charts_detected}/{len(DEMO_EXAMPLES)}")
    print(f"💵 Total API cost: ${total_cost:.4f}")
    print(f"⚡ Average cost per slide: ${total_cost/len(DEMO_EXAMPLES):.4f}")


def demo_chart_generation():
    """Demonstrate actual PowerPoint chart generation"""
    print_header("POWERPOINT CHART GENERATION DEMO")

    # Check for API key
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        print("❌ ERROR: ANTHROPIC_API_KEY environment variable not set")
        return

    # Initialize
    from anthropic import Anthropic
    client = Anthropic(api_key=api_key)
    data_intelligence = DataIntelligence(client=client)

    # Create presentation
    prs = Presentation()
    charts_created = 0

    print("🎨 Creating PowerPoint presentation with charts...\n")

    # Process first 3 examples (that should have charts)
    for example in DEMO_EXAMPLES[:4]:
        print(f"Processing: {example['title']}")

        viz_config = data_intelligence.suggest_visualization(
            text=example['text'],
            bullets=example['bullets'],
            slide_title=example['title']
        )

        if viz_config.should_visualize:
            chart_slide = data_intelligence.create_chart_slide(viz_config, prs)
            if chart_slide:
                charts_created += 1
                print(f"✅ Created {viz_config.chart_type} chart")
            else:
                print(f"❌ Failed to create chart")
        else:
            print(f"⚪ No chart recommended")

    # Save presentation
    if charts_created > 0:
        output_file = '/home/user/slidegenerator/exports/demo_charts.pptx'
        os.makedirs(os.path.dirname(output_file), exist_ok=True)
        prs.save(output_file)

        print(f"\n✅ SUCCESS!")
        print(f"   Created {charts_created} charts")
        print(f"   Total slides: {len(prs.slides)}")
        print(f"   Saved to: {output_file}")
        print(f"\n   Open the file in PowerPoint to see the charts!")
    else:
        print(f"\n⚠️  No charts were created")


def demo_integration_example():
    """Show how to use data visualization in DocumentParser"""
    print_header("INTEGRATION EXAMPLE")

    print("""
Integration with DocumentParser:

from slide_generator_pkg.document_parser import DocumentParser

# Initialize with data visualization enabled
parser = DocumentParser(
    claude_api_key='your-key',
    enable_data_visualization=True  # Enable chart detection
)

# Parse document
doc_structure = parser.parse_file('document.docx', 'document.docx')

# Charts are automatically detected and added to slides
# Access chart configuration from slides:
for slide in doc_structure.slides:
    if slide.should_visualize:
        print(f"Chart for slide: {slide.title}")
        print(f"Chart type: {slide.chart_config['chart_type']}")
        print(f"Data: {slide.chart_config['data']}")

# Generate PowerPoint with charts
from slide_generator_pkg.powerpoint_generator import SlideGenerator

slide_generator = SlideGenerator(
    data_intelligence=parser.data_intelligence  # Pass data intelligence
)

pptx_file = slide_generator.create_powerpoint(doc_structure)
print(f"Presentation with charts saved to: {pptx_file}")
    """)


def main():
    """Run all demos"""
    import argparse

    parser = argparse.ArgumentParser(description='Data Visualization Feature Demo')
    parser.add_argument('--detect', action='store_true', help='Run chart detection demo')
    parser.add_argument('--generate', action='store_true', help='Run chart generation demo')
    parser.add_argument('--integration', action='store_true', help='Show integration example')
    parser.add_argument('--all', action='store_true', help='Run all demos')

    args = parser.parse_args()

    if args.all or (not args.detect and not args.generate and not args.integration):
        demo_chart_detection()
        demo_chart_generation()
        demo_integration_example()
    else:
        if args.detect:
            demo_chart_detection()
        if args.generate:
            demo_chart_generation()
        if args.integration:
            demo_integration_example()


if __name__ == '__main__':
    main()
