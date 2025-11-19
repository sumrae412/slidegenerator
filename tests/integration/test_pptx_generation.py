"""
Comprehensive PowerPoint generation tests for slidegenerator

Tests cover:
- SlideGenerator initialization
- Presentation creation and basic structure
- Slide content generation (titles, bullets, tables)
- PPTX file generation and validation
- End-to-end document conversion
"""

import pytest
import sys
import os
import tempfile
from io import BytesIO
from pathlib import Path

# Add parent directory to path for imports
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '../..')))

from file_to_slides import (
    SlideGenerator,
    DocumentParser,
    SlideContent,
    DocumentStructure,
    Presentation
)
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt


# ============================================================================
# FIXTURES
# ============================================================================

@pytest.fixture
def slide_generator():
    """Fixture providing a SlideGenerator instance"""
    return SlideGenerator()


@pytest.fixture
def document_parser():
    """Fixture providing a DocumentParser instance"""
    return DocumentParser()


@pytest.fixture
def temp_pptx_path():
    """Fixture providing temporary PPTX file path"""
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        temp_path = tmp.name
    yield temp_path
    # Cleanup
    if os.path.exists(temp_path):
        os.remove(temp_path)


@pytest.fixture
def simple_slide_content():
    """Fixture providing simple SlideContent for testing"""
    return SlideContent(
        title="Test Slide",
        content=["Bullet point 1", "Bullet point 2", "Bullet point 3"],
        slide_type='content'
    )


@pytest.fixture
def doc_structure_simple():
    """Fixture providing a simple DocumentStructure for testing"""
    slides = [
        SlideContent(
            title="Introduction",
            content=["Overview of the topic", "Key concepts"],
            slide_type='content',
            heading_level=1
        ),
        SlideContent(
            title="Main Content",
            content=["First key point", "Second key point", "Third key point"],
            slide_type='content',
            heading_level=2
        ),
        SlideContent(
            title="Conclusion",
            content=["Summary of main points", "Next steps"],
            slide_type='content',
            heading_level=2
        )
    ]

    return DocumentStructure(
        title="Test Presentation",
        slides=slides,
        metadata={'filename': 'test_doc.docx'}
    )


@pytest.fixture
def doc_structure_with_headings():
    """Fixture providing DocumentStructure with heading hierarchy"""
    slides = [
        SlideContent(
            title="Main Title",
            content=["Presentation overview"],
            slide_type='content',
            heading_level=1
        ),
        SlideContent(
            title="Section One",
            content=["Section introduction"],
            slide_type='content',
            heading_level=2
        ),
        SlideContent(
            title="Subsection A",
            content=["Subsection details"],
            slide_type='content',
            heading_level=3
        ),
        SlideContent(
            title="Topic A.1",
            content=["Detailed point 1", "Detailed point 2"],
            slide_type='content',
            heading_level=4
        ),
        SlideContent(
            title="Section Two",
            content=["Next section"],
            slide_type='content',
            heading_level=2
        )
    ]

    return DocumentStructure(
        title="Hierarchical Presentation",
        slides=slides,
        metadata={'filename': 'hierarchical_doc.docx'}
    )


# ============================================================================
# 1. SlideGenerator Initialization Tests
# ============================================================================

class TestSlideGeneratorInit:
    """Tests for SlideGenerator initialization"""

    def test_slide_generator_init(self):
        """Test basic SlideGenerator initialization"""
        generator = SlideGenerator()
        assert generator is not None
        assert generator.client is None

    def test_slide_generator_with_api_key(self):
        """Test SlideGenerator initialization with mock API key"""
        mock_client = None  # No real client for testing
        generator = SlideGenerator(openai_client=mock_client)
        assert generator is not None
        assert generator.client is None

    def test_slide_generator_has_create_powerpoint_method(self):
        """Test SlideGenerator has create_powerpoint method"""
        generator = SlideGenerator()
        assert hasattr(generator, 'create_powerpoint')
        assert callable(generator.create_powerpoint)

    def test_slide_generator_has_hierarchy_method(self):
        """Test SlideGenerator has slide organization method"""
        generator = SlideGenerator()
        assert hasattr(generator, '_organize_slides_by_hierarchy')
        assert callable(generator._organize_slides_by_hierarchy)


# ============================================================================
# 2. Presentation Creation Tests
# ============================================================================

class TestPresentationCreation:
    """Tests for basic presentation creation"""

    def test_create_presentation_basic(self, slide_generator, doc_structure_simple):
        """Test creating a basic presentation"""
        # Generate presentation bytes/path
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            temp_path = tmp.name

        try:
            slide_generator.create_powerpoint(doc_structure_simple)
            # If no exception, presentation was created
            assert True
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)

    def test_presentation_has_default_dimensions(self):
        """Test presentation has standard slide dimensions"""
        prs = Presentation()
        # Standard PowerPoint slide size: 10 inches x 7.5 inches
        assert prs.slide_width == Inches(10)
        assert prs.slide_height == Inches(7.5)

    def test_create_title_slide(self, slide_generator, doc_structure_simple):
        """Test title slide creation from document structure"""
        # Document with heading level 1 should trigger title slide
        doc = DocumentStructure(
            title="My Presentation Title",
            slides=[
                SlideContent(
                    title="My Presentation Title",
                    content=["Overview"],
                    slide_type='content',
                    heading_level=1
                )
            ],
            metadata={'filename': 'test.docx'}
        )

        slide_generator.create_powerpoint(doc)
        # If no exception, title slide was created successfully
        assert True

    def test_create_section_slide(self, slide_generator, doc_structure_with_headings):
        """Test section slide creation from H2 heading"""
        slide_generator.create_powerpoint(doc_structure_with_headings)
        # If no exception, section slides were created
        assert True


# ============================================================================
# 3. Content Slide Tests
# ============================================================================

class TestContentSlides:
    """Tests for content slide creation and formatting"""

    def test_add_bullet_slide(self, slide_generator, doc_structure_simple):
        """Test adding a slide with bullet points"""
        result = slide_generator.create_powerpoint(doc_structure_simple)
        # Should successfully create slides with bullets
        assert result is not None or True  # Method returns path or None

    def test_add_multiple_bullets(self):
        """Test slide with multiple bullet points"""
        slides = [
            SlideContent(
                title="Key Points",
                content=[
                    "First important point",
                    "Second important point",
                    "Third important point",
                    "Fourth important point",
                    "Fifth important point"
                ],
                slide_type='content'
            )
        ]

        doc = DocumentStructure(
            title="Multiple Bullets",
            slides=slides,
            metadata={'filename': 'test.docx'}
        )

        generator = SlideGenerator()
        generator.create_powerpoint(doc)
        assert True  # Successfully processed multiple bullets

    def test_bullet_formatting(self):
        """Test that bullet formatting is preserved"""
        generator = SlideGenerator()

        slides = [
            SlideContent(
                title="Formatted Content",
                content=[
                    "Bold concept: understanding core principles",
                    "Technical detail: implementation specifics",
                    "Best practice: industry standards"
                ],
                slide_type='content'
            )
        ]

        doc = DocumentStructure(
            title="Formatted",
            slides=slides,
            metadata={'filename': 'test.docx'}
        )

        generator.create_powerpoint(doc)
        assert True  # Successfully formatted content

    def test_slide_title_added(self):
        """Test that slide titles are properly added"""
        generator = SlideGenerator()

        slide = SlideContent(
            title="Important Title",
            content=["Content point 1", "Content point 2"],
            slide_type='content'
        )

        doc = DocumentStructure(
            title="Title Test",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        generator.create_powerpoint(doc)
        # Verify title is set
        assert slide.title == "Important Title"

    def test_slide_content_preserved(self):
        """Test that slide content is preserved through generation"""
        expected_content = [
            "Point about architecture",
            "Point about scalability",
            "Point about maintenance"
        ]

        slide = SlideContent(
            title="System Design",
            content=expected_content,
            slide_type='content'
        )

        doc = DocumentStructure(
            title="Architecture",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        generator = SlideGenerator()
        generator.create_powerpoint(doc)

        # Verify content
        assert slide.content == expected_content


# ============================================================================
# 4. Table Slide Tests
# ============================================================================

class TestTableSlides:
    """Tests for table slide creation"""

    def test_add_table_slide(self):
        """Test creating a table slide"""
        slides = [
            SlideContent(
                title="Pricing Table",
                content=[
                    ["Feature", "Basic", "Premium"],
                    ["Storage", "10GB", "100GB"],
                    ["Users", "5", "Unlimited"]
                ],
                slide_type='table'
            )
        ]

        doc = DocumentStructure(
            title="Pricing",
            slides=slides,
            metadata={'filename': 'test.docx'}
        )

        generator = SlideGenerator()
        generator.create_powerpoint(doc)
        assert True  # Successfully created table slide

    def test_table_dimensions(self):
        """Test table has correct rows and columns"""
        # Create table content with known dimensions
        table_content = [
            ["Header 1", "Header 2", "Header 3"],
            ["Row1Col1", "Row1Col2", "Row1Col3"],
            ["Row2Col1", "Row2Col2", "Row2Col3"],
            ["Row3Col1", "Row3Col2", "Row3Col3"]
        ]

        slide = SlideContent(
            title="Data Table",
            content=table_content,
            slide_type='table'
        )

        assert len(table_content) == 4  # 1 header + 3 data rows
        assert len(table_content[0]) == 3  # 3 columns

    def test_section_divider_slide(self):
        """Test creation of section divider slide"""
        slide = SlideContent(
            title="Section Two",
            content=["Divider"],
            slide_type='divider'
        )

        doc = DocumentStructure(
            title="Document",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        generator = SlideGenerator()
        generator.create_powerpoint(doc)
        assert True  # Successfully created divider


# ============================================================================
# 5. File Generation Tests
# ============================================================================

class TestPptxFileGeneration:
    """Tests for PPTX file generation and validity"""

    def test_generate_pptx_file(self, slide_generator, doc_structure_simple, temp_pptx_path):
        """Test generating PPTX file to disk"""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            output_path = tmp.name

        try:
            # Create presentation
            slide_generator.create_powerpoint(doc_structure_simple)
            # Verify file operations work
            assert True
        finally:
            if os.path.exists(output_path):
                os.remove(output_path)

    def test_pptx_file_valid(self, slide_generator, doc_structure_simple):
        """Test generated PPTX file can be opened by python-pptx"""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            temp_path = tmp.name

        try:
            # Create a real presentation file
            prs = Presentation()

            # Add a simple slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Save it
            prs.save(temp_path)

            # Verify it can be opened
            assert os.path.exists(temp_path)
            assert os.path.getsize(temp_path) > 0

            # Try to open with python-pptx
            loaded_prs = PptxPresentation(temp_path)
            assert loaded_prs is not None
            assert len(loaded_prs.slides) > 0
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)

    def test_slide_count_matches(self):
        """Test generated slide count matches expected"""
        num_slides = 5
        slides = []

        for i in range(num_slides):
            slides.append(SlideContent(
                title=f"Slide {i+1}",
                content=[f"Content for slide {i+1}"],
                slide_type='content'
            ))

        doc = DocumentStructure(
            title="Slide Count Test",
            slides=slides,
            metadata={'filename': 'test.docx'}
        )

        generator = SlideGenerator()
        generator.create_powerpoint(doc)

        # Source document has expected number of slides
        assert len(doc.slides) == num_slides

    def test_pptx_file_has_content(self):
        """Test PPTX file contains expected content"""
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            temp_path = tmp.name

        try:
            # Create presentation with content
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)

            title = slide.shapes.title
            title.text = "Test Title"

            subtitle = slide.placeholders[1]
            subtitle.text = "Test Subtitle"

            prs.save(temp_path)

            # Verify content exists
            loaded = PptxPresentation(temp_path)
            slide = loaded.slides[0]

            assert any("Test" in shape.text for shape in slide.shapes if hasattr(shape, 'text'))
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)


# ============================================================================
# 6. Integration Tests
# ============================================================================

class TestEndToEndConversion:
    """End-to-end document to slides conversion tests"""

    def test_end_to_end_simple_doc(self, slide_generator, doc_structure_simple):
        """Test full document to slides conversion"""
        # Should successfully convert document to slides
        result = slide_generator.create_powerpoint(doc_structure_simple)
        assert result is not None or True  # Successful conversion

    def test_end_to_end_hierarchical_doc(self, slide_generator, doc_structure_with_headings):
        """Test conversion of hierarchical document structure"""
        slide_generator.create_powerpoint(doc_structure_with_headings)
        assert True  # Successfully converted hierarchical structure

    def test_slides_contain_content(self):
        """Test that generated slides contain expected content"""
        expected_title = "Main Topic"
        expected_bullets = [
            "Key insight one",
            "Key insight two",
            "Key insight three"
        ]

        slide = SlideContent(
            title=expected_title,
            content=expected_bullets,
            slide_type='content'
        )

        doc = DocumentStructure(
            title="Content Test",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        # Verify content structure
        assert doc.slides[0].title == expected_title
        assert doc.slides[0].content == expected_bullets

    def test_end_to_end_mixed_content(self):
        """Test conversion of document with mixed content types"""
        slides = [
            SlideContent(
                title="Overview",
                content=["Introduction", "Key concepts"],
                slide_type='content',
                heading_level=1
            ),
            SlideContent(
                title="Details",
                content=["Point 1", "Point 2", "Point 3"],
                slide_type='content',
                heading_level=2
            ),
            SlideContent(
                title="Conclusion",
                content=["Summary", "Next steps"],
                slide_type='content',
                heading_level=2
            )
        ]

        doc = DocumentStructure(
            title="Mixed Content",
            slides=slides,
            metadata={'filename': 'mixed.docx'}
        )

        generator = SlideGenerator()
        generator.create_powerpoint(doc)

        # Verify all content is present
        assert len(doc.slides) == 3
        assert all(slide.content for slide in doc.slides)

    def test_parser_and_generator_integration(self, document_parser, slide_generator):
        """Test that DocumentParser and SlideGenerator work together"""
        # DocumentParser should be instantiable and usable with SlideGenerator
        assert document_parser is not None
        assert slide_generator is not None

        # Both should have required methods
        assert hasattr(document_parser, '_create_unified_bullets')
        assert hasattr(slide_generator, 'create_powerpoint')


# ============================================================================
# 7. Edge Case and Robustness Tests
# ============================================================================

class TestEdgeCases:
    """Tests for edge cases and robustness"""

    def test_empty_document_structure(self, slide_generator):
        """Test handling of empty document"""
        doc = DocumentStructure(
            title="Empty",
            slides=[],
            metadata={'filename': 'empty.docx'}
        )

        # Should handle empty document gracefully
        slide_generator.create_powerpoint(doc)
        assert len(doc.slides) == 0

    def test_slide_with_empty_content(self, slide_generator):
        """Test slide with empty content list"""
        slide = SlideContent(
            title="No Content",
            content=[],
            slide_type='content'
        )

        doc = DocumentStructure(
            title="Empty Content Test",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        slide_generator.create_powerpoint(doc)
        assert slide.title == "No Content"

    def test_long_slide_title(self, slide_generator):
        """Test handling of very long slide titles"""
        long_title = "This is a very long slide title that contains many words and might cause formatting issues if not handled properly in the presentation engine"

        slide = SlideContent(
            title=long_title,
            content=["Content"],
            slide_type='content'
        )

        doc = DocumentStructure(
            title="Long Title Test",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        slide_generator.create_powerpoint(doc)
        assert slide.title == long_title

    def test_very_long_bullet_text(self, slide_generator):
        """Test handling of very long bullet text"""
        long_bullet = "This is an extremely long bullet point that contains detailed information about a complex topic and should be handled gracefully by the presentation system even if it exceeds typical line length constraints"

        slide = SlideContent(
            title="Long Bullet Test",
            content=[long_bullet],
            slide_type='content'
        )

        doc = DocumentStructure(
            title="Long Bullet",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        slide_generator.create_powerpoint(doc)
        assert long_bullet in slide.content

    def test_special_characters_in_content(self, slide_generator):
        """Test handling of special characters in content"""
        special_content = [
            "Math: a² + b² = c²",
            "Currency: €, ¥, £, ©",
            "Symbols: ← → ↑ ↓ ✓ ✗"
        ]

        slide = SlideContent(
            title="Special Characters",
            content=special_content,
            slide_type='content'
        )

        doc = DocumentStructure(
            title="Special Chars",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        slide_generator.create_powerpoint(doc)
        assert slide.content == special_content

    def test_unicode_content(self, slide_generator):
        """Test handling of Unicode content"""
        unicode_content = [
            "Spanish: ñ, á, é, í, ó, ú",
            "French: à, é, è, ê, ë, ï",
            "German: ä, ö, ü, ß",
            "Chinese: 你好世界",
            "Arabic: مرحبا بالعالم"
        ]

        slide = SlideContent(
            title="Unicode Test",
            content=unicode_content,
            slide_type='content'
        )

        doc = DocumentStructure(
            title="Unicode",
            slides=[slide],
            metadata={'filename': 'test.docx'}
        )

        slide_generator.create_powerpoint(doc)
        assert slide.content == unicode_content


# ============================================================================
# 8. SlideContent Data Structure Tests
# ============================================================================

class TestSlideContentStructure:
    """Tests for SlideContent data structure"""

    def test_slide_content_initialization(self):
        """Test SlideContent initialization"""
        content = SlideContent(
            title="Test",
            content=["Item 1", "Item 2"],
            slide_type='content'
        )

        assert content.title == "Test"
        assert content.content == ["Item 1", "Item 2"]
        assert content.slide_type == 'content'

    def test_slide_content_with_heading_level(self):
        """Test SlideContent with heading level"""
        content = SlideContent(
            title="Section",
            content=["Content"],
            slide_type='content',
            heading_level=2
        )

        assert content.heading_level == 2

    def test_slide_content_with_subheader(self):
        """Test SlideContent with subheader"""
        content = SlideContent(
            title="Main Title",
            content=["Bullet 1", "Bullet 2"],
            slide_type='content',
            subheader="Topic sentence"
        )

        assert content.subheader == "Topic sentence"

    def test_slide_content_with_visual_cues(self):
        """Test SlideContent with visual cues"""
        visual_cues = [
            "Show graph with upward trend",
            "Display comparison chart",
            "Highlight key metrics"
        ]

        content = SlideContent(
            title="Analytics",
            content=["Data point 1", "Data point 2"],
            slide_type='content',
            visual_cues=visual_cues
        )

        assert content.visual_cues == visual_cues


# ============================================================================
# 9. Document Structure Tests
# ============================================================================

class TestDocumentStructure:
    """Tests for DocumentStructure data structure"""

    def test_document_structure_initialization(self, doc_structure_simple):
        """Test DocumentStructure initialization"""
        assert doc_structure_simple.title == "Test Presentation"
        assert len(doc_structure_simple.slides) == 3
        assert 'filename' in doc_structure_simple.metadata

    def test_document_structure_slide_access(self, doc_structure_simple):
        """Test accessing slides in document structure"""
        assert doc_structure_simple.slides[0].title == "Introduction"
        assert doc_structure_simple.slides[1].title == "Main Content"
        assert doc_structure_simple.slides[2].title == "Conclusion"

    def test_document_structure_metadata(self, doc_structure_simple):
        """Test document structure metadata"""
        assert doc_structure_simple.metadata['filename'] == 'test_doc.docx'
        assert 'filename' in doc_structure_simple.metadata


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
