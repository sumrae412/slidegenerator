"""
Enhanced Script to Slides Generator with Modern UI and Advanced Features

A Flask web application that converts uploaded documents to professional presentations
with template support, real-time progress, and enhanced processing options.

Supported formats: DOCX, PDF, TXT
"""

import os
import json
import logging
import time
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from datetime import datetime
import re
import random
from math import cos, sin
import requests
import warnings
import hashlib
from concurrent.futures import ThreadPoolExecutor
import threading

# Semantic analysis libraries - lightweight fallback approach
try:
    import nltk
    import textstat
    from collections import Counter
    LIGHTWEIGHT_SEMANTIC = True
    
    # Try to download required NLTK data if not present
    try:
        nltk.data.find('tokenizers/punkt')
        nltk.data.find('corpora/stopwords')
        nltk.data.find('taggers/averaged_perceptron_tagger')
    except LookupError:
        logger = logging.getLogger(__name__)
        logger.info("Downloading required NLTK data...")
        try:
            nltk.download('punkt', quiet=True)
            nltk.download('stopwords', quiet=True)  
            nltk.download('averaged_perceptron_tagger', quiet=True)
        except:
            logger.warning("Could not download NLTK data - basic analysis only")
            
except ImportError:
    logging.warning("Lightweight semantic libraries not available - using basic fallback")
    LIGHTWEIGHT_SEMANTIC = False

import flask
from flask import Flask, render_template, request, jsonify, send_file, redirect, url_for, session
from werkzeug.utils import secure_filename
from flask_cors import CORS

# Document processing libraries
import docx
from docx import Document
import PyPDF2
import markdown
from bs4 import BeautifulSoup

# Presentation generation
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# Image generation for sketches
from PIL import Image, ImageDraw, ImageFont
import io
import base64

# OpenAI for bullet point generation
import openai

app = Flask(__name__)
app.secret_key = os.environ.get('FLASK_SECRET_KEY', 'enhanced-secret-key-2024')
CORS(app)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# File upload configuration
UPLOAD_FOLDER = 'uploads'
EXPORT_FOLDER = 'exports'
TEMP_FOLDER = 'temp'
ALLOWED_EXTENSIONS = {'docx', 'pdf', 'txt'}
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB

# Create necessary directories
for folder in [UPLOAD_FOLDER, EXPORT_FOLDER, TEMP_FOLDER]:
    os.makedirs(folder, exist_ok=True)

# Template configurations
TEMPLATES = {
    'custom': {
        'name': 'Custom Template',
        'template_file': 'slide_template.pptx',
        'colors': {
            'primary': RGBColor(51, 51, 51),
            'secondary': RGBColor(102, 102, 102),
            'accent': RGBColor(0, 120, 215),
            'background': RGBColor(255, 255, 255)
        },
        'fonts': {
            'title': 'Calibri',
            'body': 'Calibri Light'
        }
    },
    'professional': {
        'name': 'Professional',
        'colors': {
            'primary': RGBColor(51, 51, 51),
            'secondary': RGBColor(102, 102, 102),
            'accent': RGBColor(0, 120, 215),
            'background': RGBColor(255, 255, 255)
        },
        'fonts': {
            'title': 'Calibri',
            'body': 'Calibri Light'
        }
    },
    'creative': {
        'name': 'Creative',
        'colors': {
            'primary': RGBColor(147, 51, 234),
            'secondary': RGBColor(236, 72, 153),
            'accent': RGBColor(59, 130, 246),
            'background': RGBColor(254, 243, 255)
        },
        'fonts': {
            'title': 'Arial Black',
            'body': 'Arial'
        }
    },
    'minimal': {
        'name': 'Minimal',
        'colors': {
            'primary': RGBColor(31, 41, 55),
            'secondary': RGBColor(107, 114, 128),
            'accent': RGBColor(16, 185, 129),
            'background': RGBColor(249, 250, 251)
        },
        'fonts': {
            'title': 'Helvetica',
            'body': 'Helvetica Neue'
        }
    }
}

# Progress tracking
conversion_progress = {}
progress_lock = threading.Lock()

def update_progress(session_id: str, step: str, percentage: int, detail: str = ""):
    """Update conversion progress for real-time tracking"""
    with progress_lock:
        conversion_progress[session_id] = {
            'step': step,
            'percentage': percentage,
            'detail': detail,
            'timestamp': time.time()
        }

@dataclass
class ProcessingOptions:
    """Configuration for document processing"""
    template: str = 'custom'
    bullet_density: int = 3  # 1-5 scale
    include_visual_prompts: bool = True
    generate_speaker_notes: bool = False
    add_slide_numbers: bool = True
    slides_per_section: int = 5
    column_index: int = 0
    api_key: Optional[str] = None

@dataclass
class SlideContent:
    """Represents content for a single slide"""
    title: str
    bullets: List[str]
    visual_prompt: Optional[str] = None
    speaker_notes: Optional[str] = None
    slide_type: str = 'content'  # 'title', 'section', 'content'
    template_style: Optional[Dict] = None

class DocumentProcessor:
    """Enhanced document processor with multiple format support"""
    
    def __init__(self, options: ProcessingOptions):
        self.options = options
        self.slides = []
        
    def process_file(self, file_path: str, file_type: str, session_id: str) -> List[SlideContent]:
        """Process uploaded file based on type"""
        update_progress(session_id, "Reading document", 10, f"Processing {file_type} file")
        
        if file_type == 'docx':
            content = self._process_docx(file_path)
        elif file_type == 'pdf':
            content = self._process_pdf(file_path)
        elif file_type == 'txt':
            content = self._process_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
            
        update_progress(session_id, "Analyzing content", 30, "Extracting structure and text")
        
        # Generate slides from content
        slides = self._generate_slides(content, session_id)
        
        return slides
    
    def _process_docx(self, file_path: str) -> Dict:
        """Process DOCX file"""
        doc = Document(file_path)
        content = {
            'paragraphs': [],
            'headings': [],
            'tables': []
        }
        
        for para in doc.paragraphs:
            if para.style.name.startswith('Heading'):
                level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                content['headings'].append({
                    'level': level,
                    'text': para.text
                })
            elif para.text.strip():
                content['paragraphs'].append(para.text)
        
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            content['tables'].append(table_data)
        
        return content
    
    def _process_pdf(self, file_path: str) -> Dict:
        """Process PDF file"""
        content = {
            'paragraphs': [],
            'headings': [],
            'tables': []
        }
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        
        # Simple paragraph splitting
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        for para in paragraphs:
            # Detect headings (simple heuristic: short lines, possibly all caps)
            if len(para) < 100 and (para.isupper() or para[0].isupper()):
                content['headings'].append({
                    'level': 2 if para.isupper() else 3,
                    'text': para
                })
            else:
                content['paragraphs'].append(para)
        
        return content
    
    def _process_txt(self, file_path: str) -> Dict:
        """Process TXT file"""
        content = {
            'paragraphs': [],
            'headings': [],
            'tables': []
        }
        
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        
        # Process markdown-style headings
        lines = text.split('\n')
        for line in lines:
            if line.startswith('#'):
                level = len(line) - len(line.lstrip('#'))
                heading_text = line.lstrip('#').strip()
                if heading_text:
                    content['headings'].append({
                        'level': min(level, 4),
                        'text': heading_text
                    })
            elif line.strip():
                content['paragraphs'].append(line.strip())
        
        return content
    
    def _generate_slides(self, content: Dict, session_id: str) -> List[SlideContent]:
        """Generate slides from processed content"""
        slides = []
        
        update_progress(session_id, "Generating slides", 50, "Creating slide structure")
        
        # Add title slide if we have headings
        if content['headings']:
            title = content['headings'][0]['text'] if content['headings'] else "Presentation"
            slides.append(SlideContent(
                title=title,
                bullets=[],
                slide_type='title',
                template_style=TEMPLATES[self.options.template]
            ))
        
        # Process content into slides
        if self.options.column_index == 0:  # No table - use paragraphs
            slides.extend(self._process_paragraphs(content['paragraphs'], content['headings'], session_id))
        else:  # Use table column
            if content['tables']:
                slides.extend(self._process_table_column(content['tables'], self.options.column_index - 1, session_id))
        
        return slides
    
    def _process_paragraphs(self, paragraphs: List[str], headings: List[Dict], session_id: str) -> List[SlideContent]:
        """Process paragraphs into slides with intelligent chunking"""
        slides = []
        current_section = None
        
        # Group content by headings
        content_groups = []
        current_group = {'heading': None, 'content': []}
        
        heading_idx = 0
        for para in paragraphs:
            # Check if we should insert a heading
            if heading_idx < len(headings):
                # Simple heuristic: insert heading before related content
                current_group = {'heading': headings[heading_idx], 'content': [para]}
                content_groups.append(current_group)
                heading_idx += 1
            else:
                if current_group['content']:
                    current_group['content'].append(para)
                else:
                    current_group = {'heading': None, 'content': [para]}
                    content_groups.append(current_group)
        
        # Generate slides from groups
        for idx, group in enumerate(content_groups):
            update_progress(session_id, "Processing content", 
                          50 + int((idx / len(content_groups)) * 30),
                          f"Processing section {idx + 1}/{len(content_groups)}")
            
            if group['heading']:
                # Add section slide for major headings
                if group['heading']['level'] <= 2:
                    slides.append(SlideContent(
                        title=group['heading']['text'],
                        bullets=[],
                        slide_type='section',
                        template_style=TEMPLATES[self.options.template]
                    ))
            
            # Process content into bullet points
            if group['content']:
                bullets = self._generate_bullets(' '.join(group['content']), session_id)
                
                # Split into multiple slides if needed
                bullets_per_slide = 4 + self.options.bullet_density
                for i in range(0, len(bullets), bullets_per_slide):
                    slide_bullets = bullets[i:i + bullets_per_slide]
                    
                    # Generate visual prompt if enabled
                    visual_prompt = None
                    if self.options.include_visual_prompts:
                        visual_prompt = self._generate_visual_prompt(slide_bullets)
                    
                    # Generate speaker notes if enabled
                    speaker_notes = None
                    if self.options.generate_speaker_notes:
                        speaker_notes = self._generate_speaker_notes(slide_bullets)
                    
                    slides.append(SlideContent(
                        title=group['heading']['text'] if group['heading'] else f"Slide {len(slides) + 1}",
                        bullets=slide_bullets,
                        visual_prompt=visual_prompt,
                        speaker_notes=speaker_notes,
                        slide_type='content',
                        template_style=TEMPLATES[self.options.template]
                    ))
        
        return slides
    
    def _process_table_column(self, tables: List[List[List[str]]], column_idx: int, session_id: str) -> List[SlideContent]:
        """Process specific table column into slides"""
        slides = []
        
        for table in tables:
            for row_idx, row in enumerate(table):
                if column_idx < len(row) and row[column_idx].strip():
                    update_progress(session_id, "Processing table", 
                                  50 + int((row_idx / len(table)) * 30),
                                  f"Processing row {row_idx + 1}/{len(table)}")
                    
                    content = row[column_idx]
                    bullets = self._generate_bullets(content, session_id)
                    
                    visual_prompt = None
                    if self.options.include_visual_prompts:
                        visual_prompt = self._generate_visual_prompt(bullets)
                    
                    slides.append(SlideContent(
                        title=f"Slide {len(slides) + 1}",
                        bullets=bullets,
                        visual_prompt=visual_prompt,
                        slide_type='content',
                        template_style=TEMPLATES[self.options.template]
                    ))
        
        return slides
    
    def _generate_bullets(self, text: str, session_id: str) -> List[str]:
        """Generate bullet points from text using AI or fallback methods"""
        if self.options.api_key:
            try:
                return self._generate_bullets_ai(text)
            except Exception as e:
                logger.warning(f"AI bullet generation failed: {e}")
                return self._generate_bullets_fallback(text)
        else:
            return self._generate_bullets_fallback(text)
    
    def _generate_bullets_ai(self, text: str) -> List[str]:
        """Generate bullets using OpenAI API"""
        openai.api_key = self.options.api_key
        
        density_map = {
            1: "very concise",
            2: "concise", 
            3: "balanced",
            4: "detailed",
            5: "very detailed"
        }
        
        prompt = f"""Convert the following text into {density_map[self.options.bullet_density]} bullet points for a presentation slide:

{text}

Requirements:
- Create {3 + self.options.bullet_density} bullet points
- Each bullet should be a complete, clear statement
- Focus on key information and insights
- Make bullets engaging and informative"""

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=200,
            temperature=0.7
        )
        
        bullets = response.choices[0].message.content.strip().split('\n')
        bullets = [b.strip('- •').strip() for b in bullets if b.strip()]
        
        return bullets[:3 + self.options.bullet_density]
    
    def _generate_bullets_fallback(self, text: str) -> List[str]:
        """Generate bullets using NLP or basic text processing with improved prioritization"""
        if LIGHTWEIGHT_SEMANTIC:
            try:
                # Use NLTK for sentence extraction
                sentences = nltk.sent_tokenize(text)
                
                # Enhanced scoring system for better prioritization
                sentence_scores = {}
                
                # Define high-importance keywords and patterns
                important_keywords = {
                    # Financial/Business terms
                    'revenue', 'profit', 'earnings', 'growth', 'increase', 'decrease',
                    'sales', 'market', 'strategy', 'investment', 'cost', 'savings',
                    'performance', 'results', 'achievement', 'success', 'target',
                    'goal', 'objective', 'improvement', 'efficiency', 'productivity',
                    # Technical/Process terms
                    'analysis', 'research', 'development', 'implementation', 'solution',
                    'innovation', 'technology', 'system', 'process', 'methodology',
                    # Key business metrics
                    'percent', '%', 'million', 'billion', 'thousand', 'quarterly',
                    'annual', 'monthly', 'year', 'quarter', 'growth rate'
                }
                
                # Low-priority/trivial keywords
                trivial_keywords = {
                    'weather', 'cafeteria', 'lunch', 'menu', 'parking', 'coffee',
                    'break', 'restroom', 'temperature', 'climate', 'decoration',
                    'furniture', 'color', 'paint', 'carpet', 'lighting', 'music'
                }
                
                for sent in sentences:
                    words = nltk.word_tokenize(sent.lower())
                    words = [w for w in words if w.isalnum()]
                    
                    if not words:
                        sentence_scores[sent] = 0
                        continue
                    
                    # Base word frequency score
                    word_freq = Counter(words)
                    base_score = sum(word_freq.values()) / len(words)
                    
                    # Importance multiplier based on keywords
                    importance_multiplier = 1.0
                    
                    # Boost for important business terms
                    important_word_count = sum(1 for w in words if w in important_keywords)
                    if important_word_count > 0:
                        importance_multiplier += important_word_count * 2.0
                    
                    # Penalty for trivial terms
                    trivial_word_count = sum(1 for w in words if w in trivial_keywords)
                    if trivial_word_count > 0:
                        importance_multiplier -= trivial_word_count * 1.5
                    
                    # Boost for numerical data (likely metrics)
                    numerical_boost = 0
                    for word in words:
                        if any(char.isdigit() for char in word):
                            numerical_boost += 2.0  # Increased from 1.5 to better preserve facts
                    
                    # Position bias - earlier sentences often more important
                    position_boost = 1.0 - (sentences.index(sent) / len(sentences)) * 0.3
                    
                    # Length consideration - very short or very long sentences less important
                    length_penalty = 1.0
                    if len(words) < 5:
                        length_penalty = 0.5
                    elif len(words) > 30:
                        length_penalty = 0.8
                    
                    # Calculate final score
                    final_score = (base_score + numerical_boost) * importance_multiplier * position_boost * length_penalty
                    sentence_scores[sent] = max(0, final_score)  # Ensure non-negative
                
                # Get top sentences based on enhanced scoring
                top_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)
                
                # Select bullets with diversity consideration
                bullets = []
                used_keywords = set()
                target_count = 3 + self.options.bullet_density
                
                for sent, score in top_sentences:
                    if len(bullets) >= target_count:
                        break
                    
                    # Check for keyword diversity to avoid redundant bullets
                    sent_words = set(nltk.word_tokenize(sent.lower()))
                    overlap = len(sent_words & used_keywords)
                    
                    # Dynamic overlap threshold based on progress toward target
                    progress = len(bullets) / target_count
                    
                    if progress < 0.5:
                        # Early phase: stricter diversity for quality
                        max_overlap = 3
                    elif progress < 0.8:
                        # Mid phase: moderate diversity
                        max_overlap = 5
                    else:
                        # Late phase: relaxed diversity to meet density target
                        max_overlap = 8
                    
                    # Always accept if score is positive and overlap is acceptable
                    # Or if we're far from target and need more bullets
                    if score > 0 and (overlap <= max_overlap or len(bullets) < max(2, target_count - 2)):
                        bullets.append(sent)
                        used_keywords.update(sent_words)
                
                # Clean up bullets
                bullets = [self._clean_bullet(b) for b in bullets]
                
                return bullets
            except Exception as e:
                logger.warning(f"Enhanced NLP bullet generation failed: {e}")
        
        # Basic fallback: split into sentences with simple importance filtering
        sentences = text.replace('!', '.').replace('?', '.').split('.')
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
        
        # Simple importance filtering even for basic fallback
        important_sentences = []
        other_sentences = []
        
        for sent in sentences:
            sent_lower = sent.lower()
            # Check for numerical data or business terms
            has_numbers = any(char.isdigit() for char in sent)
            has_business_terms = any(term in sent_lower for term in 
                                   ['revenue', 'profit', 'growth', 'sales', 'market', 'percent', '%'])
            has_trivial_terms = any(term in sent_lower for term in 
                                   ['weather', 'cafeteria', 'lunch', 'parking', 'coffee'])
            
            if (has_numbers or has_business_terms) and not has_trivial_terms:
                important_sentences.append(sent)
            elif not has_trivial_terms:
                other_sentences.append(sent)
        
        # Prioritize important sentences
        num_bullets = min(3 + self.options.bullet_density, len(sentences))
        bullets = (important_sentences + other_sentences)[:num_bullets]
        
        return [self._clean_bullet(b) for b in bullets]
    
    def _clean_bullet(self, text: str) -> str:
        """Clean and format bullet point text"""
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Ensure proper capitalization
        if text and text[0].islower():
            text = text[0].upper() + text[1:]
        
        # Add period if missing
        if text and text[-1] not in '.!?':
            text += '.'
        
        # Limit length
        if len(text) > 150:
            text = text[:147] + '...'
        
        return text
    
    def _generate_visual_prompt(self, bullets: List[str]) -> str:
        """Generate a visual prompt for slide illustration"""
        key_concepts = []
        
        for bullet in bullets:
            # Extract key words (simple approach)
            words = bullet.lower().split()
            important_words = [w for w in words if len(w) > 4 and w.isalpha()]
            key_concepts.extend(important_words[:2])
        
        # Create visual prompt
        concepts = list(set(key_concepts))[:3]
        if concepts:
            return f"Visual concept: {', '.join(concepts)} - modern, professional diagram or illustration"
        else:
            return "Visual concept: Abstract professional diagram or illustration"
    
    def _generate_speaker_notes(self, bullets: List[str]) -> str:
        """Generate speaker notes for the slide"""
        notes = "Speaker Notes:\n\n"
        
        for i, bullet in enumerate(bullets, 1):
            notes += f"{i}. Expand on: {bullet}\n"
            notes += "   - Provide examples or context\n"
            notes += "   - Engage audience with questions\n\n"
        
        notes += "Transition: Connect to next topic\n"
        
        return notes

class PresentationGenerator:
    """Enhanced presentation generator with template support"""
    
    def __init__(self, slides: List[SlideContent], options: ProcessingOptions):
        self.slides = slides
        self.options = options
        self.template = TEMPLATES[options.template]
        
    def generate(self, session_id: str) -> bytes:
        """Generate PowerPoint presentation"""
        update_progress(session_id, "Creating presentation", 85, "Generating PowerPoint file")
        
        # Load from template file if specified, otherwise create new presentation
        if 'template_file' in self.template and os.path.exists(self.template['template_file']):
            prs = Presentation(self.template['template_file'])
            # Remove existing slides from template, keep only master slides
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            prs = Presentation()
            # Set presentation size (16:9 widescreen)
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
        
        for idx, slide_content in enumerate(self.slides):
            update_progress(session_id, "Adding slides", 
                          85 + int((idx / len(self.slides)) * 10),
                          f"Creating slide {idx + 1}/{len(self.slides)}")
            
            if slide_content.slide_type == 'title':
                self._add_title_slide(prs, slide_content)
            elif slide_content.slide_type == 'section':
                self._add_section_slide(prs, slide_content)
            else:
                self._add_content_slide(prs, slide_content, idx + 1)
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        update_progress(session_id, "Complete", 100, "Presentation ready for download")
        
        return output.getvalue()
    
    def _add_title_slide(self, prs: Presentation, content: SlideContent):
        """Add title slide with template styling"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = content.title
        
        # Apply template styling
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.template['fonts']['title']
                run.font.size = Pt(44)
                run.font.color.rgb = self.template['colors']['primary']
                run.font.bold = True
        
        # Add subtitle with date
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            for paragraph in subtitle.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = self.template['fonts']['body']
                    run.font.size = Pt(20)
                    run.font.color.rgb = self.template['colors']['secondary']
    
    def _add_section_slide(self, prs: Presentation, content: SlideContent):
        """Add section divider slide"""
        slide_layout = prs.slide_layouts[2]  # Section header
        slide = prs.slides.add_slide(slide_layout)
        
        # Set section title
        title = slide.shapes.title
        title.text = content.title
        
        # Apply template styling
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.template['fonts']['title']
                run.font.size = Pt(40)
                run.font.color.rgb = self.template['colors']['accent']
                run.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, prs: Presentation, content: SlideContent, slide_num: int):
        """Add content slide with bullets and visual elements"""
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = content.title
        
        # Style title
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.template['fonts']['title']
                run.font.size = Pt(32)
                run.font.color.rgb = self.template['colors']['primary']
                run.font.bold = True
        
        # Add content
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            
            # Add bullets
            for bullet in content.bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                
                # Style bullets
                p.font.name = self.template['fonts']['body']
                p.font.size = Pt(18)
                p.font.color.rgb = self.template['colors']['primary']
                p.space_after = Pt(12)
            
            # Add visual prompt as a text box if enabled
            if content.visual_prompt and self.options.include_visual_prompts:
                left = Inches(0.5)
                top = Inches(4.5)
                width = Inches(9)
                height = Inches(0.5)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = content.visual_prompt
                p.font.size = Pt(10)
                p.font.italic = True
                p.font.color.rgb = self.template['colors']['secondary']
        
        # Add slide number if enabled
        if self.options.add_slide_numbers:
            left = Inches(9)
            top = Inches(5)
            width = Inches(0.5)
            height = Inches(0.3)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = str(slide_num)
            p.font.size = Pt(10)
            p.font.color.rgb = self.template['colors']['secondary']
            p.alignment = PP_ALIGN.RIGHT
        
        # Add speaker notes if generated
        if content.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content.speaker_notes

# Flask routes

@app.route('/')
def index():
    """Render the enhanced main page"""
    return render_template('file_to_slides_enhanced.html')

@app.route('/favicon.ico')
def favicon():
    """Serve favicon"""
    return send_file('static/favicon.ico', mimetype='image/x-icon')

@app.route('/convert', methods=['POST'])
def convert():
    """Handle file conversion with progress tracking"""
    try:
        # Generate session ID for progress tracking
        session_id = hashlib.md5(f"{time.time()}".encode()).hexdigest()
        
        # Check if file is present
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Check file size
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        
        if file_size > MAX_FILE_SIZE:
            return jsonify({'error': f'File too large. Maximum size is {MAX_FILE_SIZE / 1024 / 1024}MB'}), 400
        
        # Check file extension
        file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        if file_ext not in ALLOWED_EXTENSIONS:
            return jsonify({'error': f'Invalid file type. Allowed types: {", ".join(ALLOWED_EXTENSIONS)}'}), 400
        
        # Save uploaded file
        filename = secure_filename(file.filename)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        unique_filename = f"{timestamp}_{filename}"
        file_path = os.path.join(UPLOAD_FOLDER, unique_filename)
        
        # Extract base filename without extension for output naming
        base_filename = os.path.splitext(filename)[0]
        file.save(file_path)
        
        # Parse processing options
        options = ProcessingOptions(
            template=request.form.get('template', 'professional'),
            bullet_density=int(request.form.get('bullet_density', 3)),
            include_visual_prompts=request.form.get('include_visual_prompts', 'true') == 'true',
            generate_speaker_notes=request.form.get('generate_speaker_notes', 'false') == 'true',
            add_slide_numbers=request.form.get('add_slide_numbers', 'true') == 'true',
            slides_per_section=int(request.form.get('slides_per_section', 5)),
            column_index=int(request.form.get('column_index', 0)),
            api_key=request.form.get('api_key', '').strip() or None
        )
        
        # Process document
        processor = DocumentProcessor(options)
        slides = processor.process_file(file_path, file_ext, session_id)
        
        # Generate presentation
        generator = PresentationGenerator(slides, options)
        pptx_data = generator.generate(session_id)
        
        # Clean up uploaded file
        os.remove(file_path)
        
        # Save presentation with new naming format: MMDDYYYY_<FILENAME>.pptx
        output_date = datetime.now().strftime('%m%d%Y')
        output_filename = f"{output_date}_{base_filename}.pptx"
        output_path = os.path.join(EXPORT_FOLDER, output_filename)
        
        with open(output_path, 'wb') as f:
            f.write(pptx_data)
        
        # Return file for download
        return send_file(
            io.BytesIO(pptx_data),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=output_filename
        )
        
    except Exception as e:
        logger.error(f"Conversion error: {str(e)}", exc_info=True)
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500

@app.route('/progress/<session_id>')
def get_progress(session_id):
    """Get conversion progress for real-time updates"""
    with progress_lock:
        progress = conversion_progress.get(session_id, {
            'step': 'Initializing',
            'percentage': 0,
            'detail': ''
        })
    
    return jsonify(progress)

@app.route('/sample')
def download_sample():
    """Provide sample document for testing"""
    sample_content = """# Sample Presentation Document

## Introduction Section

### About This Tool
This tool converts your documents into professional PowerPoint presentations using AI technology.

### Key Features
- Automatic slide generation from text content
- AI-powered bullet point creation  
- Multiple template options
- Visual prompt generation for graphics

## How It Works

### Step 1: Upload Your Document
Simply drag and drop your Word document, PDF, or text file into the upload area.

### Step 2: Configure Settings
Choose your preferred template and customize the output options to match your needs.

### Step 3: Generate Presentation
Our AI processes your content and creates a professional presentation ready for use.

## Advanced Features

### Template Selection
Choose from Professional, Creative, or Minimal templates to match your presentation style.

### AI Enhancement
Use OpenAI integration for more intelligent bullet point generation and content organization.

### Real-time Progress
Track the conversion progress with our live status updates.

## Conclusion

### Get Started Today
Transform your documents into engaging presentations with just a few clicks.

### Contact Information
Visit our website for more information and support."""

    # Create a temporary file
    output = io.BytesIO()
    output.write(sample_content.encode('utf-8'))
    output.seek(0)
    
    return send_file(
        output,
        mimetype='text/plain',
        as_attachment=True,
        download_name='sample_presentation_document.txt'
    )

@app.errorhandler(413)
def request_entity_too_large(error):
    """Handle file too large error"""
    return jsonify({'error': f'File too large. Maximum size is {MAX_FILE_SIZE / 1024 / 1024}MB'}), 413

@app.errorhandler(500)
def internal_server_error(error):
    """Handle internal server errors"""
    logger.error(f"Internal server error: {error}")
    return jsonify({'error': 'An internal error occurred. Please try again.'}), 500

# Cleanup old files periodically
def cleanup_old_files():
    """Remove old temporary files"""
    try:
        current_time = time.time()
        
        for folder in [UPLOAD_FOLDER, EXPORT_FOLDER, TEMP_FOLDER]:
            if os.path.exists(folder):
                for filename in os.listdir(folder):
                    file_path = os.path.join(folder, filename)
                    # Remove files older than 1 hour
                    if os.path.isfile(file_path):
                        if current_time - os.path.getmtime(file_path) > 3600:
                            os.remove(file_path)
                            logger.info(f"Cleaned up old file: {file_path}")
    except Exception as e:
        logger.error(f"Cleanup error: {e}")

# Schedule cleanup
def schedule_cleanup():
    """Schedule periodic cleanup"""
    while True:
        time.sleep(1800)  # Run every 30 minutes
        cleanup_old_files()

# Start cleanup thread
import threading
cleanup_thread = threading.Thread(target=schedule_cleanup, daemon=True)
cleanup_thread.start()

if __name__ == '__main__':
    # Run the application
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port, debug=False)