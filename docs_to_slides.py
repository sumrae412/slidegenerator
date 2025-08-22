"""
Google Docs to Slides Generator

A Flask web application that converts Google Docs to presentation slides
with support for multiple export formats (PowerPoint, Google Slides, HTML).
"""

import os
import json
import logging
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from datetime import datetime
import re

import flask
from flask import Flask, render_template, request, jsonify, send_file, redirect, url_for
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

app = Flask(__name__)
app.secret_key = os.environ.get('FLASK_SECRET_KEY', 'your-secret-key-change-this')

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Google API configuration
SCOPES = [
    'https://www.googleapis.com/auth/documents.readonly',
    'https://www.googleapis.com/auth/presentations',
    'https://www.googleapis.com/auth/drive.readonly'
]

# OAuth2 configuration
CLIENT_SECRETS_FILE = 'credentials.json'
REDIRECT_URI = 'http://localhost:5000/callback'

@dataclass
class SlideContent:
    """Represents content for a single slide"""
    title: str
    content: List[str]
    slide_type: str = 'content'  # 'title', 'content', 'image', 'bullet'
    image_url: Optional[str] = None

@dataclass
class DocumentStructure:
    """Represents the parsed structure of a Google Doc"""
    title: str
    slides: List[SlideContent]
    metadata: Dict[str, Any]

class GoogleDocsExtractor:
    """Handles extraction and parsing of Google Docs content"""
    
    def __init__(self, credentials: Credentials):
        self.credentials = credentials
        self.docs_service = build('docs', 'v1', credentials=credentials)
        self.drive_service = build('drive', 'v3', credentials=credentials)
    
    def extract_document(self, document_id: str) -> DocumentStructure:
        """Extract content from a Google Doc and structure it for slides"""
        try:
            # Get document metadata
            doc_metadata = self.drive_service.files().get(fileId=document_id).execute()
            
            # Get document content
            document = self.docs_service.documents().get(documentId=document_id).execute()
            
            doc_title = document.get('title', 'Untitled Document')
            content = document.get('body', {}).get('content', [])
            
            slides = self._parse_content_to_slides(content)
            
            metadata = {
                'doc_id': document_id,
                'original_title': doc_title,
                'created_time': doc_metadata.get('createdTime'),
                'modified_time': doc_metadata.get('modifiedTime'),
                'export_time': datetime.now().isoformat()
            }
            
            return DocumentStructure(
                title=doc_title,
                slides=slides,
                metadata=metadata
            )
            
        except HttpError as error:
            logger.error(f"Error extracting document {document_id}: {error}")
            raise
    
    def _parse_content_to_slides(self, content: List[Dict]) -> List[SlideContent]:
        """Parse document content and convert to slide structure"""
        slides = []
        current_slide = None
        slide_content = []
        
        for element in content:
            if 'paragraph' in element:
                paragraph = element['paragraph']
                text = self._extract_text_from_paragraph(paragraph)
                
                if not text.strip():
                    continue
                
                # Check if this is a heading (potential slide title)
                if self._is_heading(paragraph):
                    # Save previous slide if exists
                    if current_slide:
                        slides.append(SlideContent(
                            title=current_slide,
                            content=slide_content,
                            slide_type='content'
                        ))
                    
                    # Start new slide
                    current_slide = text.strip()
                    slide_content = []
                else:
                    # Add to current slide content
                    if text.strip():
                        slide_content.append(text.strip())
        
        # Add final slide
        if current_slide:
            slides.append(SlideContent(
                title=current_slide,
                content=slide_content,
                slide_type='content'
            ))
        
        # If no headings found, create slides from content blocks
        if not slides and slide_content:
            slides.append(SlideContent(
                title="Content",
                content=slide_content,
                slide_type='content'
            ))
        
        return slides
    
    def _extract_text_from_paragraph(self, paragraph: Dict) -> str:
        """Extract plain text from a paragraph element"""
        text_parts = []
        
        for element in paragraph.get('elements', []):
            if 'textRun' in element:
                text_parts.append(element['textRun'].get('content', ''))
        
        return ''.join(text_parts)
    
    def _is_heading(self, paragraph: Dict) -> bool:
        """Determine if a paragraph is a heading based on styling"""
        paragraph_style = paragraph.get('paragraphStyle', {})
        named_style_type = paragraph_style.get('namedStyleType', '')
        
        return named_style_type.startswith('HEADING')

class SlideGenerator:
    """Handles generation of presentation slides in various formats"""
    
    def __init__(self):
        self.template_layouts = {
            'title': 0,
            'content': 1,
            'bullet': 1,
            'image': 6
        }
    
    def create_powerpoint(self, doc_structure: DocumentStructure) -> str:
        """Generate PowerPoint presentation from document structure"""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = doc_structure.title
        subtitle.text = f"Generated from Google Docs\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Content slides
        for slide_content in doc_structure.slides:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_content.title
            
            # Add content
            if slide_content.content:
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                
                # Add content as bullet points
                for i, content_item in enumerate(slide_content.content):
                    if i == 0:
                        text_frame.text = content_item
                    else:
                        p = text_frame.add_paragraph()
                        p.text = content_item
                        p.level = 0
        
        # Save presentation
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join('exports', filename)
        os.makedirs('exports', exist_ok=True)
        prs.save(filepath)
        
        return filepath
    
    def create_html_slides(self, doc_structure: DocumentStructure) -> str:
        """Generate HTML presentation from document structure"""
        html_template = self._get_html_template()
        
        slides_html = []
        
        # Title slide
        title_slide = f"""
        <section class="slide title-slide">
            <h1>{doc_structure.title}</h1>
            <p class="subtitle">Generated from Google Docs</p>
            <p class="date">{datetime.now().strftime('%B %d, %Y')}</p>
        </section>
        """
        slides_html.append(title_slide)
        
        # Content slides
        for slide_content in doc_structure.slides:
            content_items = '</li><li>'.join(slide_content.content) if slide_content.content else ''
            
            slide_html = f"""
            <section class="slide content-slide">
                <h2>{slide_content.title}</h2>
                {f'<ul><li>{content_items}</li></ul>' if content_items else ''}
            </section>
            """
            slides_html.append(slide_html)
        
        # Combine slides with template
        full_html = html_template.replace('{{SLIDES}}', '\n'.join(slides_html))
        full_html = full_html.replace('{{TITLE}}', doc_structure.title)
        
        # Save HTML file
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        filepath = os.path.join('exports', filename)
        os.makedirs('exports', exist_ok=True)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(full_html)
        
        return filepath
    
    def _get_html_template(self) -> str:
        """Return HTML template for slides"""
        return """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{TITLE}}</title>
    <style>
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 0;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: #333;
        }
        
        .slideshow-container {
            max-width: 1000px;
            position: relative;
            margin: auto;
            height: 100vh;
            display: flex;
            align-items: center;
        }
        
        .slide {
            display: none;
            padding: 60px;
            text-align: center;
            background: white;
            border-radius: 10px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.3);
            margin: 20px;
            min-height: 500px;
            display: flex;
            flex-direction: column;
            justify-content: center;
        }
        
        .slide.active {
            display: flex;
        }
        
        .title-slide h1 {
            font-size: 3em;
            color: #667eea;
            margin-bottom: 30px;
        }
        
        .content-slide h2 {
            font-size: 2.5em;
            color: #667eea;
            margin-bottom: 30px;
            border-bottom: 3px solid #667eea;
            padding-bottom: 15px;
        }
        
        .content-slide ul {
            text-align: left;
            font-size: 1.4em;
            line-height: 1.8;
            list-style-type: none;
            padding-left: 0;
        }
        
        .content-slide li {
            margin: 15px 0;
            padding-left: 30px;
            position: relative;
        }
        
        .content-slide li::before {
            content: "→";
            color: #667eea;
            font-weight: bold;
            position: absolute;
            left: 0;
        }
        
        .subtitle {
            font-size: 1.2em;
            color: #666;
            margin: 20px 0;
        }
        
        .date {
            font-size: 1em;
            color: #999;
        }
        
        .navigation {
            position: fixed;
            bottom: 30px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 15px;
        }
        
        .nav-btn {
            background: #667eea;
            color: white;
            border: none;
            padding: 10px 20px;
            border-radius: 25px;
            cursor: pointer;
            font-size: 1em;
            transition: all 0.3s ease;
        }
        
        .nav-btn:hover {
            background: #5a67d8;
            transform: translateY(-2px);
        }
        
        .nav-btn:disabled {
            background: #ccc;
            cursor: not-allowed;
            transform: none;
        }
        
        .slide-counter {
            position: fixed;
            top: 30px;
            right: 30px;
            background: rgba(255,255,255,0.9);
            padding: 10px 15px;
            border-radius: 20px;
            font-weight: bold;
            color: #667eea;
        }
    </style>
</head>
<body>
    <div class="slideshow-container">
        {{SLIDES}}
    </div>
    
    <div class="slide-counter">
        <span id="current-slide">1</span> / <span id="total-slides"></span>
    </div>
    
    <div class="navigation">
        <button class="nav-btn" id="prev-btn" onclick="changeSlide(-1)">← Previous</button>
        <button class="nav-btn" id="next-btn" onclick="changeSlide(1)">Next →</button>
    </div>
    
    <script>
        let currentSlide = 0;
        const slides = document.querySelectorAll('.slide');
        const totalSlides = slides.length;
        
        document.getElementById('total-slides').textContent = totalSlides;
        
        function showSlide(n) {
            slides.forEach(slide => slide.classList.remove('active'));
            
            if (n >= totalSlides) currentSlide = 0;
            if (n < 0) currentSlide = totalSlides - 1;
            
            slides[currentSlide].classList.add('active');
            document.getElementById('current-slide').textContent = currentSlide + 1;
            
            document.getElementById('prev-btn').disabled = currentSlide === 0;
            document.getElementById('next-btn').disabled = currentSlide === totalSlides - 1;
        }
        
        function changeSlide(n) {
            currentSlide += n;
            showSlide(currentSlide);
        }
        
        // Keyboard navigation
        document.addEventListener('keydown', function(e) {
            if (e.key === 'ArrowLeft') changeSlide(-1);
            if (e.key === 'ArrowRight') changeSlide(1);
        });
        
        // Initialize
        showSlide(0);
    </script>
</body>
</html>
        """

# Helper functions
def check_credentials_exist():
    """Check if credentials.json file exists and is valid"""
    if not os.path.exists(CLIENT_SECRETS_FILE):
        return False, "Credentials file not found"
    
    try:
        with open(CLIENT_SECRETS_FILE, 'r') as f:
            creds_data = json.load(f)
            
        # Basic validation
        if 'web' not in creds_data:
            return False, "Invalid credentials format - missing 'web' section"
        
        web_config = creds_data['web']
        required_fields = ['client_id', 'client_secret', 'auth_uri', 'token_uri']
        
        for field in required_fields:
            if field not in web_config:
                return False, f"Invalid credentials format - missing '{field}'"
        
        return True, "Valid"
        
    except json.JSONDecodeError:
        return False, "Invalid JSON format"
    except Exception as e:
        return False, f"Error reading credentials: {str(e)}"

def get_setup_status():
    """Get overall setup status for the application"""
    creds_valid, creds_msg = check_credentials_exist()
    
    return {
        'credentials_valid': creds_valid,
        'credentials_message': creds_msg,
        'setup_complete': creds_valid,
        'next_steps': [] if creds_valid else ['Set up Google API credentials']
    }

# Flask routes
@app.route('/')
def index():
    """Main page"""
    setup_status = get_setup_status()
    return render_template('docs_to_slides.html', setup_status=setup_status)

@app.route('/setup')
def setup_page():
    """Setup wizard page"""
    setup_status = get_setup_status()
    return render_template('setup.html', setup_status=setup_status)

@app.route('/api/setup-status')
def api_setup_status():
    """API endpoint to check setup status"""
    return jsonify(get_setup_status())

@app.route('/api/validate-credentials', methods=['POST'])
def validate_credentials():
    """Validate uploaded credentials file"""
    if 'credentials' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['credentials']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename.endswith('.json'):
        return jsonify({'error': 'File must be a JSON file'}), 400
    
    try:
        # Read and validate the uploaded file
        file_content = file.read().decode('utf-8')
        creds_data = json.loads(file_content)
        
        # Validate structure
        if 'web' not in creds_data:
            return jsonify({'error': 'Invalid credentials format - missing "web" section'}), 400
        
        web_config = creds_data['web']
        required_fields = ['client_id', 'client_secret', 'auth_uri', 'token_uri']
        
        for field in required_fields:
            if field not in web_config:
                return jsonify({'error': f'Missing required field: {field}'}), 400
        
        # If validation passes, save the file
        file.seek(0)  # Reset file pointer
        file.save(CLIENT_SECRETS_FILE)
        
        return jsonify({
            'success': True,
            'message': 'Credentials file uploaded and validated successfully',
            'client_id': web_config['client_id'][:20] + '...',  # Show partial for confirmation
            'project_id': web_config.get('project_id', 'Unknown')
        })
        
    except json.JSONDecodeError:
        return jsonify({'error': 'Invalid JSON format'}), 400
    except Exception as e:
        return jsonify({'error': f'Error processing file: {str(e)}'}), 500

@app.route('/download-template')
def download_template():
    """Download credentials template file"""
    template_path = 'credentials_template.json'
    if os.path.exists(template_path):
        return send_file(template_path, as_attachment=True, 
                        download_name='credentials_template.json')
    else:
        return jsonify({'error': 'Template file not found'}), 404

@app.route('/auth')
def auth():
    """Initiate Google OAuth flow"""
    if not os.path.exists(CLIENT_SECRETS_FILE):
        return jsonify({
            'error': 'Google credentials file not found. Please add credentials.json'
        }), 400
    
    flow = Flow.from_client_secrets_file(
        CLIENT_SECRETS_FILE,
        scopes=SCOPES,
        redirect_uri=REDIRECT_URI
    )
    
    authorization_url, state = flow.authorization_url(
        access_type='offline',
        include_granted_scopes='true'
    )
    
    flask.session['state'] = state
    return redirect(authorization_url)

@app.route('/callback')
def callback():
    """Handle Google OAuth callback"""
    state = flask.session['state']
    
    flow = Flow.from_client_secrets_file(
        CLIENT_SECRETS_FILE,
        scopes=SCOPES,
        state=state,
        redirect_uri=REDIRECT_URI
    )
    
    flow.fetch_token(authorization_response=request.url)
    
    credentials = flow.credentials
    flask.session['credentials'] = {
        'token': credentials.token,
        'refresh_token': credentials.refresh_token,
        'token_uri': credentials.token_uri,
        'client_id': credentials.client_id,
        'client_secret': credentials.client_secret,
        'scopes': credentials.scopes
    }
    
    return redirect(url_for('index'))

@app.route('/convert', methods=['POST'])
def convert_document():
    """Convert Google Doc to slides"""
    if 'credentials' not in flask.session:
        return jsonify({'error': 'Not authenticated'}), 401
    
    data = request.get_json()
    document_url = data.get('document_url')
    export_format = data.get('format', 'pptx')
    
    if not document_url:
        return jsonify({'error': 'Document URL is required'}), 400
    
    # Extract document ID from URL
    doc_id = extract_doc_id(document_url)
    if not doc_id:
        return jsonify({'error': 'Invalid Google Docs URL'}), 400
    
    try:
        # Create credentials object
        credentials = Credentials(**flask.session['credentials'])
        
        # Extract document content
        extractor = GoogleDocsExtractor(credentials)
        doc_structure = extractor.extract_document(doc_id)
        
        # Generate slides
        generator = SlideGenerator()
        
        if export_format == 'pptx':
            filepath = generator.create_powerpoint(doc_structure)
        elif export_format == 'html':
            filepath = generator.create_html_slides(doc_structure)
        else:
            return jsonify({'error': 'Unsupported format'}), 400
        
        return jsonify({
            'success': True,
            'filename': os.path.basename(filepath),
            'download_url': f'/download/{os.path.basename(filepath)}'
        })
        
    except Exception as e:
        logger.error(f"Error converting document: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>')
def download_file(filename):
    """Download generated presentation file"""
    filepath = os.path.join('exports', filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    else:
        return jsonify({'error': 'File not found'}), 404

def extract_doc_id(url: str) -> Optional[str]:
    """Extract Google Docs document ID from URL"""
    patterns = [
        r'/document/d/([a-zA-Z0-9-_]+)',
        r'id=([a-zA-Z0-9-_]+)',
    ]
    
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    
    return None

if __name__ == '__main__':
    # Create necessary directories
    os.makedirs('exports', exist_ok=True)
    
    app.run(debug=True, port=5000)